# -*- coding: utf-8 -*-
# 编辑历史:
# 2026-07-20 - 小欧 - 自然单位翻页 feat:
#   1. read_pptx 增加 slide=N 按页翻页参数
#   2. llm_data 增加 current_slide 指标
#     及 warning 状态支持
#   3. 新增 slide 参数校验与单页选取逻辑
#   4. 迭代变量 sld 改名防冲突(param slide)
# 2026-07-21 - 小欧 - 文件大小安全检测: 入口加 READ_PPTX_INPUT_MAX_BYTES 字节检查，超限拒读防OOM
# 2026-07-23 - 小欧 - 三堂会审5bug修复: total_slides未反映outlimit截断+notes_data未同步截断+删死import os as _os_mod
# 2026-07-24 - 小欧 - 修复: error summary嵌入full detail → 改用truncate_summary(detail)首行; warning summary去掉三重重复detail
# 2026-07-26 - 小欧 - OOD: 确认READ_PPTX_INPUT_MAX_BYTES未落地,OOM自然抛出被except捕获(同dataanalysis模式)
# 2026-07-26 - 小欧 - 清理: 删logger死import(全文件无logger调用)
# 2026-07-26 - 小沈 - BugFix #1/#4: 截断前存_actual_slide_count,产_trunc_hint提示用slide=N读剩余页; 参数path不覆盖(Bug #3)
# 2026-07-31 - 小欧 - Bug㉑修复: 全量读取成功summary的slide_count改传total_slides(实际返回页数), 截断后"提示页数"与返回数据口径一致; 完整页数由_trunc_hint告知 | py_compile ✓
# 2026-08-13 - 小欧 - A5职责拆分: hint_* 错误提示函数/导入源改 app.tools.toolhelper.error_hints
"""
D3: read_pptx — 读取PPT文档

从document_tools.py拆分而来 — 小欧 2026-06-22
"""
# 【铁规1】helper/被调函数(以下划线_开头的函数)只返回raw dict，严禁调用build_success/build_error/build_warning和构建llm_data。
# build3+llm_data只能在tool的main函数(对外公开的函数)中包装。违反此规则的代码视为不合规。
# 【铁规2】工具返回原始data，禁止调用truncate_data_for_frontend。截断只能在前端yield层。
# 【铁规3】计时(duration_ms计算)只能在tool的主函数中，严禁在子函数/helper中计时。

import time as _time_mod
from pathlib import Path
from typing import Any, Dict, Optional

from app.tools.tool_response import build_success, build_error, build_warning
from app.tools.tool_fc_helper import _check_module
from app.tools.validate.file_type_checker import check_for_document_tool
from app.tools.toolhelper.error_hints import hint_for_read_error
from app.tools.tool_constants import ERR_DOC_READ_PPTX, READ_PPTX_OUTLIMIT_CHARS
from app.utils.text_utils import truncate_summary


def _build_read_pptx_llm_data(
    exec_code: str, duration_ms: int,
    file_path: str = "", slide_count: int = 0, slides_read: int = 0,
    text_len: int = 0, table_count: int = 0, image_count: int = 0,
    current_slide: Optional[int] = None, detail: str = "", hint: str = "",
) -> Dict[str, Any]:
    """read_pptx的llm_data构建函数 — 小健 2026-06-21 — 小欧 2026-06-22 — 小欧 2026-07-05 加hint参数 — 小欧 2026-07-20 加 current_slide(按页翻页)"""
    if exec_code == "error":
        _err_summary = truncate_summary(detail)
        return {
            "summary": f"读取PPT{file_path}，失败" + (f": {_err_summary}" if _err_summary else ""),
            "action": {"tool": "read_pptx", "tool_zh": "读取PPT", "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": "读取PPT失败", "code": ERR_DOC_READ_PPTX, "detail": detail, "hint": hint if hint else "读取失败,详见错误明细"},
            "duration_ms": duration_ms,
            "metrics": {},
        }
    # success / warning (自然单位治理 2026-07-20 小欧: 按页翻页)
    if exec_code == "warning":
        summary_str = f"读取PPT{file_path}，警告"
    elif current_slide is not None:
        summary_str = f"读取PPT{file_path}，成功: 第{current_slide}页(共{slide_count}页)，{text_len}字符"
    else:
        summary_str = f"读取PPT{file_path}，成功: {slide_count}页，{text_len}字符"
    return {
        "summary": summary_str,
        "action": {"tool": "read_pptx", "tool_zh": "读取PPT", "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": exec_code, "message": "读取PPT成功", "code": "", "detail": detail, "hint": hint},
        "duration_ms": duration_ms,
        "metrics": {
            "slide_count": {"value": slide_count, "text": f"{slide_count}页"},
            "text_len": {"value": text_len, "text": f"{text_len}字符"},
        },
    }


def read_pptx(path: str, slide: Optional[int] = None) -> Dict[str, Any]:
    """读取PPT文件 — 小沈 2026-06-19 — 小欧 2026-06-22 独立文件 — 小欧 2026-06-24 增加文件类型前置检查"""
    t0 = _time_mod.perf_counter()
    file_path = path

    # 文件类型前置检查 — 小欧 2026-06-24
    is_valid, error_detail, suggested_tool = check_for_document_tool(path)
    if not is_valid:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        if suggested_tool:
            _hint = f"建议使用{suggested_tool}工具"
        elif suggested_tool == "":
            _hint = "请检查文件路径和文件名是否正确"
        else:
            _hint = "文件类型不匹配,请使用正确的文档格式"
        llm_data = _build_read_pptx_llm_data("error", duration_ms, file_path, detail=error_detail, hint=_hint)
        return build_error(data={}, llm_data=llm_data)

    if not _check_module("pptx"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pptx_llm_data("error", duration_ms, file_path, detail="python-pptx库未安装", hint="请安装python-pptx库")
        return build_error(data={}, llm_data=llm_data)

    try:
        from pptx import Presentation

        _p = Path(file_path)
        prs = Presentation(_p)
        slides_data = []
        notes_data = []

        for slide_num, sld in enumerate(prs.slides, 1):
            slide_text = []
            tables_data = []
            for shape in sld.shapes:
                if shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_rows.append(row_data)
                        slide_text.append(" | ".join(row_data))
                    tables_data.append(table_rows)
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)

            slide_dict = {
                "slide_num": slide_num,
                "text": "\n".join(slide_text),
            }
            if tables_data:
                slide_dict["tables"] = tables_data
            slides_data.append(slide_dict)

            if sld.has_notes_slide:
                notes = sld.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_data.append({
                        "slide_num": slide_num,
                        "notes": notes,
                    })

        total_text = sum(len(s.get("text", "")) for s in slides_data)
        # 截断前存实际总页数 — 小沈 2026-07-26
        _actual_slide_count = len(slides_data)
        # outlimit: 仅全量读取(无slide参数)按幻灯片粒度截断
        if slide is None and total_text > READ_PPTX_OUTLIMIT_CHARS:
            acc = 0
            cut_idx = len(slides_data)
            for i, s in enumerate(slides_data):
                acc += len(s.get("text", ""))
                if acc > READ_PPTX_OUTLIMIT_CHARS:
                    cut_idx = i + 1
                    break
            slides_data = slides_data[:cut_idx]
            total_text = sum(len(s.get("text", "")) for s in slides_data)
            # 同步截断 notes_data 过滤已截断 slide 的笔记
            _valid_slide_nums = {s["slide_num"] for s in slides_data}
            notes_data = [n for n in notes_data if n["slide_num"] in _valid_slide_nums]
        total_slides = len(slides_data)
        # 截断提示: LLM获知实际总页数+如何读剩余页 — 小沈 2026-07-26
        _trunc_hint = ""
        if _actual_slide_count > total_slides:
            _trunc_hint = f"文档共{_actual_slide_count}页,文本超{READ_PPTX_OUTLIMIT_CHARS}字符仅返回前{total_slides}页,使用slide=N读取指定页"
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)

        # —— 自然单位治理(2026-07-20 小欧): 按幻灯片翻页(slide=N) ——
        if slide is not None:
            if slide < 1:
                llm_data = _build_read_pptx_llm_data(
                    "error", duration_ms, file_path, total_slides, detail=f"slide参数必须>=1,当前值: {slide}", hint="slide参数从1开始,表示第几页")
                return build_error(data={}, llm_data=llm_data)
            if slide > total_slides:
                llm_data = _build_read_pptx_llm_data(
                    "warning", duration_ms, file_path, total_slides, detail=f"slide={slide}超出范围(共{total_slides}页),返回空内容", hint=f"PPT共{total_slides}页,slide取值1..{total_slides}")
                result_data = {"slides": []}
                if notes_data:
                    result_data["notes"] = notes_data
                return build_warning(data=result_data, llm_data=llm_data)
            slides_data = [slides_data[slide - 1]]
            sel_text = len(slides_data[0].get("text", ""))
            result_data = {"slides": slides_data}
            if notes_data:
                result_data["notes"] = notes_data
            llm_data = _build_read_pptx_llm_data(
                "success", duration_ms, file_path, total_slides, text_len=sel_text, current_slide=slide)
            return build_success(data=result_data, llm_data=llm_data)

        result_data = {
            "slides": slides_data,
        }
        if notes_data:
            result_data["notes"] = notes_data

        # 2026-07-31 小欧: Bug㉑修复 — slide_count传total_slides(实际返回页数), 与text_len/返回数据口径一致; 完整页数由_trunc_hint告知
        llm_data = _build_read_pptx_llm_data("success", duration_ms, file_path, total_slides, text_len=total_text, hint=_trunc_hint)
        # =============================================================================
        # 数据设计：slide_count 从 data 移除，通过 llm_data.metrics 传入 summary
        # summary 示例: "读取PPT成功: 10页, 5000字符"
        # — 小欧 2026-07-06 18:46:13; 2026-07-20 加 slide 参数(按页翻页)
        # =============================================================================
        # ---- observation_formatter route -------------------------------------------
        # branch: #16 slides items
        # trigger: "slides" in data — slides 是 List[dict], 每项含 slide_num/text/tables
        # handler: _format_slides(data)
        # file:    observation_formatter.py:_format_slides
        # ------------------------------------------------------------------------------
        return build_success(data=result_data, llm_data=llm_data)

    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pptx_llm_data("error", duration_ms, file_path, detail=str(e), hint=hint_for_read_error(e, file_path))
        return build_error(data={}, llm_data=llm_data)
