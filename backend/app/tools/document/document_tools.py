# -*- coding: utf-8 -*-
"""
文档读写工具函数模块
【设计说明 2026-06-17 北京老陈确认】本文件是按工具分类聚合的实现文件，文件大是正常设计。后续审查关注功能逻辑本身的代码10大规范遵守和最优美简洁性，禁止以"文件过大"作为问题提出。

【创建时间】2026-05-02 小沈
【设计依据】按文档第8.3节 Tool 80-82 定义
【重构 2026-05-18 小健】8个旧函数抽取为内部函数

【重要】新函数增加规范 - 小沈 2026-05-04
新增函数时必须同步修改以下3个文件:
1. *_tools.py: 函数实现(必须有详细注释)
2. *_schema.py: Pydantic 模型(输入参数定义)
3. *_register.py: 显式注册(description + examples + input_model)

包含:
- _read_pdf / _read_docx / _read_xlsx / _read_xls / _read_pptx: 内部读取函数
- _write_docx / _write_xlsx / _write_pdf / _write_pptx: 内部写入函数

Author: 小沈 - 2026-05-02
【新增 2026-05-05 小沈】write_pdf, convert_document
【重构 2026-05-18 小健】8合2路由重构
"""

import csv
import os
import shutil
import subprocess
import tempfile
import time as _time_mod
from typing import Dict, Any, List, Optional, Literal, Union, Tuple
from pathlib import Path


from app.tools.tool_response import build_success, build_error, build_warning
from app.tools.toolhelper.common_helper import _check_module
from app.tools.tool_constants import SUBPROCESS_TIMEOUT_LONG
from app.utils.json_utils import coerce_json
from app.constants import (
    ERR_DOC_CHART_GENERATE,
    ERR_DOC_CONVERT_FAILED,
    ERR_DOC_NO_OPENPYXL,
    ERR_DOC_NO_PPTX,
    ERR_DOC_READ_CSV,
    ERR_DOC_READ_DOCX,
    ERR_DOC_READ_PDF,
    ERR_DOC_READ_PPTX,
    ERR_DOC_READ_XLSX,
    ERR_DOC_WRITE_PPTX,
    ERR_NO_DOCX,
    ERR_NO_PDFPLUMBER,
    ERR_NO_REPORTLAB,
    ERR_WRITE_DOCX,
    ERR_WRITE_PDF,
    ERR_WRITE_XLSX,
)


def _build_read_pdf_llm_data(exec_code, duration_ms, file_path="", page_count=0, pages_read=0,
                              text_len=0, table_count=0, image_count=0, detail=""):
    """read_pdf的llm_data构建函数 — 小健 2026-06-21"""
    if exec_code == "error":
        return {
            "summary": f"读取PDF失败: {detail}",
            "action": {"tool": "read_pdf", "tool_zh": "读取PDF", "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": "读取PDF失败", "code": ERR_DOC_READ_PDF, "detail": detail, "hint": "请检查文件路径和格式"},
            "duration_ms": duration_ms, "metrics": {},
        }
    return {
        "summary": f"读取PDF成功: {pages_read}/{page_count}页, {text_len}字符",
        "action": {"tool": "read_pdf", "tool_zh": "读取PDF", "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": "success", "message": "读取PDF成功", "code": "", "detail": "", "hint": ""},
        "duration_ms": duration_ms,
        "metrics": {"page_count": {"value": page_count, "text": f"{page_count}页"}, "pages_read": {"value": pages_read, "text": f"读取{pages_read}页"}, "text_len": {"value": text_len, "text": f"{text_len}字符"}},
    }


def _build_read_docx_llm_data(exec_code, duration_ms, file_path="", para_count=0, text_len=0, detail=""):
    """read_docx的llm_data构建函数 — 小健 2026-06-21"""
    if exec_code == "error":
        return {
            "summary": f"读取Word失败: {detail}",
            "action": {"tool": "read_docx", "tool_zh": "读取Word", "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": "读取Word失败", "code": ERR_DOC_READ_DOCX, "detail": detail, "hint": "请检查文件路径和格式"},
            "duration_ms": duration_ms, "metrics": {},
        }
    return {
        "summary": f"读取Word成功: {para_count}段, {text_len}字符",
        "action": {"tool": "read_docx", "tool_zh": "读取Word", "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": "success", "message": "读取Word成功", "code": "", "detail": "", "hint": ""},
        "duration_ms": duration_ms,
        "metrics": {"para_count": {"value": para_count, "text": f"{para_count}段"}, "text_len": {"value": text_len, "text": f"{text_len}字符"}},
    }


def _build_read_xlsx_llm_data(exec_code, duration_ms, file_path="", row_count=0, sheet_count=0, detail=""):
    """read_xlsx的llm_data构建函数 — 小健 2026-06-21"""
    if exec_code == "error":
        return {
            "summary": f"读取Excel失败: {detail}",
            "action": {"tool": "read_xlsx", "tool_zh": "读取Excel", "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": "读取Excel失败", "code": ERR_DOC_READ_XLSX, "detail": detail, "hint": "请检查文件路径和格式"},
            "duration_ms": duration_ms, "metrics": {},
        }
    return {
        "summary": f"读取Excel成功: {row_count}行, {sheet_count}个工作表",
        "action": {"tool": "read_xlsx", "tool_zh": "读取Excel", "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": "success", "message": "读取Excel成功", "code": "", "detail": "", "hint": ""},
        "duration_ms": duration_ms,
        "metrics": {"row_count": {"value": row_count, "text": f"{row_count}行"}, "sheet_count": {"value": sheet_count, "text": f"{sheet_count}个表"}},
    }


def _build_read_pptx_llm_data(exec_code, duration_ms, file_path="", slide_count=0, text_len=0, detail=""):
    """read_pptx的llm_data构建函数 — 小健 2026-06-21"""
    if exec_code == "error":
        return {
            "summary": f"读取PPT失败: {detail}",
            "action": {"tool": "read_pptx", "tool_zh": "读取PPT", "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": "读取PPT失败", "code": ERR_DOC_READ_PPTX, "detail": detail, "hint": "请检查文件路径和格式"},
            "duration_ms": duration_ms, "metrics": {},
        }
    return {
        "summary": f"读取PPT成功: {slide_count}页, {text_len}字符",
        "action": {"tool": "read_pptx", "tool_zh": "读取PPT", "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": "success", "message": "读取PPT成功", "code": "", "detail": "", "hint": ""},
        "duration_ms": duration_ms,
        "metrics": {"slide_count": {"value": slide_count, "text": f"{slide_count}页"}, "text_len": {"value": text_len, "text": f"{text_len}字符"}},
    }


def _build_write_doc_llm_data(exec_code, duration_ms, tool_name, tool_zh, file_path="", detail=""):
    """write_doc/xlsx/pdf/pptx的通用llm_data构建函数 — 小健 2026-06-21"""
    if exec_code == "error":
        return {
            "summary": f"写入{tool_zh}失败: {detail}",
            "action": {"tool": tool_name, "tool_zh": tool_zh, "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": f"写入{tool_zh}失败", "code": "", "detail": detail, "hint": "请检查路径和权限"},
            "duration_ms": duration_ms, "metrics": {},
        }
    return {
        "summary": f"写入{tool_zh}成功: {file_path}",
        "action": {"tool": tool_name, "tool_zh": tool_zh, "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": "success", "message": f"写入{tool_zh}成功", "code": "", "detail": "", "hint": ""},
        "duration_ms": duration_ms, "metrics": {},
    }





def _check_pdf_readable(file_path: str) -> Dict[str, Any]:
    """检查PDF文件是否可读(内部helper) - 小沈 2026-05-18,从env_check_tools.py迁入
    【2026-06-21 小健】builder改造，新3字段result
    """
    path = Path(file_path)
    if not path.exists():
        llm_data = _build_read_pdf_llm_data("error", 0, file_path, detail=f"文件不存在: {file_path}")
        return build_error(data={"file_path": file_path}, llm_data=llm_data)
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
        return build_success(data={"readable": True, "page_count": page_count}, llm_data=_build_read_pdf_llm_data("success", 0, file_path, page_count))
    except ImportError:
        llm_data = _build_read_pdf_llm_data("error", 0, file_path, detail="pdfplumber库未安装")
        return build_error(data={"error": "pdfplumber库未安装"}, llm_data=llm_data)
    except Exception as e:
        llm_data = _build_read_pdf_llm_data("error", 0, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _check_docx_readable(file_path: str) -> Dict[str, Any]:
    """检查Word文件是否可读(内部helper) - 小沈 2026-05-18,从env_check_tools.py迁入
    【2026-06-21 小健】builder改造，新3字段result
    """
    path = Path(file_path)
    if not path.exists():
        llm_data = _build_read_docx_llm_data("error", 0, file_path, detail=f"文件不存在: {file_path}")
        return build_error(data={"file_path": file_path}, llm_data=llm_data)
    try:
        import docx
        doc = docx.Document(path)
        para_count = len(doc.paragraphs)
        return build_success(data={"readable": True, "paragraph_count": para_count}, llm_data=_build_read_docx_llm_data("success", 0, file_path, para_count))
    except ImportError:
        llm_data = _build_read_docx_llm_data("error", 0, file_path, detail="python-docx库未安装")
        return build_error(data={"error": "python-docx库未安装"}, llm_data=llm_data)
    except Exception as e:
        llm_data = _build_read_docx_llm_data("error", 0, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _check_xlsx_readable(file_path: str) -> Dict[str, Any]:
    """检查Excel文件是否可读(内部helper) - 小沈 2026-05-18,从env_check_tools.py迁入
    【2026-06-21 小健】builder改造，新3字段result
    """
    path = Path(file_path)
    if not path.exists():
        llm_data = _build_read_xlsx_llm_data("error", 0, file_path, detail=f"文件不存在: {file_path}")
        return build_error(data={"file_path": file_path}, llm_data=llm_data)
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True)
        sheet_names = wb.sheetnames
        wb.close()
        return build_success(data={"readable": True, "sheet_names": sheet_names}, llm_data=_build_read_xlsx_llm_data("success", 0, file_path, 0, len(sheet_names)))
    except ImportError:
        llm_data = _build_read_xlsx_llm_data("error", 0, file_path, detail="openpyxl库未安装")
        return build_error(data={"error": "openpyxl库未安装"}, llm_data=llm_data)
    except Exception as e:
        llm_data = _build_read_xlsx_llm_data("error", 0, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)





def _build_validate_chart_llm_data(exec_code, errors=None):
    """_validate_chart_data的llm_data构建函数 — 小健 2026-06-21"""
    errors = errors or []
    if exec_code == "error":
        return {
            "summary": f"图表数据格式有{len(errors)}个问题",
            "action": {"tool": "_validate_chart_data", "tool_zh": "验证图表", "target": "", "params": {}},
            "status": {"exec_code": "error", "message": "验证失败", "code": ERR_DOC_CHART_GENERATE, "detail": "; ".join(errors), "hint": "请检查labels和values字段"},
            "duration_ms": 0, "metrics": {},
        }
    return {
        "summary": "图表数据格式正确",
        "action": {"tool": "_validate_chart_data", "tool_zh": "验证图表", "target": "", "params": {}},
        "status": {"exec_code": "success", "message": "验证通过", "code": "", "detail": "", "hint": ""},
        "duration_ms": 0, "metrics": {},
    }


def _validate_chart_data(data: Dict[str, Any]) -> Dict[str, Any]:
    """验证图表数据格式(内部helper) - 小沈 2026-05-18,从env_check_tools.py迁入
    【2026-06-21 小健】builder改造，新3字段result
    """
    errors = []
    if not isinstance(data, dict):
        errors.append("data必须是字典类型")
    else:
        if "labels" not in data:
            errors.append("缺少labels字段")
        elif not isinstance(data["labels"], list):
            errors.append("labels必须是数组类型")
        if "values" not in data:
            errors.append("缺少values字段")
        elif not isinstance(data["values"], list):
            errors.append("values必须是数组类型")
        if "labels" in data and "values" in data:
            if isinstance(data["labels"], list) and isinstance(data["values"], list):
                if len(data["labels"]) != len(data["values"]):
                    errors.append(f"labels和values长度不一致: labels={len(data['labels'])}, values={len(data['values'])}")
    is_valid = len(errors) == 0
    if is_valid:
        return build_success(data={"valid": True}, llm_data=_build_validate_chart_llm_data("success"))
    else:
        return build_error(data={"valid": False, "errors": errors}, llm_data=_build_validate_chart_llm_data("error", errors))





def _parse_pages(pages_str: str) -> List[int]:
    """解析页码字符串为页码列表 - 小沈 2026-05-02"""
    result = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            try:
                result.extend(range(int(start), int(end) + 1))
            except ValueError:
                pass
        else:
            try:
                result.append(int(part))
            except ValueError:
                pass
    return sorted(set(result))


# ============================================================
# 内部读取函数(原 read_pdf/read_docx/read_xlsx/read_pptx 逻辑)
# ============================================================

def _process_page(page, page_num: int,
                  extract_tables: bool, extract_images: bool
                  ) -> Tuple[str, List[Dict], List[Dict]]:
    """处理单页PDF:返回 (text, tables_data, images_data)

    小沈 2026-05-25 重构拆分
    """
    text = page.extract_text() or ""

    tables = []
    if extract_tables:
        for idx, t in enumerate(page.extract_tables() or []):
            tables.append({"page": page_num, "table_idx": idx, "data": t})

    images = []
    if extract_images:
        for idx, img in enumerate(page.images or []):
            images.append({
                "page": page_num, "image_idx": idx,
                "x0": float(img.get("x0", 0)), "y0": float(img.get("y0", 0)),
                "x1": float(img.get("x1", 0)), "y1": float(img.get("y1", 0)),
                "width": float(img.get("width", 0)), "height": float(img.get("height", 0)),
            })

    return text, tables, images


def _build_pdf_result(fp: str, pages_read: List[int], all_text: List[str],
                      tables_data: List[Dict], images_data: List[Dict],
                      page_count: int, duration_ms: int) -> dict:
    """构建统一的 PDF 读取返回结果
    【2026-06-21 小健】builder改造，新3字段result
    """
    full_text = "\n\n".join(all_text)
    result = {"text": full_text, "page_count": page_count, "pages_read": pages_read}
    if tables_data:
        result["tables"] = tables_data
        result["table_count"] = len(tables_data)
    if images_data:
        result["images"] = images_data
        result["image_count"] = len(images_data)
    llm_data = _build_read_pdf_llm_data("success", duration_ms, fp, page_count, len(pages_read),
                                          len(full_text), len(tables_data), len(images_data))
    return build_success(data=result, llm_data=llm_data)


def _read_pdf(
    file_path: str,
    pages: str = None,
    extract_images: bool = False,
    extract_tables: bool = False
) -> Dict[str, Any]:
    """读取PDF文件并提取文本内容(内部函数) - 小健 2026-05-18
    【2026-05-25 小沈重构】拆分:逐页提取 → _process_page,结果构建 → _build_pdf_result
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("pdfplumber"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pdf_llm_data("error", duration_ms, file_path, detail="pdfplumber库未安装")
        return build_error(data={"error": "pdfplumber库未安装"}, llm_data=llm_data)

    try:
        import pdfplumber

        path = Path(file_path)
        if not path.exists():
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_pdf_llm_data("error", duration_ms, file_path, detail=f"文件不存在: {file_path}")
            return build_error(data={"file_path": file_path}, llm_data=llm_data)

        all_text, pages_read, tables_data, images_data = [], [], [], []
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            target = _parse_pages(pages) if pages else list(range(1, page_count + 1))
            target = [p for p in target if 1 <= p <= page_count]

            for pn in target:
                page = pdf.pages[pn - 1]
                text, tables, images = _process_page(page, pn, extract_tables, extract_images)
                all_text.append(f"--- 第 {pn} 页 ---\n{text}")
                pages_read.append(pn)
                tables_data.extend(tables)
                images_data.extend(images)

        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        return _build_pdf_result(file_path, pages_read, all_text, tables_data, images_data, page_count, duration_ms)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pdf_llm_data("error", duration_ms, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _read_docx(
    file_path: str,
    extract_tables: bool = False
) -> Dict[str, Any]:
    """读取Word文档并提取文本内容(内部函数) - 小健 2026-05-18
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("docx"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_docx_llm_data("error", duration_ms, file_path, detail="python-docx库未安装")
        return build_error(data={"error": "python-docx库未安装"}, llm_data=llm_data)

    try:
        import docx

        path = Path(file_path)
        if not path.exists():
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_docx_llm_data("error", duration_ms, file_path, detail=f"文件不存在: {file_path}")
            return build_error(data={"file_path": file_path}, llm_data=llm_data)

        doc = docx.Document(path)
        paragraphs = [para.text for para in doc.paragraphs]
        text = "\n".join(paragraphs)
        
        result_data = {
            "text": text,
            "paragraph_count": len(paragraphs),
        }
        
        if extract_tables:
            tables_data = []
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_rows.append(row_data)
                tables_data.append(table_rows)
            result_data["tables"] = tables_data
            result_data["table_count"] = len(tables_data)
        
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_docx_llm_data("success", duration_ms, file_path, len(paragraphs), len(text))
        return build_success(data=result_data, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_docx_llm_data("error", duration_ms, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _read_xlsx(
    file_path: str,
    sheet_name: str = None,
    max_rows: int = 1000,
    header: bool = True
) -> Dict[str, Any]:
    """读取Excel文件并提取表格数据(内部函数) - 小健 2026-05-18
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("openpyxl"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail="openpyxl库未安装")
        return build_error(data={"error": "openpyxl库未安装"}, llm_data=llm_data)

    try:
        from openpyxl import load_workbook

        path = Path(file_path)
        if not path.exists():
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=f"文件不存在: {file_path}")
            return build_error(data={"file_path": file_path}, llm_data=llm_data)

        wb = load_workbook(path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames

        if sheet_name:
            if sheet_name not in sheet_names:
                wb.close()
                duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
                llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=f"工作表不存在: {sheet_name}")
                return build_error(data={"sheet_name": sheet_name, "available_sheets": sheet_names}, llm_data=llm_data)
            ws = wb[sheet_name]
        else:
            ws = wb.active

        rows = []
        headers = []
        row_count = 0

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows + (1 if header else 0):
                break
            row_data = [
                None if val is None else val
                for val in row
            ]
            if i == 0 and header:
                headers = [str(h) if h is not None else f"column_{j}" for j, h in enumerate(row_data)]
            else:
                if i == 0 and not header:
                    headers = [f"column_{j}" for j in range(len(row_data))]
                rows.append(row_data)
                row_count += 1

        wb.close()
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("success", duration_ms, file_path, row_count, len(sheet_names))
        return build_success(data={"headers": headers, "rows": rows, "row_count": row_count, "sheet_names": sheet_names}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _read_xls(
    file_path: str,
    max_rows: int = 1000,
) -> Dict[str, Any]:
    """读取旧版Excel(.xls)文件(内部函数) - 小沈 2026-06-19
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("xlrd"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail="xlrd库未安装")
        return build_error(data={"error": "xlrd库未安装", "file_path": file_path}, llm_data=llm_data)

    try:
        import xlrd

        path = Path(file_path)
        if not path.exists():
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=f"文件不存在: {file_path}")
            return build_error(data={"file_path": file_path}, llm_data=llm_data)

        wb = xlrd.open_workbook_xls(str(path))
        sheet_names = wb.sheet_names()
        ws = wb.sheet_by_index(0)

        headers = []
        rows = []
        for i in range(min(ws.nrows, max_rows)):
            row_data = [ws.cell_value(i, j) for j in range(ws.ncols)]
            if i == 0:
                headers = [str(h) if h else f"column_{j}" for j, h in enumerate(row_data)]
            else:
                rows.append(row_data)

        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("success", duration_ms, file_path, len(rows), len(sheet_names))
        return build_success(data={"headers": headers, "rows": rows, "row_count": len(rows), "sheet_names": sheet_names}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _read_csv_stdlib(
    file_path: str,
    encoding: str = "utf-8",
    delimiter: str = ",",
    has_header: bool = True,
    max_rows: int = 1000
) -> Dict[str, Any]:
    """使用标准库csv读取CSV文件(内部函数)— 小健 2026-05-18
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    try:
        path = Path(file_path)
        if not path.exists():
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=f"文件不存在: {file_path}")
            return build_error(data={"file_path": file_path}, llm_data=llm_data)
        
        rows = []
        headers = []
        encodings_to_try = [encoding, "gbk", "gb2312", "latin-1"] if encoding == "utf-8" else [encoding, "utf-8", "latin-1"]
        read_ok = False
        for enc in encodings_to_try:
            try:
                with open(path, "r", encoding=enc, newline="") as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    for i, row in enumerate(reader):
                        if i >= max_rows:
                            break
                        if i == 0:
                            if has_header:
                                headers = row
                            else:
                                headers = [f"col_{j}" for j in range(len(row))]
                                rows.append(row)
                        else:
                            rows.append(row)
                read_ok = True
                break
            except UnicodeDecodeError:
                continue
        if not read_ok:
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=f"编码不匹配(尝试了{encodings_to_try})")
            return build_error(data={"file_path": file_path, "encodings_tried": encodings_to_try}, llm_data=llm_data)
        
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("success", duration_ms, file_path, len(rows))
        return build_success(data={"headers": headers, "rows": rows, "row_count": len(rows)}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_xlsx_llm_data("error", duration_ms, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)





def _read_pptx(
    file_path: str,
    extract_notes: bool = False
) -> Dict[str, Any]:
    """读取PPT幻灯片(内部函数) - 小健 2026-05-18
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("pptx"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pptx_llm_data("error", duration_ms, file_path, detail="python-pptx库未安装")
        return build_error(data={"error": "python-pptx库未安装"}, llm_data=llm_data)

    try:
        from pptx import Presentation

        path = Path(file_path)
        if not path.exists():
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_pptx_llm_data("error", duration_ms, file_path, detail=f"文件不存在: {file_path}")
            return build_error(data={"file_path": file_path}, llm_data=llm_data)

        prs = Presentation(path)
        slides_data = []
        notes_data = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
            
            slides_data.append({
                "slide_num": slide_num,
                "text": "\n".join(slide_text)
            })
            
            if extract_notes and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_data.append({
                        "slide_num": slide_num,
                        "notes": notes
                    })

        result_data = {
            "slide_count": len(prs.slides),
            "slides": slides_data,
        }

        if extract_notes and notes_data:
            result_data["notes"] = notes_data

        _total_text = sum(len(s.get("text", "")) for s in slides_data)
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pptx_llm_data("success", duration_ms, file_path, len(prs.slides), _total_text)
        return build_success(data=result_data, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_read_pptx_llm_data("error", duration_ms, file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


# ============================================================
# 公用函数 — paragraphs 解析 + 内容元素处理
# ============================================================

def _resolve_paragraphs(paragraphs):
    """解析paragraphs参数: 统一输出为 (title_from_dict, content_list) - 小欧 2026-06-19"""
    title = None
    items = []

    if isinstance(paragraphs, str):
        items = [paragraphs]
    elif isinstance(paragraphs, list):
        items = paragraphs
    elif isinstance(paragraphs, dict):
        title = paragraphs.get("title")
        content = paragraphs.get("content", [])
        items = content if isinstance(content, list) else [content]

    return title, items


def _add_docx_content_item(doc, item):
    """处理单个Word内容元素(list中的str或dict) - 小欧 2026-06-19"""
    if isinstance(item, str):
        if item.strip():
            doc.add_paragraph(item)
    elif isinstance(item, dict):
        item_type = item.get("type", "paragraph")
        item_text = item.get("text", "")
        if item_type in ("h1", "heading"):
            doc.add_heading(item_text, item.get("level", 1))
        elif item_type in ("h2", "h3", "h4", "h5"):
            doc.add_heading(item_text, level=int(item_type[1]))
        elif item_type == "paragraph":
            doc.add_paragraph(item_text)
        elif item_type == "table":
            rows_data = item.get("rows", [])
            if rows_data:
                t = doc.add_table(rows=len(rows_data), cols=len(rows_data[0]))
                for ri, rd in enumerate(rows_data):
                    for ci, cv in enumerate(rd):
                        t.rows[ri].cells[ci].text = str(cv)
    else:
        text = str(item)
        if text.strip():
            doc.add_paragraph(text)


# ============================================================
# 内部写入函数(原 write_docx/write_xlsx/write_pdf/write_pptx 逻辑)
# ============================================================

def _write_docx(
    file_path: str,
    title: str = None,
    paragraphs=None,
) -> Dict[str, Any]:
    """写入Word文档(内部函数) - 小欧 2026-06-19 重构为3参数
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("docx"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_docx", "写入Word", file_path, detail="python-docx库未安装")
        return build_error(data={"error": "python-docx库未安装"}, llm_data=llm_data)

    try:
        from docx import Document

        doc = Document()

        if title:
            doc.add_heading(title, 0)

        if paragraphs is not None:
            dict_title, items = _resolve_paragraphs(paragraphs)
            if dict_title and not title:
                doc.add_heading(dict_title, 0)
            for item in items:
                _add_docx_content_item(doc, item)

        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(path)

        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("success", duration_ms, "write_docx", "写入Word", str(path))
        return build_success(data={"file_path": str(path)}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_docx", "写入Word", file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _write_xlsx(
    file_path: str,
    data: dict,
    sheet_name: str = "Sheet1"
) -> Dict[str, Any]:
    """写入Excel文件(内部函数) - 小健 2026-05-18
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("openpyxl"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_xlsx", "写入Excel", file_path, detail="openpyxl库未安装")
        return build_error(data={"error": "openpyxl库未安装"}, llm_data=llm_data)

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        
        if "headers" in data and data["headers"]:
            headers = data["headers"]
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal="center")
        
        if "rows" in data and data["rows"]:
            for row_idx, row_data in enumerate(data["rows"], 2):
                for col_idx, cell_data in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell_data)
        
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(path)
        
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("success", duration_ms, "write_xlsx", "写入Excel", str(path))
        return build_success(data={"file_path": str(path), "row_count": len(data.get("rows", []))}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_xlsx", "写入Excel", file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _write_pdf(
    file_path: str,
    title: str = None,
    paragraphs=None,
) -> Dict[str, Any]:
    """写入PDF文档(内部函数) — 小欧 2026-06-19 重构嵌套函数
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("reportlab"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_pdf", "写入PDF", file_path, detail="reportlab库未安装")
        return build_error(data={"error": "reportlab库未安装"}, llm_data=llm_data)

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
        from reportlab.lib import colors
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except ImportError:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_pdf", "写入PDF", file_path, detail="reportlab库导入失败")
        return build_error(data={"error": "reportlab库导入失败"}, llm_data=llm_data)

    def _add_content_item(elements, item, chinese_style, title_style):
        """处理单个PDF内容元素 — 小欧 2026-06-19"""
        if isinstance(item, str):
            if item.strip():
                elements.append(Paragraph(item, chinese_style))
                elements.append(Spacer(1, 3*mm))
        elif isinstance(item, dict):
            item_type = item.get("type", "paragraph")
            item_text = item.get("text", "")
            if item_type in ("h1", "heading"):
                heading_style = ParagraphStyle('h1', parent=chinese_style,
                    fontSize=18, spaceBefore=12, spaceAfter=6)
                elements.append(Paragraph(item_text, heading_style))
                elements.append(Spacer(1, 3*mm))
            elif item_type in ("h2", "h3", "h4", "h5"):
                fs = max(18 - int(item_type[1]) * 2, 12)
                h_style = ParagraphStyle(f'h{item_type[1]}', parent=chinese_style,
                    fontSize=fs, spaceBefore=8, spaceAfter=4)
                elements.append(Paragraph(item_text, h_style))
                elements.append(Spacer(1, 2*mm))
            elif item_type == "paragraph":
                elements.append(Paragraph(item_text, chinese_style))
                elements.append(Spacer(1, 3*mm))
            elif item_type == "table":
                rows_data = item.get("rows", [])
                if rows_data and len(rows_data) > 0:
                    try:
                        t = Table(rows_data)
                        t.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                            ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ]))
                        elements.append(t)
                        elements.append(Spacer(1, 5*mm))
                    except Exception:
                        pass
        else:
            text = str(item)
            if text.strip():
                elements.append(Paragraph(text, chinese_style))
                elements.append(Spacer(1, 3*mm))

    try:
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        doc = SimpleDocTemplate(str(path), pagesize=A4)
        styles = getSampleStyleSheet()

        try:
            font_path = "C:/Windows/Fonts/simsun.ttc"
            pdfmetrics.registerFont(TTFont('SimSun', font_path, subfontIndex=0))
            chinese_style = ParagraphStyle(
                'Chinese', parent=styles['Normal'],
                fontName='SimSun', fontSize=10, leading=14,
                wordWrap='CJK'
            )
            title_style = ParagraphStyle(
                'ChineseTitle', parent=styles['Title'],
                fontName='SimSun', fontSize=18, leading=24,
                wordWrap='CJK'
            )
        except Exception:
            chinese_style = styles['Normal']
            title_style = styles['Title']

        elements = []

        if title:
            elements.append(Paragraph(title, title_style))
            elements.append(Spacer(1, 10*mm))

        if paragraphs is not None:
            dict_title, items = _resolve_paragraphs(paragraphs)
            if dict_title and not title:
                elements.append(Paragraph(dict_title, title_style))
                elements.append(Spacer(1, 10*mm))
            for item in items:
                _add_content_item(elements, item, chinese_style, title_style)

        if not elements:
            elements.append(Paragraph(" ", chinese_style))

        doc.build(elements)

        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("success", duration_ms, "write_pdf", "写入PDF", str(path))
        return build_success(data={"file_path": str(path)}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_pdf", "写入PDF", file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


def _select_layout(prs, slide_type):
    """选布局 — 小欧 2026-06-19"""
    m = {0: 0, "cover": 0, 1: 1, "content": 1, 2: 2, "two": 2}
    return prs.slide_layouts[m.get(slide_type, 1)]


def _add_pptx_content(slide, content):
    """处理正文到 content placeholder — 小欧 2026-06-19, 小健 2026-06-20 DRY修复"""
    body = None
    for sh in slide.placeholders:
        idx = sh.placeholder_format.idx
        if idx == 1:
            continue
        if idx >= 2:
            body = sh.text_frame; break
    if body is None:
        return

    def _get_para(is_first: bool):
        return body.paragraphs[0] if is_first else body.add_paragraph()

    if isinstance(content, str):
        body.paragraphs[0].text = content
    elif isinstance(content, list):
        first_slot = True
        for item in content:
            if isinstance(item, str):
                p = _get_para(first_slot); first_slot = False
                p.text = item
            elif isinstance(item, dict):
                t = item.get("type", "paragraph")
                txt = item.get("text", "")
                if t == "paragraph":
                    p = _get_para(first_slot); first_slot = False
                    p.text = txt
                elif t == "bullets":
                    for b in item.get("items", []):
                        p = _get_para(first_slot); first_slot = False
                        p.text = str(b); p.level = 1


def _add_pptx_table(slide, table_data):
    """添加表格到幻灯片(独立shape) — 小欧 2026-06-19, 小健 2026-06-20 空检查修复"""
    if not table_data or not table_data[0] or len(table_data[0]) == 0:
        return
    rows, cols = len(table_data), len(table_data[0])
    from pptx.util import Inches
    left, top, width, height = Inches(1), Inches(2), Inches(6), Inches(0.4 * rows)
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for ri, row in enumerate(table_data):
        for ci, val in enumerate(row):
            tbl.cell(ri, ci).text = str(val)


def _add_pptx_slide(prs, slide_data):
    """添加一页幻灯片 — 小欧 2026-06-19"""
    slide_type = slide_data.get("type", 1)
    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    content = slide_data.get("content")
    tables = slide_data.get("tables")

    layout = _select_layout(prs, slide_type)
    slide = prs.slides.add_slide(layout)

    if title and slide.shapes.title:
        slide.shapes.title.text = title

    if subtitle and slide_type in (0, "cover"):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = subtitle; break

    if content:
        _add_pptx_content(slide, content)

    if tables:
        for td in tables:
            _add_pptx_table(slide, td)


def _build_pptx_presentation(slides: list):
    """构建全部幻灯片 — 小欧 2026-06-19（重写，无title参数）"""
    from pptx import Presentation
    prs = Presentation()

    if slides:
        for slide_data in slides:
            _add_pptx_slide(prs, slide_data)

    return prs


def _write_pptx(
    file_path: str,
    slides: list = None
) -> Dict[str, Any]:
    """写入PPT幻灯片(内部函数) - 小欧 2026-06-19 重构为2参数(无title)
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    if not _check_module("pptx"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_pptx", "写入PPT", file_path, detail="python-pptx库未安装")
        return build_error(data={"error": "python-pptx库未安装"}, llm_data=llm_data)

    try:
        prs = _build_pptx_presentation(slides)
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)

        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("success", duration_ms, "write_pptx", "写入PPT", str(path))
        return build_success(data={"file_path": str(path), "slide_count": len(prs.slides)}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_doc_llm_data("error", duration_ms, "write_pptx", "写入PPT", file_path, detail=str(e))
        return build_error(data={"error": str(e), "file_path": file_path}, llm_data=llm_data)


# ============================================================
# 路由函数 — LLM调用的统一入口
# ============================================================

def read_pdf(file_name: str) -> Dict[str, Any]:
    """读取PDF文件 — 小沈 2026-06-19"""
    check = _check_pdf_readable(file_name)
    if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
        return check
    result = _read_pdf(file_name, extract_images=True, extract_tables=True)
    return result


def read_docx(file_name: str) -> Dict[str, Any]:
    """读取Word文档 — 小沈 2026-06-19 加pandoc转.doc
    【2026-06-21 小健】builder改造，新3字段result
    """
    t0 = _time_mod.perf_counter()
    path = Path(file_name)
    suffix = path.suffix.lower()

    if suffix == ".doc":
        pandoc = shutil.which("pandoc")
        if not pandoc:
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_docx_llm_data("error", duration_ms, file_name, detail="读取.doc文件需要安装pandoc转换工具")
            return build_error(data={"file_name": file_name}, llm_data=llm_data)
        tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        tmp_path = tmp.name
        tmp.close()
        try:
            subprocess.run([pandoc, file_name, "-o", tmp_path], capture_output=True, text=True,
                         timeout=SUBPROCESS_TIMEOUT_LONG, check=True)
            check = _check_docx_readable(tmp_path)
            if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
                return check
            result = _read_docx(tmp_path, extract_tables=True)
            return result
        except subprocess.CalledProcessError as e:
            duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
            llm_data = _build_read_docx_llm_data("error", duration_ms, file_name, detail=f"pandoc转换.doc失败: {e.stderr}")
            return build_error(data={"stderr": e.stderr, "file_name": file_name}, llm_data=llm_data)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    check = _check_docx_readable(file_name)
    if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
        return check
    result = _read_docx(file_name, extract_tables=True)
    return result


def read_pptx(file_name: str) -> Dict[str, Any]:
    """读取PPT文件 — 小沈 2026-06-19"""
    result = _read_pptx(file_name, extract_notes=True)
    return result


def read_xlsx(file_name: str) -> Dict[str, Any]:
    """读取Excel/CSV文件 — 小沈 2026-06-19 修.xls检查"""
    path = Path(file_name)
    suffix = path.suffix.lower()

    if suffix == ".csv":
        result = _read_csv_stdlib(file_name, encoding="utf-8", delimiter=",", has_header=True, max_rows=10000)
    elif suffix == ".xls":
        result = _read_xls(file_name, max_rows=10000)
    else:
        check = _check_xlsx_readable(file_name)
        if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
            return check
        result = _read_xlsx(file_name, max_rows=10000)

    return result


def write_docx(
    file_name: str,
    title: Optional[str] = None,
    paragraphs: Optional[Union[str, List, Dict]] = None,
) -> Dict[str, Any]:
    """写入Word文档 — 小欧 2026-06-19 改3参数, 小健 2026-06-20 加coerce_json防御"""
    paragraphs = coerce_json(paragraphs)
    Path(file_name).parent.mkdir(parents=True, exist_ok=True)
    result = _write_docx(file_name, title=title, paragraphs=paragraphs)
    return result


def write_xlsx(
    file_name: str,
    data: Optional[Union[Dict[str, Any], List]] = None,
    sheet_name: str = "Sheet1",
) -> Dict[str, Any]:
    """写入Excel文件 — 小沈 2026-06-16, 小健 2026-06-20 加coerce_json防御"""
    data = coerce_json(data)
    Path(file_name).parent.mkdir(parents=True, exist_ok=True)
    if data is None:
        data = {"headers": [], "rows": []}
    elif isinstance(data, list):
        if len(data) > 0 and isinstance(data[0], list):
            first_row = data[0]
            data = {"headers": first_row, "rows": data[1:]}
        elif len(data) > 0 and isinstance(data[0], dict):
            headers = list(data[0].keys())
            rows = [list(row.values()) for row in data]
            data = {"headers": headers, "rows": rows}
        else:
            data = {"headers": [], "rows": data}
    elif isinstance(data, dict) and "headers" not in data and "rows" not in data:
        headers = list(data.keys())
        rows = [list(data.values())]
        data = {"headers": headers, "rows": rows}
    result = _write_xlsx(file_name, data=data, sheet_name=sheet_name)
    return result


def write_pdf(
    file_name: str,
    title: Optional[str] = None,
    paragraphs: Optional[Union[str, List, Dict]] = None,
) -> Dict[str, Any]:
    """写入PDF文件 — 小欧 2026-06-19 改3参数, 小健 2026-06-20 加coerce_json防御"""
    paragraphs = coerce_json(paragraphs)
    Path(file_name).parent.mkdir(parents=True, exist_ok=True)
    result = _write_pdf(file_name, title=title, paragraphs=paragraphs)
    return result


def write_pptx(
    file_name: str,
    slides: Optional[List[Dict]] = None,
) -> Dict[str, Any]:
    """写入PPT文件 — 小欧 2026-06-19 改2参数(无title), 小健 2026-06-20 加coerce_json防御"""
    slides = coerce_json(slides)
    Path(file_name).parent.mkdir(parents=True, exist_ok=True)
    result = _write_pptx(file_name, slides=slides)
    return result

