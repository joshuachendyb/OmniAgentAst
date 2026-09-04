# -*- coding: utf-8 -*-
# 编辑历史:
# 2026-08-21 - 小欧 - 11.6.1: success分支调 with_artifact_file 声明产出物
"""
D8: write_pptx — 写入PPT文档

从document_tools.py拆分而来 — 小欧 2026-06-22
内聚: _select_layout / _add_pptx_content / _add_pptx_table / _add_pptx_slide / _build_pptx_presentation 辅助函数
"""
# 【铁规1】helper/被调函数(以下划线_开头的函数)只返回raw dict，严禁调用build_success/build_error/build_warning和构建llm_data。
# build3+llm_data只能在tool的main函数(对外公开的函数)中包装。违反此规则的代码视为不合规。
# 【铁规2】工具返回原始data，禁止调用truncate_data_for_frontend。截断只能在前端yield层。
# 【铁规3】计时(duration_ms计算)只能在tool的主函数中，严禁在子函数/helper中计时。
# 2026-07-31 - 小欧 - CRITICAL: _add_pptx_table 中 Inches 导入在使用之后(NameError崩溃)。移动 import 到函数顶部，确保 Inches(2) fallback 可用
# 2026-07-31 - 小欧 - Bug④/⑥/⑭修复: 正文占位符按类型(BODY/OBJECT/VERTICAL_BODY)选择, 防封面SUBTITLE被content覆盖; falsy-zero改is not None; 表高超版心时跳过防负高度
# 2026-08-13 - 小欧 - A5职责拆分: permission_error_hint/hint_for_write_error 导入源从 app.tools.validate.file_path_checker 改为 app.tools.toolhelper.error_hints

import time as _time_mod
from pathlib import Path
from typing import Any, Dict, List, Optional

from app.logger import logger
from app.tools.tool_response import build_success, build_error, with_artifact_file
from app.tools.tool_fc_helper import _check_module
from app.tools.validate.file_type_checker import check_office_file
from app.tools.validate.file_safety_checker import check_content_safety
from app.tools.tool_constants import ERR_DOC_WRITE_PPTX  # 2026-07-31 小欧: 移除未使用 ERR_DOC_NO_PPTX
from app.tools.toolhelper.error_hints import permission_error_hint, hint_for_write_error
from app.utils.json_utils import coerce_json  # 2026-07-31 小欧: 移除未使用 logger
from app.utils.table_helper import calculate_column_widths, get_table_header_style_config, dict_table_to_rows
from app.tools.document.md_inline_utils import _parse_inline_md


def _set_pptx_paragraph_text(paragraph, text):
    """设置段落文本，支持行内Markdown格式 — 小欧 2026-07-08"""
    from pptx.util import Pt
    paragraph.clear()
    for seg_text, bold, italic, code, link_url in _parse_inline_md(text):
        run = paragraph.add_run()
        run.text = seg_text
        run.font.bold = bold
        run.font.italic = italic
        if code:
            run.font.name = 'Courier New'
            run.font.size = Pt(9)


def _build_write_pptx_llm_data(
    exec_code: str, duration_ms: int,
    file_path: str = "", slide_count: int = 0, detail: str = "", hint: str = "",
) -> Dict[str, Any]:
    """write_pptx的llm_data构建函数 — 小欧 2026-06-22 — 小欧 2026-07-05 新增hint参数"""
    if exec_code == "error":
        return {
            "summary": f"写入PPT{file_path}，失败: {detail}",
            "action": {"tool": "write_pptx", "tool_zh": "写入PPT", "target": file_path, "params": {"file_path": file_path}},
            "status": {"exec_code": "error", "message": "写入PPT失败", "code": ERR_DOC_WRITE_PPTX, "detail": detail, "hint": hint if hint else "请检查路径和权限"},
            "duration_ms": duration_ms,
            "metrics": {},
        }
    return {
        "summary": f"写入PPT{file_path}，成功: {slide_count}页",
        "action": {"tool": "write_pptx", "tool_zh": "写入PPT", "target": file_path, "params": {"file_path": file_path}},
        "status": {"exec_code": "success", "message": "写入PPT成功", "code": "", "detail": "", "hint": ""},
        "duration_ms": duration_ms,
        "metrics": {
            "slide_count": {"value": slide_count, "text": f"{slide_count}页"},
        },
    }


def _select_layout(prs, slide_type):
    """选布局 — 小欧 2026-06-19"""
    m = {0: 0, "cover": 0, 1: 1, "content": 1, 2: 2, "two": 2}
    return prs.slide_layouts[m.get(slide_type, 1)]


def _add_pptx_content(slide, content):
    """处理正文到 content placeholder — 小欧 2026-06-19; 小健 2026-06-24 修复跳过idx=1的bug
       2026-07-31 小欧: Bug④修复 — 按占位符类型(BODY/OBJECT/VERTICAL_BODY)选择正文占位符,
       避免封面(Title Slide)布局中 idx==1 的 SUBTITLE 被 content 覆盖, 也避免误写进 DATE/FOOTER 等占位符
    """
    from pptx.enum.shapes import PP_PLACEHOLDER
    body = None
    for sh in slide.placeholders:
        ptype = sh.placeholder_format.type
        if ptype in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.VERTICAL_BODY):
            body = sh.text_frame
            break
    if body is None:
        logger.warning("[write_pptx] 当前布局无正文占位符, content被忽略(封面/标题类布局)")
        return

    def _get_para(is_first: bool):
        return body.paragraphs[0] if is_first else body.add_paragraph()

    if isinstance(content, str):
        _set_pptx_paragraph_text(body.paragraphs[0], content)
    elif isinstance(content, list):
        first_slot = True
        for item in content:
            if isinstance(item, str):
                p = _get_para(first_slot)
                first_slot = False
                _set_pptx_paragraph_text(p, item)
            elif isinstance(item, dict):
                t = item.get("type", "paragraph")
                txt = item.get("text", "")
                if t == "paragraph":
                    p = _get_para(first_slot)
                    first_slot = False
                    _set_pptx_paragraph_text(p, txt)
                elif t == "bullets":
                    for b in item.get("items", []):
                        p = _get_para(first_slot)
                        first_slot = False
                        _set_pptx_paragraph_text(p, str(b))
                        p.level = 1
    elif isinstance(content, dict):
        t = content.get("type", "paragraph")
        items = content.get("items", [])
        first_slot = True
        if t == "bullets":
            for b in items:
                p = _get_para(first_slot)
                first_slot = False
                _set_pptx_paragraph_text(p, str(b))
                p.level = 1
        elif t == "paragraph":
            txt = content.get("text", "")
            p = _get_para(first_slot)
            _set_pptx_paragraph_text(p, txt)


def _add_pptx_table(slide, table_data, start_top=None):
    """添加表格到幻灯片(独立shape) — 小欧 2026-06-19; 小健 2026-06-24 修复多表格重叠和超出边界、列宽自适应、表头样式"""
    from pptx.util import Inches, Emu, Pt  # 2026-07-31 小欧: 修移至函数顶部(Inches 原在行130, 但行128先使用导致NameError)
    from pptx.dml.color import RGBColor
    if not table_data or not table_data[0] or len(table_data[0]) == 0:
        return start_top if start_top is not None else Inches(2)  # 2026-07-31 小欧: Bug⑥ falsy-zero → is not None
    
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(1)
    top = start_top if start_top is not None else Inches(2)  # 2026-07-31 小欧: Bug⑥ falsy-zero → is not None
    width = Inches(6)
    
    slide_height = Inches(7.5)
    max_height = int(slide_height - top - Inches(0.5))
    # 2026-07-31 小欧: Bug⑭ — 上方表格已把 top 推出版心时 max_height 为负, add_table 负高度崩溃/损坏XML。
    # 可用高度不足0.5英寸时跳过该表格(保留当前 top, 供后续表格继续判断)
    if max_height < int(Inches(0.5)):
        logger.warning("[write_pptx] 表格超出版心范围, 跳过该表格")
        return Emu(int(top))
    row_height = min(int(Inches(0.4)), max_height // rows)
    height = Emu(row_height * rows)
    
    if height > max_height:
        height = Emu(max_height)
        row_height = max_height // rows
    
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    col_widths = calculate_column_widths(table_data, total_width=6.0)
    total_width_emu = int(width)
    for ci, w in enumerate(col_widths):
        col_width_emu = int(total_width_emu * w / 6.0)
        tbl.columns[ci].width = Emu(col_width_emu)
    
    header_config = get_table_header_style_config()
    for ri, row in enumerate(table_data):
        for ci, val in enumerate(row):
            if ci < cols:
                cell = tbl.cell(ri, ci)
                text = str(val)
                if '\n' in text:
                    lines = text.split('\n')
                    _set_pptx_paragraph_text(cell.text_frame.paragraphs[0], lines[0])
                    for line in lines[1:]:
                        _set_pptx_paragraph_text(cell.text_frame.add_paragraph(), line)
                else:
                    _set_pptx_paragraph_text(cell.text_frame.paragraphs[0], text)
                
                if ri == 0:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.bold = header_config["bold"]
                            run.font.size = Pt(header_config["font_size"])
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
    
    return Emu(int(top) + int(height) + int(Inches(0.3)))


def _dict_table_to_rows(dict_table):
    """把dict型表格{headers,rows}转成list[list] — 小欧 2026-07-08"""
    return dict_table_to_rows(dict_table)


def _extract_tables_from_content(content):
    """从content提取表格数据，返回(text_to_render, extracted_tables) — 小欧 2026-07-08

    覆盖C1~C11全部11种content结构:
      C1 str / C2 list[str] / C3 list[dict] / C4 dict → 原样返回
      C5 list[list]纯二维数组 → 全转表格
      C6 混合list(str+dict+list) → 拆出list行合并成表
      C7 dict type=table → 转表格
      C8 list[dict]含type=table → 拆出转表
      C9 [] / C10 None → 跳过
      C11 含None元素 → 过滤None
    """
    if content is None:
        return None, []
    if isinstance(content, str):
        return content, []
    if isinstance(content, dict):
        if content.get("type") == "table":
            rows = _dict_table_to_rows(content)
            return None, [rows] if rows else []
        return content, []
    if isinstance(content, list):
        if not content:
            return None, []
        has_table_item = any(
            isinstance(item, list)
            or (isinstance(item, dict) and item.get("type") == "table")
            for item in content
        )
        if not has_table_item:
            return content, []
        text_parts = []
        extracted = []
        current_table = []
        for item in content:
            if item is None:
                continue
            if isinstance(item, list):
                current_table.append(item)
            elif isinstance(item, dict) and item.get("type") == "table":
                if current_table:
                    extracted.append(current_table)
                    current_table = []
                rows = _dict_table_to_rows(item)
                if rows:
                    extracted.append(rows)
            else:
                if current_table:
                    extracted.append(current_table)
                    current_table = []
                text_parts.append(item)
        if current_table:
            extracted.append(current_table)
        return text_parts if text_parts else None, extracted
    return None, []


def _normalize_tables(tables):
    """归一化tables为list[list[list]]标准格式 — 小欧 2026-07-08

    覆盖T1~T6全部6种tables结构:
      T1 list[list[list]] → 原样
      T2 list[list]少包一层 → 自动包成[list]
      T3 dict{headers,rows} → 转[list[list]]
      T4 list[dict{headers,rows}] → 逐个转
      T5 [] / T6 None → 返回[]
    """
    if not tables:
        return []
    if isinstance(tables, dict):
        rows = _dict_table_to_rows(tables)
        return [rows] if rows else []
    if isinstance(tables, list):
        if not tables:
            return []
        first = tables[0]
        if isinstance(first, list):
            if not first:
                return []
            if isinstance(first[0], list):
                return tables
            return [tables]
        if isinstance(first, dict):
            result = []
            for td in tables:
                rows = _dict_table_to_rows(td)
                if rows:
                    result.append(rows)
            return result
    return []


def _normalize_text(value, default=""):
    """文本标准化 — 小健 2026-06-24"""
    if isinstance(value, str):
        return value
    return str(value) if value is not None else default


def _add_pptx_slide(prs, slide_data):
    """添加一页幻灯片 — 小欧 2026-06-19; 小健 2026-06-24 增加参数验证"""
    if not isinstance(slide_data, dict):
        return  # 跳过无效slide
    
    slide_type = slide_data.get("type", 1)
    title = _normalize_text(slide_data.get("title", ""))
    subtitle = _normalize_text(slide_data.get("subtitle", ""))
    content = slide_data.get("content")
    tables = slide_data.get("tables")
    
    layout = _select_layout(prs, slide_type)
    slide = prs.slides.add_slide(layout)

    if title and slide.shapes.title:
        slide.shapes.title.text = title

    if subtitle and slide_type in (0, "cover"):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = subtitle
                break

    # === 归一化 content + tables，覆盖 C1~C11 / T1~T6 / G1~G6 === — 小欧 2026-07-08
    text_content, table_list = _extract_tables_from_content(content)
    table_list.extend(_normalize_tables(tables))

    if text_content is not None:
        _add_pptx_content(slide, text_content)

    if table_list:
        from pptx.util import Inches
        table_top = Inches(2)
        for td in table_list:
            table_top = _add_pptx_table(slide, td, table_top)


def _build_pptx_presentation(slides: list):
    """构建全部幻灯片 — 小欧 2026-06-19"""
    from pptx import Presentation
    prs = Presentation()

    if slides:
        for slide_data in slides:
            _add_pptx_slide(prs, slide_data)

    return prs


def write_pptx(
    path: str,
    slides: Optional[List[Dict]] = None,
) -> Dict[str, Any]:
    """写入PPT文件 — 小欧 2026-06-19 — 小欧 2026-06-22 独立文件 — 小欧 2026-06-24 增加文件类型前置检查"""
    t0 = _time_mod.perf_counter()

    # 文件类型前置检查（含路径检查+类型检查+模块安全检查）— 北京老陈 2026-07-09
    is_valid, error_detail, hint = check_office_file(path, allow_create=True)
    if not is_valid:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_pptx_llm_data("error", duration_ms, path, detail=error_detail, hint=hint)
        return build_error(data={}, llm_data=llm_data)

    slides = coerce_json(slides)

    cs_error, safe_slides = check_content_safety(slides, "pptx", param_name="slides")
    if cs_error:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_pptx_llm_data("error", duration_ms, path, detail=cs_error, hint=f"请检查slides参数(当前类型: {type(slides).__name__})")
        return build_error(data={}, llm_data=llm_data)
    slides = safe_slides

    if not _check_module("pptx"):
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_pptx_llm_data("error", duration_ms, path, detail="python-pptx库未安装", hint="请安装python-pptx库")
        return build_error(data={}, llm_data=llm_data)

    try:
        prs = _build_pptx_presentation(slides)
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)

        slide_count = len(prs.slides)
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        llm_data = _build_write_pptx_llm_data("success", duration_ms, str(path), slide_count)
        with_artifact_file(llm_data, str(path))   # 11.6.1 产出物声明 — 小欧 2026-08-21
        # =============================================================================
        # 数据设计：slide_count/file_path 从 data 移除，通过 llm_data.metrics/summary
        # 传入 LLM observation。summary 已包含文件路径和页数：
        #   "写入PPT成功: /path.pptx, 5页"
        # data 为空 dict 时 formatter 不追加详情，避免冗余。
        # — 小欧 2026-07-06 18:46:13
        # =============================================================================
        return build_success(data={}, llm_data=llm_data)
    except PermissionError as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        hint = permission_error_hint(path)
        # 小欧 2026-07-12: 此处path经Path()重赋值为WindowsPath,须str()化后传入builder,
        # 避免action.target持有Path对象触发观察格式化len()崩溃
        llm_data = _build_write_pptx_llm_data("error", duration_ms, str(path), detail=str(e), hint=hint)
        return build_error(data={}, llm_data=llm_data)
    except Exception as e:
        duration_ms = int((_time_mod.perf_counter() - t0) * 1000)
        hint = hint_for_write_error(e, path)
        # 小欧 2026-07-12: 此处path经Path()重赋值为WindowsPath,须str()化后传入builder,
        # 避免action.target持有Path对象触发观察格式化len()崩溃
        llm_data = _build_write_pptx_llm_data("error", duration_ms, str(path), detail=str(e), hint=hint)
        return build_error(data={}, llm_data=llm_data)

