# -*- coding: utf-8 -*-
"""
文档读写工具函数模块

【创建时间】2026-05-02 小沈
【设计依据】按文档第8.3节 Tool 80-82 定义
【重构 2026-05-18 小健】8个旧函数抽取为内部函数，新增read_document/write_document路由函数

【重要】新函数增加规范 - 小沈 2026-05-04
新增函数时必须同步修改以下3个文件：
1. *_tools.py: 函数实现（必须有详细注释）
2. *_schema.py: Pydantic 模型（输入参数定义）
3. *_register.py: 显式注册（description + examples + input_model）

包含：
- _read_pdf / _read_docx / _read_xlsx / _read_pptx: 内部读取函数
- _write_docx / _write_xlsx / _write_pdf / _write_pptx: 内部写入函数
- read_document: 统一读取路由（按后缀自动选择解析器）
- write_document: 统一写入路由（按后缀自动选择写入器）
- convert_document: 文档格式转换

Author: 小沈 - 2026-05-02
【新增 2026-05-05 小沈】write_pdf, convert_document
【重构 2026-05-18 小健】8合2路由重构
"""

import csv
import json
import os
import tempfile
from typing import Dict, Any, List, Optional, Literal, Union, Tuple
from pathlib import Path

from app.services.tools.document.document_schema import (
    ReadDocumentInput,
    WriteDocumentInput,
)
from app.services.tools.tool_result_utils import build_next_actions
from app.services.tools._response import build_success, build_error, build_warning
from app.services.tools.toolhelper.common_helper import _check_module
from app.services.tools.toolhelper.data_helper import _serialize_rows





def _check_pdf_readable(file_path: str) -> Dict[str, Any]:
    """检查PDF文件是否可读（内部helper） - 小沈 2026-05-18，从env_check_tools.py迁入"""
    path = Path(file_path)
    if not path.exists():
        return build_error(ERR_DOC_READ_PDF, f"文件不存在: {file_path}")
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
        return build_success({"readable": True, "page_count": page_count},
                f"PDF文件可读，共 {page_count} 页")
    except ImportError:
        return build_error(ERR_NO_PDFPLUMBER,
                "pdfplumber库未安装，请先执行: pip install pdfplumber")
    except Exception as e:
        return build_error(ERR_DOC_READ_PDF, f"PDF文件不可读: {str(e)}")


def _check_docx_readable(file_path: str) -> Dict[str, Any]:
    """检查Word文件是否可读（内部helper） - 小沈 2026-05-18，从env_check_tools.py迁入"""
    path = Path(file_path)
    if not path.exists():
        return build_error(ERR_DOC_READ_DOCX, f"文件不存在: {file_path}")
    try:
        import docx
        doc = docx.Document(path)
        para_count = len(doc.paragraphs)
        return build_success({"readable": True, "paragraph_count": para_count},
                f"Word文件可读，共 {para_count} 段")
    except ImportError:
        return build_error(ERR_NO_DOCX,
                "python-docx库未安装，请先执行: pip install python-docx")
    except Exception as e:
        return build_error(ERR_DOC_READ_DOCX, f"Word文件不可读: {str(e)}")


def _check_xlsx_readable(file_path: str) -> Dict[str, Any]:
    """检查Excel文件是否可读（内部helper） - 小沈 2026-05-18，从env_check_tools.py迁入"""
    path = Path(file_path)
    if not path.exists():
        return build_error(ERR_DOC_READ_XLSX, f"文件不存在: {file_path}")
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True)
        sheet_names = wb.sheetnames
        wb.close()
        return build_success({"readable": True, "sheet_names": sheet_names},
                f"Excel文件可读，共 {len(sheet_names)} 个工作表")
    except ImportError:
        return build_error(ERR_DOC_NO_OPENPYXL,
                "openpyxl库未安装，请先执行: pip install openpyxl")
    except Exception as e:
        return build_error(ERR_DOC_READ_XLSX, f"Excel文件不可读: {str(e)}")


def _validate_csv_format(file_path: str) -> Dict[str, Any]:
    """验证CSV文件格式（内部helper） - 小沈 2026-05-18，从env_check_tools.py迁入"""
    path = Path(file_path)
    if not path.exists():
        return build_error(ERR_DOC_READ_CSV, f"文件不存在: {file_path}")
    if not path.suffix.lower() in (".csv", ".tsv", ".txt"):
        return build_error(ERR_DOC_READ_CSV, "文件扩展名不是CSV格式")

    errors = []
    try:
        with open(path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            row_count = 0
            col_count = None
            for i, row in enumerate(reader):
                row_count += 1
                if col_count is None:
                    col_count = len(row)
                elif len(row) != col_count:
                    if len(row) != 0:
                        errors.append(f"第{i+1}行列数不一致: 期望{col_count}列，实际{len(row)}列")
                if row_count > 1000:
                    break
    except UnicodeDecodeError:
        try:
            with open(path, "r", encoding="gbk", newline="") as f:
                reader = csv.reader(f)
                row_count = 0
                for _ in reader:
                    row_count += 1
                    if row_count > 1000:
                        break
        except Exception as e:
            errors.append(f"文件编码错误: {str(e)}")
    except Exception as e:
        errors.append(f"文件读取错误: {str(e)}")

    is_valid = len(errors) == 0
    if is_valid:
        return build_success({"valid": True}, "CSV格式正确")
    else:
        return build_error(ERR_DOC_READ_CSV, f"CSV格式有 {len(errors)} 个问题",
                data={"valid": False, "errors": errors})


def _validate_chart_data(data: Dict[str, Any]) -> Dict[str, Any]:
    """验证图表数据格式（内部helper） - 小沈 2026-05-18，从env_check_tools.py迁入"""
    errors = []
    if not isinstance(data, dict):
        errors.append("data必须是字典类型")
    else:
        if "labels" not in data:
            errors.append("缺少labels字段")
        elif not isinstance(data["labels"], list):
            errors.append("labels必须是数组类型")
        if "values" not in data:
            errors.append("缺少values字段")
        elif not isinstance(data["values"], list):
            errors.append("values必须是数组类型")
        if "labels" in data and "values" in data:
            if isinstance(data["labels"], list) and isinstance(data["values"], list):
                if len(data["labels"]) != len(data["values"]):
                    errors.append(f"labels和values长度不一致: labels={len(data['labels'])}, values={len(data['values'])}")
    is_valid = len(errors) == 0
    if is_valid:
        return build_success({"valid": True}, "图表数据格式正确")
    else:
        return build_error(ERR_DOC_CHART_GENERATE, f"图表数据格式有 {len(errors)} 个问题",
                data={"valid": False, "errors": errors})


def _read_csv_pandas(
    file_path: str,
    encoding: str = "utf-8",
    delimiter: str = ",",
    has_header: bool = True,
    max_rows: int = 1000,
) -> Dict[str, Any]:
    """使用pandas读取CSV文件 - 小沈 2026-05-18（从data_analysis迁入）"""
    if not _check_module("pandas"):
        return build_error(ERR_NO_PANDAS, "pandas库未安装，请先执行: pip install pandas")
    try:
        import pandas as pd
        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_READ_CSV_DATAFRAME, f"文件不存在: {file_path}")
        header = 0 if has_header else None
        df = pd.read_csv(path, encoding=encoding, delimiter=delimiter, header=header, nrows=max_rows)
        columns = df.columns.tolist()
        serialized_rows = _serialize_rows(df)
        dtypes = {col: str(dtype) for col, dtype in df.dtypes.items()}
        return build_success({"columns": columns, "rows": serialized_rows, "row_count": len(serialized_rows), "dtypes": dtypes}, f"成功读取CSV文件: {file_path}，共 {len(serialized_rows)} 行数据")
    except Exception as e:
        return build_error(ERR_READ_CSV_DATAFRAME, f"读取CSV文件失败: {str(e)}")


def _read_excel_pandas(
    file_path: str,
    sheet_name: Optional[str] = None,
    max_rows: int = 1000,
) -> Dict[str, Any]:
    """使用pandas读取Excel文件 - 小沈 2026-05-18（从data_analysis迁入）"""
    if not _check_module("pandas"):
        return build_error(ERR_NO_PANDAS, "pandas库未安装，请先执行: pip install pandas openpyxl")
    if not _check_module("openpyxl"):
        return build_error(ERR_DOC_NO_OPENPYXL, "openpyxl库未安装，请先执行: pip install openpyxl")
    try:
        import pandas as pd
        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_READ_EXCEL_DATAFRAME, f"文件不存在: {file_path}")
        df = pd.read_excel(path, sheet_name=sheet_name if sheet_name else 0, nrows=max_rows, engine="openpyxl")
        columns = df.columns.tolist()
        serialized_rows = _serialize_rows(df)
        dtypes = {col: str(dtype) for col, dtype in df.dtypes.items()}
        actual_sheet = sheet_name if sheet_name else "Sheet1"
        return build_success({"columns": columns, "rows": serialized_rows, "row_count": len(serialized_rows), "dtypes": dtypes, "sheet_name": actual_sheet}, f"成功读取Excel文件: {file_path}，共 {len(serialized_rows)} 行数据")
    except Exception as e:
        return build_error(ERR_READ_EXCEL_DATAFRAME, f"读取Excel文件失败: {str(e)}")


def _parse_pages(pages_str: str) -> List[int]:
    """解析页码字符串为页码列表 - 小沈 2026-05-02"""
    result = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            try:
                result.extend(range(int(start), int(end) + 1))
            except ValueError:
                pass
        else:
            try:
                result.append(int(part))
            except ValueError:
                pass
    return sorted(set(result))


# ============================================================
# 内部读取函数（原 read_pdf/read_docx/read_xlsx/read_pptx 逻辑）
# ============================================================

def _process_page(page, page_num: int,
                  extract_tables: bool, extract_images: bool
                  ) -> Tuple[str, List[Dict], List[Dict]]:
    """处理单页PDF：返回 (text, tables_data, images_data)

    小沈 2026-05-25 重构拆分
    """
    text = page.extract_text() or ""

    tables = []
    if extract_tables:
        for idx, t in enumerate(page.extract_tables() or []):
            tables.append({"page": page_num, "table_idx": idx, "data": t})

    images = []
    if extract_images:
        for idx, img in enumerate(page.images or []):
            images.append({
                "page": page_num, "image_idx": idx,
                "x0": float(img.get("x0", 0)), "y0": float(img.get("y0", 0)),
                "x1": float(img.get("x1", 0)), "y1": float(img.get("y1", 0)),
                "width": float(img.get("width", 0)), "height": float(img.get("height", 0)),
            })

    return text, tables, images


def _build_pdf_result(fp: str, pages_read: List[int], all_text: List[str],
                      tables_data: List[Dict], images_data: List[Dict],
                      page_count: int) -> dict:
    """构建统一的 PDF 读取返回结果

    小沈 2026-05-25 重构拆分
    """
    full_text = "\n\n".join(all_text)
    result = {"text": full_text, "page_count": page_count, "pages_read": pages_read}
    _llm = {"文件": fp, "页数": f"{page_count}页(读取{len(pages_read)}页)",
            "文本长度": f"{len(full_text)}字符", "内容": full_text}
    msg = f"成功读取PDF文件: {fp}，共读取 {len(pages_read)} 页"
    if tables_data:
        result["tables"] = tables_data
        result["table_count"] = len(tables_data)
        _llm["表格数"] = len(tables_data)
        msg += f"，{len(tables_data)} 个表格"
    if images_data:
        result["images"] = images_data
        result["image_count"] = len(images_data)
        _llm["图片数"] = len(images_data)
        msg += f"，{len(images_data)} 张图片"
    return build_success(result, msg, llm_data=_llm)


def _read_pdf(
    file_path: str,
    pages: str = None,
    extract_images: bool = False,
    extract_tables: bool = False
) -> Dict[str, Any]:
    """读取PDF文件并提取文本内容（内部函数） - 小健 2026-05-18
    【2026-05-25 小沈重构】拆分：逐页提取 → _process_page，结果构建 → _build_pdf_result
    """
    if not _check_module("pdfplumber"):
        return build_error(ERR_NO_PDFPLUMBER,
                "pdfplumber库未安装，请先执行: pip install pdfplumber")

    try:
        import pdfplumber

        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_DOC_READ_PDF, f"文件不存在: {file_path}")

        all_text, pages_read, tables_data, images_data = [], [], [], []
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            target = _parse_pages(pages) if pages else list(range(1, page_count + 1))
            target = [p for p in target if 1 <= p <= page_count]

            for pn in target:
                page = pdf.pages[pn - 1]
                text, tables, images = _process_page(page, pn, extract_tables, extract_images)
                all_text.append(f"--- 第 {pn} 页 ---\n{text}")
                pages_read.append(pn)
                tables_data.extend(tables)
                images_data.extend(images)

        return _build_pdf_result(file_path, pages_read, all_text, tables_data, images_data, page_count)
    except Exception as e:
        return build_error(ERR_DOC_READ_PDF, f"读取PDF文件失败: {str(e)}")


def _read_docx(
    file_path: str,
    extract_tables: bool = False
) -> Dict[str, Any]:
    """读取Word文档并提取文本内容（内部函数） - 小健 2026-05-18"""
    if not _check_module("docx"):
        return build_error(ERR_NO_DOCX,
                "python-docx库未安装，请先执行: pip install python-docx")

    try:
        import docx

        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_DOC_READ_DOCX, f"文件不存在: {file_path}")

        doc = docx.Document(path)
        paragraphs = [para.text for para in doc.paragraphs]
        text = "\n".join(paragraphs)
        
        result_data = {
            "text": text,
            "paragraph_count": len(paragraphs),
        }
        
        if extract_tables:
            tables_data = []
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_rows.append(row_data)
                tables_data.append(table_rows)
            result_data["tables"] = tables_data
            result_data["table_count"] = len(tables_data)
        
        _llm = {
            "文件": file_path,
            "段落数": len(paragraphs),
            "文本长度": f"{len(text)}字符",
            "内容": text,
        }
        return build_success(result_data, f"成功读取Word文档: {file_path}，共 {len(paragraphs)} 段",
                llm_data=_llm)
    except Exception as e:
        return build_error(ERR_DOC_READ_DOCX, f"读取Word文档失败: {str(e)}")


def _read_xlsx(
    file_path: str,
    sheet_name: str = None,
    max_rows: int = 1000,
    header: bool = True
) -> Dict[str, Any]:
    """读取Excel文件并提取表格数据（内部函数） - 小健 2026-05-18"""
    if not _check_module("openpyxl"):
        return build_error(ERR_DOC_NO_OPENPYXL,
                "openpyxl库未安装，请先执行: pip install openpyxl")

    try:
        from openpyxl import load_workbook

        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_DOC_READ_XLSX, f"文件不存在: {file_path}")

        wb = load_workbook(path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames

        if sheet_name:
            if sheet_name not in sheet_names:
                wb.close()
                return build_error(ERR_DOC_READ_XLSX, f"工作表不存在: {sheet_name}，可用工作表: {sheet_names}")
            ws = wb[sheet_name]
        else:
            ws = wb.active

        rows = []
        headers = []
        row_count = 0

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows + (1 if header else 0):
                break
            row_data = [
                None if val is None else val
                for val in row
            ]
            if i == 0 and header:
                headers = [str(h) if h is not None else f"column_{j}" for j, h in enumerate(row_data)]
            else:
                if i == 0 and not header:
                    headers = [f"column_{j}" for j in range(len(row_data))]
                rows.append(row_data)
                row_count += 1

        wb.close()

        return build_success({
                "headers": headers,
                "rows": rows,
                "row_count": row_count,
                "sheet_names": sheet_names,
            }, f"成功读取Excel文件: {file_path}，工作表: {ws.title}，共 {row_count} 行数据")
    except Exception as e:
        return build_error(ERR_DOC_READ_XLSX, f"读取Excel文件失败: {str(e)}")


def _read_csv_stdlib(
    file_path: str,
    encoding: str = "utf-8",
    delimiter: str = ",",
    has_header: bool = True,
    max_rows: int = 1000
) -> Dict[str, Any]:
    """使用标准库csv读取CSV文件（内部函数）— 小健 2026-05-18"""
    import csv
    
    try:
        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_DOC_READ_CSV, f"文件不存在: {file_path}")
        
        rows = []
        columns = []
        # 小健 2026-05-19: 多编码尝试回退(utf-8→gbk→gb2312→latin-1)
        encodings_to_try = [encoding, "gbk", "gb2312", "latin-1"] if encoding == "utf-8" else [encoding, "utf-8", "latin-1"]
        read_ok = False
        for enc in encodings_to_try:
            try:
                with open(path, "r", encoding=enc, newline="") as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    for i, row in enumerate(reader):
                        if i >= max_rows:
                            break
                        if i == 0:
                            if has_header:
                                columns = row
                            else:
                                columns = [f"col_{j}" for j in range(len(row))]
                                rows.append(row)
                        else:
                            rows.append(row)
                read_ok = True
                break
            except UnicodeDecodeError:
                continue
        if not read_ok:
            return build_error(ERR_DOC_READ_CSV, f"读取CSV文件失败: 编码不匹配(尝试了{encodings_to_try})")
        
        return build_success({
                "columns": columns,
                "rows": rows,
                "row_count": len(rows),
            }, f"成功读取CSV文件: {file_path}，共 {len(rows)} 行数据")
    except Exception as e:
        return build_error(ERR_DOC_READ_CSV, f"读取CSV文件失败: {str(e)}")


def _auto_convert_to_pdf(input_path: str, suffix: str) -> Optional[str]:
    """自动将旧版Office格式转换为PDF，返回PDF路径或None"""
    try:
        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf_path = tmp.name
        tmp.close()
        result = convert_document(input_path, output_path=pdf_path)
        if result.get("code") != "SUCCESS":
            try:
                os.unlink(pdf_path)
            except:
                pass
            return None
        return pdf_path
    except Exception:
        return None


def _cleanup_temp_pdf(pdf_path: str) -> None:
    """清理临时PDF文件"""
    try:
        if os.path.isfile(pdf_path):
            os.unlink(pdf_path)
    except:
        pass


def _read_pptx(
    file_path: str,
    extract_notes: bool = False
) -> Dict[str, Any]:
    """读取PPT幻灯片（内部函数） - 小健 2026-05-18"""
    if not _check_module("pptx"):
        return build_error(ERR_DOC_NO_PPTX,
                "python-pptx库未安装，请先执行: pip install python-pptx")

    try:
        from pptx import Presentation

        path = Path(file_path)
        if not path.exists():
            return build_error(ERR_DOC_READ_PPTX, f"文件不存在: {file_path}")

        prs = Presentation(path)
        slides_data = []
        notes_data = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
            
            slides_data.append({
                "slide_num": slide_num,
                "text": "\n".join(slide_text)
            })
            
            if extract_notes and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_data.append({
                        "slide_num": slide_num,
                        "notes": notes
                    })

        result_data = {
            "slide_count": len(prs.slides),
            "slides": slides_data,
        }

        if extract_notes and notes_data:
            result_data["notes"] = notes_data

        _total_text = sum(len(s.get("text", "")) for s in slides_data)
        _llm = {
            "文件": file_path,
            "幻灯片数": len(prs.slides),
            "文本总长度": f"{_total_text}字符",
            "幻灯片内容": [{"页码": s["slide_num"], "文本": s.get("text", "")} for s in slides_data],
        }
        if extract_notes and notes_data:
            _llm["备注数"] = len(notes_data)

        return build_success(result_data, f"成功读取PPT文件: {file_path}，共 {len(prs.slides)} 页",
                llm_data=_llm)
    except Exception as e:
        return build_error(ERR_DOC_READ_PPTX, f"读取PPT文件失败: {str(e)}")


# ============================================================
# 内部写入函数（原 write_docx/write_xlsx/write_pdf/write_pptx 逻辑）
# ============================================================

def _write_docx(
    file_path: str,
    content: str = None,
    paragraphs: list = None,
    title: str = None,
    table_data: list = None
) -> Dict[str, Any]:
    """写入Word文档（内部函数） - 小健 2026-05-18"""
    if not _check_module("docx"):
        return build_error(ERR_NO_DOCX,
                "python-docx库未安装，请先执行: pip install python-docx")

    try:
        import docx
        from docx import Document
        from docx.shared import Inches, Pt

        doc = Document()
        
        if title:
            doc.add_heading(title, 0)
        
        if paragraphs:
            for para in paragraphs:
                doc.add_paragraph(para)
        
        if content:
            doc.add_paragraph(content)
        
        if table_data:
            for tbl in table_data:
                if tbl and len(tbl) > 0:
                    rows = len(tbl)
                    cols = len(tbl[0]) if tbl[0] else 0
                    if rows > 0 and cols > 0:
                        table = doc.add_table(rows=rows, cols=cols)
                        for i, row_data in enumerate(tbl):
                            for j, cell_data in enumerate(row_data):
                                table.rows[i].cells[j].text = str(cell_data if cell_data is not None else "")
        
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(path)
        
        return build_success({"file_path": str(path)}, f"成功写入Word文档: {file_path}")
    except Exception as e:
        return build_error(ERR_WRITE_DOCX, f"写入Word文档失败: {str(e)}")


def _write_xlsx(
    file_path: str,
    data: dict,
    sheet_name: str = "Sheet1"
) -> Dict[str, Any]:
    """写入Excel文件（内部函数） - 小健 2026-05-18"""
    if not _check_module("openpyxl"):
        return build_error(ERR_DOC_NO_OPENPYXL,
                "openpyxl库未安装，请先执行: pip install openpyxl")

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        
        if "headers" in data and data["headers"]:
            headers = data["headers"]
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal="center")
        
        if "rows" in data and data["rows"]:
            for row_idx, row_data in enumerate(data["rows"], 2):
                for col_idx, cell_data in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell_data)
        
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(path)
        
        return build_success({"file_path": str(path), "row_count": len(data.get("rows", []))},
                f"成功写入Excel文件: {file_path}")
    except Exception as e:
        return build_error(ERR_WRITE_XLSX, f"写入Excel文件失败: {str(e)}")


def _write_pdf(
    file_path: str,
    title: str = None,
    content: str = None,
    paragraphs: list = None,
    table_data: list = None
) -> Dict[str, Any]:
    """写入PDF文档（内部函数） - 小健 2026-05-18"""
    if not _check_module("reportlab"):
        return build_error(ERR_NO_REPORTLAB,
                "reportlab库未安装，请先执行: pip install reportlab")

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
        from reportlab.lib import colors
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        doc = SimpleDocTemplate(str(path), pagesize=A4)
        styles = getSampleStyleSheet()

        try:
            font_path = "C:/Windows/Fonts/simsun.ttc"
            pdfmetrics.registerFont(TTFont('SimSun', font_path, subfontIndex=0))
            chinese_style = ParagraphStyle(
                'Chinese', parent=styles['Normal'],
                fontName='SimSun', fontSize=10, leading=14,
                wordWrap='CJK'
            )
            title_style = ParagraphStyle(
                'ChineseTitle', parent=styles['Title'],
                fontName='SimSun', fontSize=18, leading=24,
                wordWrap='CJK'
            )
        except Exception:
            chinese_style = styles['Normal']
            title_style = styles['Title']

        elements = []

        if title:
            elements.append(Paragraph(title, title_style))
            elements.append(Spacer(1, 10*mm))

        if content:
            elements.append(Paragraph(content, chinese_style))
            elements.append(Spacer(1, 5*mm))

        if paragraphs:
            for para in paragraphs:
                elements.append(Paragraph(str(para), chinese_style))
                elements.append(Spacer(1, 3*mm))

        if table_data:
            for tbl in table_data:
                if tbl and len(tbl) > 0:
                    try:
                        t = Table(tbl)
                        t.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                            ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ]))
                        elements.append(t)
                        elements.append(Spacer(1, 5*mm))
                    except Exception:
                        pass

        if not elements:
            elements.append(Paragraph(" ", chinese_style))

        doc.build(elements)

        return build_success({"file_path": str(path)}, f"成功写入PDF文档: {file_path}")
    except Exception as e:
        return build_error(ERR_WRITE_PDF, f"写入PDF文档失败: {str(e)}")


def _write_pptx(
    file_path: str,
    title: str = None,
    slides: list = None
) -> Dict[str, Any]:
    """写入PPT幻灯片（内部函数） - 小健 2026-05-18"""
    if not _check_module("pptx"):
        return build_error(ERR_DOC_NO_PPTX,
                "python-pptx库未安装，请先执行: pip install python-pptx")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        if title:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
        
        if slides:
            for slide_data in slides:
                slide_title = slide_data.get("title", "幻灯片")
                content = slide_data.get("content", "")
                
                content_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_layout)
                
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                
                for shape in slide.shapes:
                    if shape.has_text_frame and not shape.text_frame.text.strip():
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = content
                        p.font.size = Pt(18)
                        break
        
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        
        return build_success({"file_path": str(path), "slide_count": len(prs.slides)},
                f"成功写入PPT文件: {file_path}，共 {len(prs.slides)} 页")
    except Exception as e:
        return build_error(ERR_DOC_WRITE_PPTX, f"写入PPT文件失败: {str(e)}")


# ============================================================
# 路由函数 — LLM调用的统一入口
# ============================================================

def read_document(
    file_path: str,
    pages: Optional[str] = None,
    extract_tables: bool = False,
    sheet_name: Optional[str] = None,
    max_rows: int = 1000,
    header: bool = True,
    encoding: str = "utf-8",
    delimiter: Optional[str] = None,
) -> Dict[str, Any]:
    """读取文档内容 — 小健 2026-05-18
    合并 read_pdf + read_docx + read_pptx + read_xlsx + read_csv
    按文件后缀自动路由到对应解析器
    支持 .doc 后缀（通过 win32com 降级处理）
    支持 .csv/.tsv 后缀（use_pandas=True时使用pandas，否则使用标准库csv）
    """
    path = Path(file_path)
    suffix = path.suffix.lower()
    
    if suffix == ".pdf":
        check = _check_pdf_readable(file_path)
        if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
            return check
        result = _read_pdf(file_path, pages=pages, extract_tables=extract_tables)
    elif suffix == ".docx":
        check = _check_docx_readable(file_path)
        if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
            return check
        result = _read_docx(file_path, extract_tables=extract_tables)
    elif suffix == ".doc":
        # 【自动转换 .doc/.xls 2026-05-20 小健】自动调用convert_document转PDF后读取
        pdf_path = _auto_convert_to_pdf(file_path, suffix)
        if pdf_path is None:
            return build_error(ERR_DOC_CONVERT_FAILED, f"自动转换{suffix}为PDF失败，请手动使用convert_document工具")
        check = _check_pdf_readable(pdf_path)
        if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
            return check
        result = _read_pdf(pdf_path, pages=pages, extract_tables=extract_tables)
        _cleanup_temp_pdf(pdf_path)
    elif suffix == ".pptx":
        result = _read_pptx(file_path)
    elif suffix == ".xlsx":
        check = _check_xlsx_readable(file_path)
        if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
            return check
        result = _read_xlsx(file_path, sheet_name=sheet_name, max_rows=max_rows, header=header)
    elif suffix == ".xls":
        # 【自动转换 .doc/.xls 2026-05-20 小健】自动调用convert_document转PDF后读取
        pdf_path = _auto_convert_to_pdf(file_path, suffix)
        if pdf_path is None:
            return build_error(ERR_DOC_CONVERT_FAILED, f"自动转换{suffix}为PDF失败，请手动使用convert_document工具")
        check = _check_pdf_readable(pdf_path)
        if check["code"] != "SUCCESS" or not check["data"].get("readable", False):
            return check
        result = _read_pdf(pdf_path, pages=pages, extract_tables=extract_tables)
        _cleanup_temp_pdf(pdf_path)
    elif suffix in (".csv", ".tsv"):
        actual_delimiter = "\t" if suffix == ".tsv" else (delimiter or ",")
        result = _read_csv_stdlib(file_path, encoding=encoding, delimiter=actual_delimiter, has_header=header, max_rows=max_rows)
    elif suffix == ".json":
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                json_data = json.load(f)
            if isinstance(json_data, list) and len(json_data) > 0 and isinstance(json_data[0], dict):
                columns = list(json_data[0].keys())
                rows = [[item.get(c) for c in columns] for item in json_data[:max_rows]]
                result = build_success({"format": "json_table", "columns": columns, "rows": rows, "row_count": len(rows)}, f"读取JSON文件成功: {file_path}")
            else:
                result = build_success({"format": "json", "content": json_data}, f"读取JSON文件成功: {file_path}")
        except Exception as e:
            result = build_error(ERR_DOC_READ_JSON, f"读取JSON文件失败: {str(e)}")
    else:
        return build_error(ERR_DOC_FORMAT_NOT_SUPPORTED, f"不支持的格式: {suffix}。支持: .pdf/.docx/.xlsx/.pptx/.csv/.tsv/.json")
    
    if result.get("code") == "SUCCESS":
        result["next_actions"] = build_next_actions([
            ("write_document", "修改文档", "需要编辑文档时"),
            ("convert_document", "转换格式", "需要转PDF/DOCX时"),
            ("analyze_data", "分析数据", "读取的是数据文件时"),
        ])
    return result


def write_document(
    file_path: str,
    content: Optional[str] = None,
    paragraphs: Optional[List[str]] = None,
    title: Optional[str] = None,
    table_data: Optional[List] = None,
    data: Optional[Union[Dict[str, Any], List]] = None,
    sheet_name: str = "Sheet1",
    slides: Optional[List[Dict]] = None,
) -> Dict[str, Any]:
    """写入文档 — 小健 2026-05-18
    合并 write_docx + write_xlsx + write_pdf + write_pptx
    按文件后缀自动路由到对应写入器
    """
    path = Path(file_path)
    suffix = path.suffix.lower()
    path.parent.mkdir(parents=True, exist_ok=True)
    
    if suffix == ".docx":
        result = _write_docx(file_path, content=content, paragraphs=paragraphs, title=title, table_data=table_data)
    elif suffix == ".xlsx":
        if data is None:
            data = {"headers": [], "rows": []}
        elif isinstance(data, list):
            if len(data) > 0 and isinstance(data[0], list):
                first_row = data[0]
                data = {"headers": first_row, "rows": data[1:]}
            elif len(data) > 0 and isinstance(data[0], dict):
                headers = list(data[0].keys())
                rows = [list(row.values()) for row in data]
                data = {"headers": headers, "rows": rows}
            else:
                data = {"headers": [], "rows": data}
        elif isinstance(data, dict) and "headers" not in data and "rows" not in data:
            headers = list(data.keys())
            rows = [list(data.values())]
            data = {"headers": headers, "rows": rows}
        result = _write_xlsx(file_path, data=data, sheet_name=sheet_name)
    elif suffix == ".pdf":
        result = _write_pdf(file_path, title=title, content=content, paragraphs=paragraphs, table_data=table_data)
    elif suffix == ".pptx":
        result = _write_pptx(file_path, title=title, slides=slides)
    else:
        return build_error(ERR_DOC_FORMAT_NOT_SUPPORTED, f"不支持的输出格式: {suffix}。支持: .docx/.xlsx/.pdf/.pptx")
    
    if result.get("code") == "SUCCESS":
        result["next_actions"] = build_next_actions([
            ("read_document", "验证写入结果", "需要确认内容时"),
        ])
    return result


def convert_document(
    input_path: str,
    output_format: Literal["pdf"] = "pdf",
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """文档格式转换 - 小健 2026-05-06 output_format改可选对齐Schema"""
    try:
        src = Path(input_path)
        if not src.exists():
            return build_error(ERR_DOC_CONVERT_FAILED, f"文件不存在: {input_path}")
        
        src_ext = src.suffix.lower()
        supported_inputs = ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.odt', '.ods']
        
        if src_ext not in supported_inputs:
            return build_error(ERR_DOC_CONVERT_FAILED, f"不支持的输入格式: {src_ext}，支持: {supported_inputs}")
        
        if output_format.lower() != 'pdf':
            return build_error(ERR_DOC_CONVERT_FAILED, "当前仅支持转换为PDF格式")
        
        if output_path is None:
            output_path = str(src.with_suffix('.pdf'))
        
        import subprocess
         
        soffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        import platform
        if platform.system() != 'Windows':
            soffice_paths = ["/usr/bin/soffice", "/usr/local/bin/soffice"]
        
        soffice = None
        for p in soffice_paths:
            if Path(p).exists():
                soffice = p
                break
        
        if not soffice:
            return build_error(ERR_NO_LIBREOFFICE, "LibreOffice未安装，无法转换。请安装LibreOffice: https://www.libreoffice.org/download/")
        
        out_dir = str(Path(output_path).parent)
        Path(out_dir).mkdir(parents=True, exist_ok=True)
        
        cmd = [
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", out_dir,
            str(src)
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode != 0:
            return build_error(ERR_DOC_CONVERT_FAILED, f"LibreOffice转换失败: {result.stderr}")
        
        expected_pdf = Path(out_dir) / src.with_suffix('.pdf').name
        if not expected_pdf.exists():
            return build_error(ERR_DOC_CONVERT_FAILED, "转换后PDF文件未生成")
        
        if output_path != str(expected_pdf):
            import shutil

            shutil.move(str(expected_pdf), output_path)
        
        return build_success({"input_path": str(src), "output_path": output_path},
                f"成功转换: {src_ext} → .pdf",
                next_actions=build_next_actions([
                    ("read_document", "读取转换后文件", "需要查看结果时"),
                ]))
    except Exception as e:
        return build_error(ERR_DOC_CONVERT_FAILED, f"文档转换失败: {str(e)}")


# === 公开接口定义 — 小健 2026-05-18 ===
__all__ = [
    "read_document",
    "write_document",
    "convert_document",
]
from app.constants import (
    ERR_DOC_CHART_GENERATE,
    ERR_DOC_CONVERT_FAILED,
    ERR_DOC_FORMAT_NOT_SUPPORTED,
    ERR_DOC_NO_OPENPYXL,
    ERR_DOC_NO_PPTX,
    ERR_DOC_READ_CSV,
    ERR_DOC_READ_DOCX,
    ERR_DOC_READ_JSON,
    ERR_DOC_READ_PDF,
    ERR_DOC_READ_PPTX,
    ERR_DOC_READ_XLSX,
    ERR_DOC_WRITE_PPTX,
    ERR_NO_DOCX,
    ERR_NO_LIBREOFFICE,
    ERR_NO_PANDAS,
    ERR_NO_PDFPLUMBER,
    ERR_NO_REPORTLAB,
    ERR_READ_CSV_DATAFRAME,
    ERR_READ_EXCEL_DATAFRAME,
    ERR_WRITE_DOCX,
    ERR_WRITE_PDF,
    ERR_WRITE_XLSX,
)
