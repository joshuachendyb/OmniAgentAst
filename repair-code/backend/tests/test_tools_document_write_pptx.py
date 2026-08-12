# -*- coding: utf-8 -*-
"""test"""
import json
import os
import tempfile
from pathlib import Path

import pytest

from app.tools.document.write_pptx import (
    write_pptx, _add_pptx_content, _add_pptx_slide, _build_pptx_presentation,
    _dict_table_to_rows, _extract_tables_from_content, _normalize_tables,
)


# ============================================================
# 单元测试 _add_pptx_content content 处理
# ============================================================

def _make_slide(content_or_none):
    """test"""
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[1]  # content layout
    slide = prs.slides.add_slide(layout)
    return slide


class TestAddPptxContent:
    """TestAddPptxContent"""

    def test_content_none(self):
        """content none"""
        prs = _build_pptx_presentation([{"title": "T", "content": "Hello World"}])
        slide = prs.slides[0]
        texts = [p.text for p in slide.placeholders[1].text_frame.paragraphs if p.text.strip()]
        assert any("Hello World" in t for t in texts)

    def test_content_empty_str(self):
        """content empty str"""
        content = [
            {"type": "paragraph", "text": "\u7b2c\u4e00\u6bb5\u843d"},
            {"type": "bullets", "items": ["\u7b2c\u4e00\u9879", "\u7b2c\u4e8c\u9879", "\u7b2c\u4e09\u9879"]},
        ]
        prs = _build_pptx_presentation([{"title": "T", "content": content}])
        assert len(prs.slides) == 1

    def test_content_list_strings(self):
        """content list strings"""
        prs = _build_pptx_presentation([{"title": "T", "content": ["\u4e00", "\u4e8c", "\u4e09"]}])
        assert len(prs.slides) == 1

    def test_content_dict_bullets(self):
        """content dict bullets"""
        prs = _build_pptx_presentation([{"title": "T", "content": {"type": "paragraph", "text": "\u6d4b\u8bd5\u6bb5\u843d"}}])
        assert len(prs.slides) == 1

    def test_content_dict_bullets_cover(self):
        """content dict bullets cover"""
        prs = _build_pptx_presentation([{
            "title": "\u62a5\u544a\u6807\u9898", "subtitle": "\u524d\u8a00\u9884\u89c8", "type": "cover",
            "content": {"type": "bullets", "items": ["\u8981\u70b91", "\u8981\u70b92", "\u8981\u70b93", "\u8981\u70b94"]}
        }])
        assert len(prs.slides) == 1

    def test_content_dict_bullets_two(self):
        """content dict bullets two"""
        prs = _build_pptx_presentation([{
            "title": "\u5bf9\u6bd4\u6807\u9898", "type": "two",
            "content": {"type": "bullets", "items": ["\u5de6\u680fA", "\u5de6\u680fB", "\u5de6\u680fC"]},
            "tables": [[["Title1", "Title2", "Title3"], ["React", "^18", "UI"]]],
        }])
        assert len(prs.slides) == 1

    # --- 内容 / 异常写入 ---

    def test_no_content_but_tables(self):
        """no content but tables"""
        prs = _build_pptx_presentation([{
            "title": "\u4ec5\u8868\u683c\u5e7b\u706f\u7247",
            "tables": [
                [["\u5934", "\u503c"], ["v1", "v2"], ["v3", "v4"]],
                [["A", "B"], ["C", "D"]],
            ]
        }])
        assert len(prs.slides) == 1

    def test_dict_bullets_single_item(self):
        """dict bullets single item"""
        prs = _build_pptx_presentation([{"title": "T", "content": {"type": "bullets", "items": ["\u552f\u4e00\u9879"]}}])
        assert len(prs.slides) == 1

    def test_dict_bullets_many_items(self):
        """dict bullets many items"""
        prs = _build_pptx_presentation([{"title": "T", "content": {"type": "bullets", "items": [f"\u9879{i}" for i in range(20)]}}])
        assert len(prs.slides) == 1

    def test_dict_bullets_empty_items(self):
        """dict bullets empty items"""
        prs = _build_pptx_presentation([{"title": "T", "content": {"type": "bullets", "items": []}}])
        assert len(prs.slides) == 1

    def test_list_mixed_many_types(self):
        """list mixed many types"""
        content = [
            {"type": "paragraph", "text": "\u5f02\u5e38\u5904\u7406"},
            "\u6807\u51c6\u63a5\u53e3\u5b9e\u73b0\u4e8c",
            {"type": "bullets", "items": ["\u5b9e\u73b0\u98791", "\u5b9e\u73b0\u98792"]},
            {"type": "paragraph", "text": "\u7ed3\u675f\u8bed"},
        ]
        prs = _build_pptx_presentation([{"title": "T", "content": content}])
        assert len(prs.slides) == 1

    def test_cover_without_content(self):
        """cover without content"""
        prs = _build_pptx_presentation([{"title": "\u8c22\u8c22", "subtitle": "Q&A", "type": "cover"}])
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "\u8c22\u8c22"

    def test_cover_with_subtitle(self):
        """cover with subtitle"""
        prs = _build_pptx_presentation([{"type": "cover", "title": "\u62a5\u544a\u6807\u9898", "subtitle": "2026-07-02"}])
        assert len(prs.slides) == 1


# ============================================================
# 集成测试 write_pptx 完整写入
# ============================================================

class TestWritePptxIntegration:
    """TestWritePptxIntegration"""

    def test_write_empty_slides(self):
        """write empty slides"""
        result = write_pptx("E:\\test_dir\\_tmp_empty.pptx", [])
        assert result["llm_data"]["status"].get("detail", "") != ""

    def test_write_single_slide(self, tmp_path):
        """write single slide"""
        p = tmp_path / "single.pptx"
        result = write_pptx(str(p), [{"title": "Hello", "content": "World"}])
        assert p.exists()
        assert result["llm_data"]["status"]["exec_code"] == "success"

    def test_write_all_patterns(self, tmp_path):
        """write all patterns"""
        slides = [
            {"type": "cover", "title": "\u9879\u76ee\u62a5\u544a", "subtitle": "2026-07-02", "content": ""},
            {"title": "\u62a5\u544a\u6982\u89c8", "content": {"type": "bullets", "items": ["\u62a5\u544aA", "\u62a5\u544aB", "\u62a5\u544aC"]}},
            {"title": "\u8be6\u7ec6\u8bf4\u660e", "type": "content", "content": {"type": "bullets", "items": ["\u8981\u70b9X", "\u8981\u70b9Y"]}},
            {"title": "\u6982\u89c8", "type": "cover", "subtitle": "\u6458\u8981", "content": {"type": "bullets", "items": ["\u7b2c\u4e00\u9879", "\u7b2c\u4e8c\u9879", "\u7b2c\u4e09\u9879", "\u7b2c\u56db\u9879"]}},
            {"title": "\u5bf9\u6bd4\u5206\u6790", "type": "two", "content": {"type": "bullets", "items": ["\u7ef4\u5ea6A", "\u7ef4\u5ea6B"]}, "tables": [[["\u6307\u6807", "\u503c"], ["\u901a\u8fc7", "\u662f"], ["\u8d28\u91cf", "\u4f18"]]]},
            {"title": "\u603b\u7ed3", "type": "content", "content": [{"type": "paragraph", "text": "\u4ee5\u4e0a\u603b\u7ed3\u5982\u4e0b\uff1a"}, {"type": "bullets", "items": ["\u8981\u70b91", "\u8981\u70b92"]}]},
            {"title": "\u4ee3\u7801\u7247\u6bb5", "content": ["\u7b2c\u4e00\u6bb5", "\u7b2c\u4e8c\u6bb5", "\u7b2c\u4e09\u6bb5"]},
            {"title": "\u6570\u636e\u8868", "tables": [[["\u540d\u79f0", "\u6570\u91cf"], ["A", "10"], ["B", "20"]]]},
            {"type": "cover", "title": "\u8c22\u8c22", "subtitle": "Q&A"},
            {"title": "\u603b\u7ed3", "content": "\u9879\u76ee\u603b\u7ed3\u62a5\u544a\u6458\u8981"},
        ]
        p = tmp_path / "all_patterns.pptx"
        result = write_pptx(str(p), slides)
        assert p.exists()
        assert result["llm_data"]["status"]["exec_code"] == "success"
        from pptx import Presentation
        prs = Presentation(str(p))
        assert len(prs.slides) == len(slides)

    def test_write_file_with_special_chars(self, tmp_path):
        """write file with special chars"""
        p = tmp_path / "test_report_v1.0.pptx"
        result = write_pptx(str(p), [{"title": "T", "content": "OK"}])
        assert result["llm_data"]["status"]["exec_code"] == "success"

    def test_write_to_nonexistent_dir(self):
        """write to nonexistent dir"""
        d = tempfile.mkdtemp()
        p = Path(d) / "newdir" / "sub" / "report.pptx"
        result = write_pptx(str(p), [{"title": "T", "content": "OK"}])
        assert p.exists()
        import shutil
        shutil.rmtree(d)

    def test_write_with_coerce_json(self):
        """write with coerce json"""
        slides_str = '[{"title": "T1", "content": {"type": "bullets", "items": ["a", "b"]}}, {"title": "T2", "content": "plain"}]'
        d = tempfile.mkdtemp()
        p = Path(d) / "coerced.pptx"
        result = write_pptx(str(p), slides_str)
        assert p.exists()
        assert result["llm_data"]["status"]["exec_code"] == "success"
        import shutil
        shutil.rmtree(d)

    def test_massive_slides(self, tmp_path):
        """massive slides"""
        slides = []
        for i in range(50):
            slides.append({
                "title": "\u7b2c%d\u9875" % (i+1),
                "content": [
                    {"type": "paragraph", "text": "\u8fd9\u662f\u7b2c%d\u9875\u7684\u5185\u5bb9" % (i+1)},
                    {"type": "bullets", "items": ["\u9879\u76ee%d" % (i+1), "\u9879\u76ee%d" % (i+2), "\u9879\u76ee%d" % (i+3)]},
                ]
            })
        p = tmp_path / "massive.pptx"
        result = write_pptx(str(p), slides)
        assert p.exists()
        assert result["llm_data"]["status"]["exec_code"] == "success"
        from pptx import Presentation
        prs = Presentation(str(p))
        assert len(prs.slides) == 50

    def test_multiple_tables(self, tmp_path):
        """multiple tables"""
        slides = [
            {"title": "\u7b2c\u4e00\u5f20\u8868", "tables": [[["A", "B"], ["1", "2"]], [["C", "D"], ["3", "4"]]]},
            {"title": "\u7b2c\u4e8c\u5f20\u8868", "tables": [[["X", "Y"], ["5", "6"]]]},
        ]
        p = tmp_path / "multi_tables.pptx"
        result = write_pptx(str(p), slides)
        assert p.exists()
        assert result["llm_data"]["status"]["exec_code"] == "success"
        from pptx import Presentation
        prs = Presentation(str(p))
        assert len(prs.slides) == len(slides)

    def test_p1_01_report_style(self, tmp_path):
        """p1 01 report style"""
        slides = [
            {"type": "cover", "title": "\u9879\u76ee\u5206\u6790\u62a5\u544a", "subtitle": "myproject-react \u9879\u76ee\u68c0\u67e5\u62a5\u544a"},
            {"title": "\u9879\u76ee\u6982\u89c8", "content": {"type": "bullets", "items": [
                "\u9879\u76ee\u540d\u79f0: myproject-react",
                "\u9879\u76ee\u8def\u5f84: E:\\test_dir\\myproject-react",
                "React \u7248\u672c: 18.2.0",
                "\u6267\u884c\u65f6\u95f4: 2026-07-02",
            ]}},
            {"title": "\u9636\u6bb5\u603b\u7ed3", "type": "content", "content": [
                {"type": "paragraph", "text": "\u672c\u6b21\u5206\u6790\u5171\u5206\u4e3a12 \u4e2a\u6d41\u7a0b,\u5206\u4e3a3 \u4e2a\u9636\u6bb5\uff1a"},
                {"type": "bullets", "items": ["\u9636\u6bb5\u4e00: \u6846\u67b6\u6784\u5efa(\u6d41\u7a0b1-2)", "\u9636\u6bb5\u4e8c: \u6838\u5fc3\u529f\u80fd\u5b8c\u6210(\u6d41\u7a0b3-10)", "\u9636\u6bb5\u4e09: \u6d4b\u8bd5\u62a5\u544a\u751f\u6210(\u6d41\u7a0b11-12)"]},
            ]},
            {"title": "\u4f9d\u8d56\u5217\u8868", "tables": [["\u4f9d\u8d56", "React", "React"], ["React", "^18.2.0", "UI\u6846\u67b6"], ["React-DOM", "^18.2.0", "DOM\u6e32\u67d3"], ["React-Scripts", "5.0.1", "\u6784\u5efa\u5de5\u5177"]]},
            {"type": "cover", "title": "\u8c22\u8c22", "subtitle": "Q&A"},
        ]
        p = tmp_path / "p1_01_report.pptx"
        result = write_pptx(str(p), slides)
        assert p.exists()
        assert result["llm_data"]["status"]["exec_code"] == "success"
        from pptx import Presentation
        prs = Presentation(str(p))
        assert len(prs.slides) == len(slides)


# ============================================================
# 单元测试: _dict_table_to_rows — 覆盖dict型表格全部格式
# ============================================================

class TestDictTableToRows:
    """TestDictTableToRows"""

    def test_normal(self):
        """headers+rows标准dict"""
        result = _dict_table_to_rows({"headers": ["a", "b"], "rows": [["1", "2"], ["3", "4"]]})
        assert result == [["a", "b"], ["1", "2"], ["3", "4"]]

    def test_no_headers(self):
        """无headers"""
        result = _dict_table_to_rows({"rows": [["1", "2"], ["3", "4"]]})
        assert result == [["1", "2"], ["3", "4"]]

    def test_no_rows(self):
        """无rows"""
        result = _dict_table_to_rows({"headers": ["a", "b"]})
        assert result == [["a", "b"]]

    def test_empty(self):
        """全空"""
        result = _dict_table_to_rows({})
        assert result == []

    def test_single_row_non_list(self):
        """单行非list"""
        result = _dict_table_to_rows({"headers": ["a"], "rows": ["1"]})
        assert result == [["a"], ["1"]]

    def test_empty_headers_list(self):
        """headers为空列表"""
        result = _dict_table_to_rows({"headers": [], "rows": [["1"]]})
        assert result == [["1"]]


# ============================================================
# 单元测试: _extract_tables_from_content — 覆盖C1~C11
# ============================================================

class TestExtractTablesFromContent:
    """TestExtractTablesFromContent"""

    # --- C1~C4: 纯文本类，原样返回 ---

    def test_c1_str(self):
        """C1 str → 原样"""
        text, tables = _extract_tables_from_content("纯文本")
        assert text == "纯文本"
        assert tables == []

    def test_c2_list_str(self):
        """C2 list[str] → 原样"""
        data = ["行1", "行2"]
        text, tables = _extract_tables_from_content(data)
        assert text == data
        assert tables == []

    def test_c3_list_dict(self):
        """C3 list[dict] → 原样"""
        data = [{"type": "paragraph", "text": "段1"}, {"type": "bullets", "items": ["a", "b"]}]
        text, tables = _extract_tables_from_content(data)
        assert text == data
        assert tables == []

    def test_c4_dict_paragraph(self):
        """C4 dict(paragraph) → 原样"""
        data = {"type": "paragraph", "text": "段落"}
        text, tables = _extract_tables_from_content(data)
        assert text == data
        assert tables == []

    # --- C5: 纯二维数组 → 全转表格 ---

    def test_c5_pure_2d_array(self):
        """C5 list[list]纯二维数组 → 全转表格"""
        data = [["能力大类", "通过数"], ["文件系统", "9"], ["文档处理", "4"]]
        text, tables = _extract_tables_from_content(data)
        assert text is None
        assert len(tables) == 1
        assert tables[0] == data

    def test_c5_single_row(self):
        """C5 单行二维数组"""
        data = [["能力大类", "通过率"]]
        text, tables = _extract_tables_from_content(data)
        assert text is None
        assert tables == [[["能力大类", "通过率"]]]

    # --- C6: 混合列表 → 拆出list行合并成表 ---

    def test_c6_mixed_text_then_table(self):
        """C6 文本在前，表格在后"""
        data = ["说明文字", ["h1", "h2"], ["r1", "r2"]]
        text, tables = _extract_tables_from_content(data)
        assert text == ["说明文字"]
        assert tables == [[["h1", "h2"], ["r1", "r2"]]]

    def test_c6_mixed_table_then_text(self):
        """C6 表格在前，文本在后"""
        data = [["h1", "h2"], ["r1", "r2"], "注"]
        text, tables = _extract_tables_from_content(data)
        assert text == ["注"]
        assert tables == [[["h1", "h2"], ["r1", "r2"]]]

    def test_c6_mixed_interleaved(self):
        """C6 交错: 文本+表格+文本"""
        data = ["说明", ["h1", "h2"], ["r1", "r2"], "注"]
        text, tables = _extract_tables_from_content(data)
        assert text == ["说明", "注"]
        assert tables == [[["h1", "h2"], ["r1", "r2"]]]

    def test_c6_mixed_multiple_tables(self):
        """C6 多张表"""
        data = ["说明", ["h1", "h2"], ["r1", "r2"], "中间", ["a", "b"], "结束"]
        text, tables = _extract_tables_from_content(data)
        assert text == ["说明", "中间", "结束"]
        assert tables == [[["h1", "h2"], ["r1", "r2"]], [["a", "b"]]]

    # --- C7: dict type=table → 转表格 ---

    def test_c7_dict_table(self):
        """C7 dict type=table"""
        data = {"type": "table", "headers": ["能力大类", "通过数"], "rows": [["文件系统", "9"]]}
        text, tables = _extract_tables_from_content(data)
        assert text is None
        assert tables == [[["能力大类", "通过数"], ["文件系统", "9"]]]

    def test_c7_dict_table_empty_rows(self):
        """C7 dict type=table 空rows"""
        data = {"type": "table", "headers": ["a", "b"]}
        text, tables = _extract_tables_from_content(data)
        assert text is None
        assert tables == [[["a", "b"]]]

    # --- C8: list[dict]含type=table ---

    def test_c8_list_with_table_dict(self):
        """C8 list中含dict type=table"""
        data = ["说明", {"type": "table", "headers": ["a", "b"], "rows": [["1", "2"]]}]
        text, tables = _extract_tables_from_content(data)
        assert text == ["说明"]
        assert tables == [[["a", "b"], ["1", "2"]]]

    def test_c8_multiple_table_dicts(self):
        """C8 多个dict type=table"""
        data = [
            {"type": "table", "headers": ["a"], "rows": [["1"]]},
            {"type": "table", "headers": ["b"], "rows": [["2"]]},
        ]
        text, tables = _extract_tables_from_content(data)
        assert text is None
        assert tables == [[["a"], ["1"]], [["b"], ["2"]]]

    # --- C9~C10: 空/None ---

    def test_c9_empty_list(self):
        """C9 [] → 跳过"""
        text, tables = _extract_tables_from_content([])
        assert text is None
        assert tables == []

    def test_c10_none(self):
        """C10 None → 跳过"""
        text, tables = _extract_tables_from_content(None)
        assert text is None
        assert tables == []

    # --- C11: 含None元素 ---

    def test_c11_with_none(self):
        """C11 含None元素 → 过滤None"""
        data = ["a", None, ["h1", "h2"]]
        text, tables = _extract_tables_from_content(data)
        assert text == ["a"]
        assert tables == [[["h1", "h2"]]]

    # --- 混合: C5 + C7 在同一list中 ---

    def test_c5_c7_list_rows_and_dict_table(self):
        """C5 list行 + C7 dict table 混排"""
        data = [["h1", "h2"], {"type": "table", "headers": ["a", "b"], "rows": [["1", "2"]]}, ["r1", "r2"]]
        text, tables = _extract_tables_from_content(data)
        assert text is None
        assert len(tables) == 3
        assert tables[0] == [["h1", "h2"]]
        assert tables[1] == [["a", "b"], ["1", "2"]]
        assert tables[2] == [["r1", "r2"]]


# ============================================================
# 单元测试: _normalize_tables — 覆盖T1~T6
# ============================================================

class TestNormalizeTables:
    """TestNormalizeTables"""

    def test_t1_standard(self):
        """T1 list[list[list]] → 原样"""
        data = [[["h1", "h2"], ["r1", "r2"]]]
        result = _normalize_tables(data)
        assert result == data

    def test_t1_multi_tables(self):
        """T1 多张标准表"""
        data = [[["a", "b"], ["1", "2"]], [["c", "d"], ["3", "4"]]]
        result = _normalize_tables(data)
        assert result == data

    def test_t2_wrap_once(self):
        """T2 list[list]少包一层 → 自动包成[list]"""
        data = [["h1", "h2"], ["r1", "r2"]]
        result = _normalize_tables(data)
        assert result == [[["h1", "h2"], ["r1", "r2"]]]

    def test_t2_single_row(self):
        """T2 单行list[list]"""
        data = [["h1", "h2"]]
        result = _normalize_tables(data)
        assert result == [[["h1", "h2"]]]

    def test_t3_dict(self):
        """T3 dict{headers,rows} → 转"""
        data = {"headers": ["a", "b"], "rows": [["1", "2"]]}
        result = _normalize_tables(data)
        assert result == [[["a", "b"], ["1", "2"]]]

    def test_t3_dict_no_headers(self):
        """T3 dict 无headers"""
        data = {"rows": [["1", "2"]]}
        result = _normalize_tables(data)
        assert result == [[["1", "2"]]]

    def test_t4_list_dict(self):
        """T4 list[dict] → 逐个转"""
        data = [
            {"headers": ["a"], "rows": [["1"]]},
            {"headers": ["b"], "rows": [["2"]]},
        ]
        result = _normalize_tables(data)
        assert result == [[["a"], ["1"]], [["b"], ["2"]]]

    def test_t5_empty_list(self):
        """T5 [] → []"""
        result = _normalize_tables([])
        assert result == []

    def test_t6_none(self):
        """T6 None → []"""
        result = _normalize_tables(None)
        assert result == []


# ============================================================
# 集成测试: G1~G6 组合场景 — 通过_build_pptx_presentation验证
# ============================================================

class TestNormalizeIntegration:
    """TestNormalizeIntegration"""

    def test_g1_text_only(self):
        """G1: content纯文本，无tables"""
        prs = _build_pptx_presentation([{"title": "T", "content": "纯文本"}])
        slide = prs.slides[0]
        texts = [p.text for p in slide.placeholders[1].text_frame.paragraphs if p.text.strip()]
        assert any("纯文本" in t for t in texts)

    def test_g2_text_with_tables(self):
        """G2: content文本 + tables标准表格"""
        prs = _build_pptx_presentation([{
            "title": "T",
            "content": "说明文字",
            "tables": [[["h1", "h2"], ["r1", "r2"]]],
        }])
        slide = prs.slides[0]
        texts = [p.text for p in slide.placeholders[1].text_frame.paragraphs if p.text.strip()]
        assert any("说明文字" in t for t in texts)

    def test_g3_content_2d_array(self):
        """G3: content纯二维数组（原问题场景）"""
        prs = _build_pptx_presentation([{
            "title": "能力验证结果",
            "content": [
                ["能力大类", "子能力数", "通过数", "失败数", "通过率"],
                ["文件系统操作", "9", "9", "0", "100%"],
                ["合计", "45", "45", "0", "100%"],
            ],
        }])
        assert len(prs.slides) == 1

    def test_g4_both_have_tables(self):
        """G4: content二维数组 + tables字段 → 两张表不重叠"""
        prs = _build_pptx_presentation([{
            "title": "两张表",
            "content": [["表1行1"], ["表1行2"]],
            "tables": [[["表2行1"]]],
        }])
        assert len(prs.slides) == 1

    def test_g5_mixed_with_tables(self):
        """G5: content混合list + tables字段"""
        prs = _build_pptx_presentation([{
            "title": "混排+表",
            "content": ["说明", ["t1h", "t1v"], ["t1h2", "t1v2"]],
            "tables": [[["t2h", "t2v"]]],
        }])
        assert len(prs.slides) == 1

    def test_g6_dict_table_with_tables(self):
        """G6: content dict type=table + tables字段"""
        prs = _build_pptx_presentation([{
            "title": "dict表+表",
            "content": {"type": "table", "headers": ["a"], "rows": [["1"]]},
            "tables": [[["b", "2"]]],
        }])
        assert len(prs.slides) == 1
