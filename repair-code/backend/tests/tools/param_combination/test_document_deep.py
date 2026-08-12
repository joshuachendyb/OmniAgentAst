# -*- coding: utf-8 -*-
"""
文档工具参数组合深度测试 - 小健 2026-06-25

测试覆盖8个文档工具:
- read_docx, read_pdf, read_pptx, read_xlsx
- write_docx, write_pdf, write_pptx, write_xlsx

测试类别:
1. ParameterCombinations - 参数组合测试
2. SingleFunction - 单函数功能测试
3. MixedContent - 中英文混合内容测试
4. RealScenarios - 真实场景测试
5. Boundary - 边界条件测试
6. Negative - 为面用例测试
"""

import os
import csv
import pytest
import tempfile
import shutil
from pathlib import Path
from typing import Dict, Any, List

from app.tools.tool_response import is_success, is_error

# ============================================================
# 导入被测工具
# ============================================================
from app.tools.document.read_docx import read_docx
from app.tools.document.read_pdf import read_pdf
from app.tools.document.read_pptx import read_pptx
from app.tools.document.read_xlsx import read_xlsx
from app.tools.document.write_docx import write_docx
from app.tools.document.write_pdf import write_pdf
from app.tools.document.write_pptx import write_pptx
from app.tools.document.write_xlsx import write_xlsx


# ============================================================
# 测试fixtures
# ============================================================

@pytest.fixture
def temp_dir():
    """临时目录fixture"""
    d = tempfile.mkdtemp()
    yield d
    shutil.rmtree(d, ignore_errors=True)


@pytest.fixture
def sample_docx_path(temp_dir):
    """创建测试用docx文件"""
    from docx import Document
    path = os.path.join(temp_dir, "test_report.docx")
    doc = Document()
    doc.add_heading("项目技术文档", 0)
    doc.add_heading("一,系统架构概述", level=1)
    doc.add_paragraph(
        "本系统采用微服务架构设计,前里使用React 18+TypeScript构建;"
        "在里使用FastAPI框架提供RESTful API服务.数据库层采用SQLAlchemy ORM配合SQLite存储,"
        "支持异步操作以提升并发性能.系统支持多模态AI Agent的智能对话和任务执行能力,"
        "包括文件操作,代码生成,数据分析等核心功能模块."
    )
    doc.add_heading("1.1 技术栈详情", level=2)
    doc.add_paragraph("前里框架: React 18 + TypeScript 5 + Vite")
    doc.add_paragraph("在里框架: FastAPI + Uvicorn + SQLAlchemy")
    doc.add_paragraph("数据库: SQLite + aiosqlite异步驱动")
    doc.add_paragraph("测试框架: Vitest + Playwright + pytest")
    doc.add_heading("二,核心功能模块", level=1)
    doc.add_paragraph(
        "Agent系统是本项目的核心,实现了ReAct循环推理机制,"
        "支持工具调用,多轮对话,上下文管理等功能."
        "系统内置了文件操作,Shell命令,网络请求,桌面操作等多种工具类型,"
        "每种工具都实现了安全检查和权限控制机制."
    )
    table = doc.add_table(rows=4, cols=3)
    table.style = 'Table Grid'
    table.rows[0].cells[0].text = "模块名称"
    table.rows[0].cells[1].text = "功能描述"
    table.rows[0].cells[2].text = "状态"
    table.rows[1].cells[0].text = "Agent核心"
    table.rows[1].cells[1].text = "ReAct循环推理引擎"
    table.rows[1].cells[2].text = "已完成"
    table.rows[2].cells[0].text = "工具系统"
    table.rows[2].cells[1].text = "多类型工具注册与调用"
    table.rows[2].cells[2].text = "已完成"
    table.rows[3].cells[0].text = "安全模块"
    table.rows[3].cells[1].text = "工具执行前安全检查"
    table.rows[3].cells[2].text = "已完成"
    doc.add_paragraph("三,性能指标")
    doc.add_paragraph("平均响应时间: 156ms")
    doc.add_paragraph("并发连接数: 1000+")
    doc.add_paragraph("内存占用: 256MB稳定运行")
    doc.save(path)
    return path


@pytest.fixture
def sample_pdf_path(temp_dir):
    """创建测试用PDF文件"""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    path = os.path.join(temp_dir, "test_document.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4)
    styles = getSampleStyleSheet()

    try:
        pdfmetrics.registerFont(TTFont('SimSun', 'C:/Windows/Fonts/simsun.ttc', subfontIndex=0))
        chinese_style = styles['Normal']
    except Exception:
        chinese_style = styles['Normal']

    elements = []
    elements.append(Paragraph("系统技术案范文档", styles['Title']))
    elements.append(Spacer(1, 10 * mm))
    elements.append(Paragraph("第一章 系统概述", styles['Heading1']))
    elements.append(Paragraph(
        "本系统是一个基于人工智能的自动化助手平台,支持多种文档格式的读写操作,"
        "包括Word,PDF,Excel,PowerPoint等主流办公文档格式."
        "系统采用模块化设计,每种文档类型都有独立的读写工具实现,"
        "支持参数化配置和批量处理功能.系统具备完善的安全检查机制,"
        "在执行任何文件操作前都会进行文件类型验证和权限检查.",
        chinese_style
    ))
    elements.append(Spacer(1, 5 * mm))
    elements.append(Paragraph("第二章 功能模块", styles['Heading1']))
    elements.append(Paragraph(
        "2.1 文档读取模块: 支持读取docx,pdf,pptx,xlsx格式文档,提取文本,表格,图片等内容."
        "2.2 文档写入模块: 支持创建docx,pdf,pptx,xlsx格式文档,支持Markdown表格渲染."
        "2.3 安全检查模块: 对所有文件操作进行类型检查和权限验证."
        "2.4 错误处理模块: 统一的错误码定义和用户友好的错误提示.",
        chinese_style
    ))
    elements.append(Spacer(1, 5 * mm))
    elements.append(Paragraph("第三章 性能指标", styles['Heading1']))
    elements.append(Paragraph(
        "系统经过严格的性能测试,文档读取操作平均耗时120ms,写入操作平均耗时85ms."
        "支持并发处理多个文档操作请求,最大并发数可达50个."
        "内存占用稳定在256MB以内,CPU使用率在正常为载下不超过30%.",
        chinese_style
    ))
    doc.build(elements)
    return path


@pytest.fixture
def sample_pptx_path(temp_dir):
    """创建测试用pptx文件"""
    from pptx import Presentation
    from pptx.util import Inches

    path = os.path.join(temp_dir, "test_presentation.pptx")
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    title = slide1.shapes.title
    title.text = "AI智能助手系统介绍"
    subtitle = slide1.placeholders[1]
    subtitle.text = "技术架构与功能概览"

    slide_layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide_layout2)
    title2 = slide2.shapes.title
    title2.text = "系统核心功能"
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "智能对话与任务执行"
    p = tf.add_paragraph()
    p.text = "多模态文档处理"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "代码生成与分析"
    p.level = 1

    slide_layout3 = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(slide_layout3)
    title3 = slide3.shapes.title
    title3.text = "技术架构详情"
    rows, cols = 4, 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "层级"
    table.cell(0, 1).text = "技术栈"
    table.cell(0, 2).text = "说明"
    table.cell(1, 0).text = "前里"
    table.cell(1, 1).text = "React + TypeScript"
    table.cell(1, 2).text = "用户界面层"
    table.cell(2, 0).text = "在里"
    table.cell(2, 1).text = "FastAPI + SQLAlchemy"
    table.cell(2, 2).text = "API服务层"
    table.cell(3, 0).text = "数据库"
    table.cell(3, 1).text = "SQLite"
    table.cell(3, 2).text = "数据持久化层"

    prs.save(path)
    return path


@pytest.fixture
def sample_xlsx_path(temp_dir):
    """创建测试用xlsx文件"""
    from openpyxl import Workbook

    path = os.path.join(temp_dir, "test_data.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "员工信息表"
    ws.append(["工号", "姓名", "部门", "职位", "入职日期"])
    ws.append(["E001", "张三", "技术部", "高级工程师", "2022-03-15"])
    ws.append(["E002", "李四", "产品部", "产品经理", "2021-08-20"])
    ws.append(["E003", "王五", "测试部", "测试工程师", "2023-01-10"])
    ws.append(["E004", "赵六", "技术部", "架构师", "2020-06-01"])
    ws.append(["E005", "钱七", "运维部", "运维工程师", "2022-11-25"])

    ws2 = wb.create_sheet("项目统计")
    ws2.append(["项目名称", "为责人", "进度", "预算"])
    ws2.append(["智能助手系统", "张三", "85%", "500000"])
    ws2.append(["数据分析平台", "李四", "60%", "300000"])
    ws2.append(["自动化测试工具", "王五", "90%", "150000"])

    wb.save(path)
    return path


@pytest.fixture
def sample_csv_path(temp_dir):
    """创建测试用CSV文件"""
    path = os.path.join(temp_dir, "test_records.csv")
    with open(path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(["日期", "类型", "描述", "金额", "备注"])
        writer.writerow(["2026-01-15", "收入", "项目收款", "150000", "合同尾款"])
        writer.writerow(["2026-01-20", "支出", "服务器租赁", "12000", "阿里云ECS"])
        writer.writerow(["2026-02-01", "收入", "咨询服务", "50000", "技术咨询"])
        writer.writerow(["2026-02-10", "支出", "员工工资", "85000", "2月份工资"])
        writer.writerow(["2026-02-15", "支出", "办公设备", "35000", "笔记本电脑x5"])
        writer.writerow(["2026-03-01", "收入", "软件授权", "200000", "年度授权费"])
    return path


# ============================================================
# 1. ParameterCombinations 参数组合测试
# ============================================================

@pytest.mark.timeout(60)
class TestWriteDocxParameterCombinations:
    """write_docx参数组合测试"""

    def test_title_only(self, temp_dir):
        """仅传入title参数"""
        path = os.path.join(temp_dir, "title_only.docx")
        result = write_docx(path=path, title="测试文档标题")
        assert is_success(result), f"仅title参数应成功: {result}"
        assert os.path.exists(path), "文件应已创建"

    def test_content_only(self, temp_dir):
        """仅传入content参数"""
        path = os.path.join(temp_dir, "content_only.docx")
        content = "这是一段测试内容,用于验证write_docx工具的content参数功能."
        result = write_docx(path=path, content=content)
        assert is_success(result), f"仅content参数应成功: {result}"

    def test_table_data_only(self, temp_dir):
        """仅传入table_data参数"""
        path = os.path.join(temp_dir, "table_only.docx")
        table_data = [
            ["姓名", "年龄", "职业"],
            ["张三", "28", "工程师"],
            ["李四", "35", "产品经理"],
        ]
        result = write_docx(path=path, table_data=table_data)
        assert is_success(result), f"仅table_data参数应成功: {result}"

    def test_title_and_content(self, temp_dir):
        """同时传入title和content参数"""
        path = os.path.join(temp_dir, "title_content.docx")
        result = write_docx(
            path=path,
            title="技术报告",
            content="# 第一章\n这是第一章的内容.\n## 1.1 小节\n这是小节内容."
        )
        assert is_success(result), f"title+content参数应成功: {result}"

    def test_content_and_table_data_mutual_exclusion(self, temp_dir):
        """content和table_data同时传入时table_data被忽略(已知行为)"""
        path = os.path.join(temp_dir, "content_table.docx")
        table_data = [["列1", "列2"], ["数据1", "数据2"]]
        result = write_docx(
            path=path,
            content="这是content内容",
            table_data=table_data
        )
        assert is_success(result), f"content+table_data组合应成功: {result}"

    def test_all_params_combined(self, temp_dir):
        """同时传入title,content,table_data"""
        path = os.path.join(temp_dir, "all_params.docx")
        result = write_docx(
            path=path,
            title="完整文档",
            content="这是正文内容",
            table_data=[["A", "B"], ["1", "2"]]
        )
        assert is_success(result), f"所有参数组合应成功: {result}"

    def test_markdown_table_in_content(self, temp_dir):
        """content中嵌入markdown表格"""
        path = os.path.join(temp_dir, "md_table.docx")
        content = """# 数据报告

## 统计表
| 指标 | 数值 | 说明 |
|------|------|------|
| 总用户数 | 1000 | 活跃用户 |
| 日活用户 | 250 | 25%活跃率 |
| 月活用户 | 800 | 80%月活率 |

## 总结

以上为本季度数据统计."""
        result = write_docx(path=path, content=content)
        assert is_success(result), f"Markdown表格应成功: {result}"


@pytest.mark.timeout(60)
class TestWritePdfParameterCombinations:
    """write_pdf参数组合测试"""

    def test_title_only(self, temp_dir):
        """仅传入title参数"""
        path = os.path.join(temp_dir, "title_only.pdf")
        result = write_pdf(path=path, title="PDF文档标题")
        assert is_success(result), f"仅title参数应成功: {result}"
        assert os.path.exists(path), "文件应已创建"

    def test_content_only(self, temp_dir):
        """仅传入content参数"""
        path = os.path.join(temp_dir, "content_only.pdf")
        content = "这是一段PDF测试内容,验证write_pdf工具的content参数功能."
        result = write_pdf(path=path, content=content)
        assert is_success(result), f"仅content参数应成功: {result}"

    def test_table_data_only(self, temp_dir):
        """仅传入table_data参数"""
        path = os.path.join(temp_dir, "table_only.pdf")
        table_data = [
            ["产品名称", "价格", "库存"],
            ["笔记本电脑", "5999", "100"],
            ["机械键盘", "299", "500"],
        ]
        result = write_pdf(path=path, table_data=table_data)
        assert is_success(result), f"仅table_data参数应成功: {result}"

    def test_title_and_content(self, temp_dir):
        """同时传入title和content"""
        path = os.path.join(temp_dir, "title_content.pdf")
        result = write_pdf(
            path=path,
            title="技术案范文档",
            content="# 第一章 系统概述\n系统采用微服务架构设计."
        )
        assert is_success(result), f"title+content参数应成功: {result}"

    def test_content_and_table_data_mutual_exclusion(self, temp_dir):
        """content和table_data同时传入时table_data被忽略(已知行为)"""
        path = os.path.join(temp_dir, "content_table.pdf")
        table_data = [["列1", "列2"], ["数据1", "数据2"]]
        result = write_pdf(
            path=path,
            content="PDF正文内容",
            table_data=table_data
        )
        assert is_success(result), f"content+table_data组合应成功: {result}"

    def test_markdown_table_in_content(self, temp_dir):
        """content中嵌入markdown表格"""
        path = os.path.join(temp_dir, "md_table.pdf")
        content = """# 销售报表
## 月度统计

| 月份 | 销售额 | 增长率 |
|------|--------|--------|
| 1月 | 150000 | +12% |
| 2月 | 180000 | +20% |
| 3月 | 200000 | +11% |

以上为第一季度销售数据."""
        result = write_pdf(path=path, content=content)
        assert is_success(result), f"Markdown表格应成功: {result}"


@pytest.mark.timeout(60)
class TestWritePptxParameterCombinations:
    """write_pptx参数组合测试"""

    def test_single_cover_slide(self, temp_dir):
        """创建单页封面幻灯片"""
        path = os.path.join(temp_dir, "cover.pptx")
        slides = [{"type": 0, "title": "项目报告", "subtitle": "2025年度总结"}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"封面幻灯片应成功: {result}"

    def test_single_content_slide(self, temp_dir):
        """创建单页内容幻灯片"""
        path = os.path.join(temp_dir, "content.pptx")
        slides = [{"type": 1, "title": "核心功能", "content": "系统支持多种文档格式的读写操作."}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"内容幻灯片应成功: {result}"

    def test_slide_with_table(self, temp_dir):
        """幻灯片中包含表格"""
        path = os.path.join(temp_dir, "with_table.pptx")
        slides = [{
            "type": 2,
            "title": "数据统计",
            "tables": [["指标", "数值"], ["用户数", "1000"], ["日活", "250"]]
        }]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"含表格幻灯片应成功: {result}"

    def test_multiple_slides(self, temp_dir):
        """创建多页幻灯片"""
        path = os.path.join(temp_dir, "multi.pptx")
        slides = [
            {"type": 0, "title": "封面", "subtitle": "技术分享"},
            {"type": 1, "title": "目录", "content": "1. 架构设计\n2. 功能模块\n3. 性能优化"},
            {"type": 2, "title": "详情", "content": "详细技术说明"},
        ]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"多页幻灯片应成功: {result}"

    def test_slide_with_bullets(self, temp_dir):
        """幻灯片中包含列表"""
        path = os.path.join(temp_dir, "bullets.pptx")
        slides = [{
            "type": 1,
            "title": "功能列表",
            "content": [
                {"type": "bullets", "items": ["文件操作", "代码生成", "数据分析"]}
            ]
        }]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"含列表幻灯片应成功: {result}"

    def test_empty_tables_list(self, temp_dir):
        """幻灯片中包含空表格列表"""
        path = os.path.join(temp_dir, "empty_tables.pptx")
        slides = [{"type": 1, "title": "空表格", "tables": []}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"空表格列表应成功: {result}"

    def test_slide_type_string_cover(self, temp_dir):
        """使用字符串类型的slide_type"""
        path = os.path.join(temp_dir, "str_type.pptx")
        slides = [{"type": "cover", "title": "封面", "subtitle": "副标题"}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"字符串类型封面应成功: {result}"


@pytest.mark.timeout(60)
class TestWriteXlsxParameterCombinations:
    """write_xlsx参数组合测试"""

    def test_data_only(self, temp_dir):
        """仅传入data参数"""
        path = os.path.join(temp_dir, "data_only.xlsx")
        data = [{"姓名": "张三", "年龄": 28}, {"姓名": "李四", "年龄": 35}]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"仅data参数应成功: {result}"

    def test_custom_sheet_name(self, temp_dir):
        """自定义工作表名"""
        path = os.path.join(temp_dir, "custom_sheet.xlsx")
        data = [{"产品": "A", "销量": 100}]
        result = write_xlsx(path=path, data=data, sheet_name="销售数据")
        assert is_success(result), f"自定义工作表名应成功: {result}"

    def test_empty_data(self, temp_dir):
        """空数据列表 - 内容安全检查应拒绝空数据"""
        path = os.path.join(temp_dir, "empty.xlsx")
        result = write_xlsx(path=path, data=[])
        assert is_error(result), f"空数据应被内容安全检查拒绝: {result}"

    def test_none_data(self, temp_dir):
        """None数据 - 内容安全检查应拒绝None"""
        path = os.path.join(temp_dir, "none_data.xlsx")
        result = write_xlsx(path=path, data=None)
        assert is_error(result), f"None数据应被内容安全检查拒绝: {result}"

    def test_single_row_data(self, temp_dir):
        """单行数据"""
        path = os.path.join(temp_dir, "single_row.xlsx")
        data = [{"ID": 1, "Name": "Test", "Status": "Active"}]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"单行数据应成功: {result}"

    def test_mixed_type_values(self, temp_dir):
        """混合类型值"""
        path = os.path.join(temp_dir, "mixed.xlsx")
        data = [
            {"名称": "项目A", "预算": 500000, "进度": 0.85, "状态": "进行中"},
            {"名称": "项目B", "预算": 300000, "进度": 1.0, "状态": "已完成"},
        ]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"混合类型应成功: {result}"


# ============================================================
# 2. SingleFunction 单函数功能测试
# ============================================================

@pytest.mark.timeout(60)
class TestReadDocxSingleFunction:
    """read_docx单函数功能测试"""

    def test_read_valid_docx(self, sample_docx_path):
        """读取有效docx文件"""
        result = read_docx(path=sample_docx_path)
        assert is_success(result), f"读取有效docx应成功: {result}"
        data = result.get("data", {})
        assert "text" in data, "应包含text字段"
        assert len(data["text"]) > 0, "文本内容不应为空"

    def test_read_docx_tables(self, sample_docx_path):
        """读取docx中的表格"""
        result = read_docx(path=sample_docx_path)
        data = result.get("data", {})
        assert "tables" in data, "应包含tables字段"
        assert len(data["tables"]) > 0, "应至少有1个表格"

    def test_read_docx_metrics(self, sample_docx_path):
        """验证llm_data中的metrics"""
        result = read_docx(path=sample_docx_path)
        llm_data = result.get("llm_data", {})
        metrics = llm_data.get("metrics", {})
        assert "para_count" in metrics, "应包含para_count指标"
        assert "text_len" in metrics, "应包含text_len指标"
        assert metrics["para_count"]["value"] > 0, "段落数应>0"

    def test_read_nonexistent_file(self, temp_dir):
        """读取不存在的文件"""
        path = os.path.join(temp_dir, "nonexistent.docx")
        result = read_docx(path=path)
        assert is_error(result), f"读取不存在文件应失败: {result}"

    def test_read_invalid_extension(self, temp_dir):
        """读取非法扩展名文件"""
        path = os.path.join(temp_dir, "test.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_docx(path=path)
        assert is_error(result), f"读取.txt文件应失败: {result}"

    def test_read_unsupported_format(self, temp_dir):
        """读取不支持的文档格式"""
        path = os.path.join(temp_dir, "test.doc")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_docx(path=path)
        assert is_error(result), f"读取.doc文件应失败: {result}"

    def test_read_docx_content_accuracy(self, temp_dir):
        """验证读取内容的准认性"""
        from docx import Document
        path = os.path.join(temp_dir, "accuracy.docx")
        doc = Document()
        expected_text = "验证文本内容准认性测试"
        doc.add_paragraph(expected_text)
        doc.save(path)
        result = read_docx(path=path)
        data = result.get("data", {})
        assert expected_text in data["text"], f"应包含预期文本: {data['text']}"


@pytest.mark.timeout(60)
class TestReadPdfSingleFunction:
    """read_pdf单函数功能测试"""

    def test_read_valid_pdf(self, sample_pdf_path):
        """读取有效PDF文件"""
        result = read_pdf(path=sample_pdf_path)
        assert is_success(result), f"读取有效PDF应成功: {result}"
        data = result.get("data", {})
        assert "text" in data, "应包含text字段"
        metrics = result.get("llm_data", {}).get("metrics", {})
        assert "page_count" in metrics, "应包含page_count指标"
        assert metrics["page_count"]["value"] >= 1, "页数应>=1"

    def test_read_pdf_pages_read(self, sample_pdf_path):
        """验证pages_read字段"""
        result = read_pdf(path=sample_pdf_path)
        metrics = result.get("llm_data", {}).get("metrics", {})
        assert "pages_read" in metrics, "应包含pages_read指标"
        assert metrics["pages_read"]["value"] >= 1, "已读页数应>=1"

    def test_read_pdf_metrics(self, sample_pdf_path):
        """验证llm_data中的metrics"""
        result = read_pdf(path=sample_pdf_path)
        llm_data = result.get("llm_data", {})
        metrics = llm_data.get("metrics", {})
        assert "page_count" in metrics, "应包含page_count指标"
        assert "text_len" in metrics, "应包含text_len指标"

    def test_read_nonexistent_pdf(self, temp_dir):
        """读取不存在的PDF文件"""
        path = os.path.join(temp_dir, "nonexistent.pdf")
        result = read_pdf(path=path)
        assert is_error(result), f"读取不存在PDF应失败: {result}"

    def test_read_invalid_extension_pdf(self, temp_dir):
        """读取非法扩展名"""
        path = os.path.join(temp_dir, "test.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_pdf(path=path)
        assert is_error(result), f"读取.txt文件应失败: {result}"

    def test_read_pdf_content_check(self, sample_pdf_path):
        """验证PDF内容包含预期文本"""
        result = read_pdf(path=sample_pdf_path)
        data = result.get("data", {})
        assert "Word" in data["text"] or "PDF" in data["text"] or "Excel" in data["text"], "PDF内容应包含预期文本"

    def test_read_pdf_text_not_empty(self, sample_pdf_path):
        """验证PDF文本内容不为空"""
        result = read_pdf(path=sample_pdf_path)
        data = result.get("data", {})
        assert len(data["text"]) > 100, "PDF文本长度应>100字符"


@pytest.mark.timeout(60)
class TestReadPptxSingleFunction:
    """read_pptx单函数功能测试"""

    def test_read_valid_pptx(self, sample_pptx_path):
        """读取有效PPT文件"""
        result = read_pptx(path=sample_pptx_path)
        assert is_success(result), f"读取有效PPT应成功: {result}"
        data = result.get("data", {})
        assert "slides" in data, "应包含slides字段"
        metrics = result.get("llm_data", {}).get("metrics", {})
        assert "slide_count" in metrics, "应包含slide_count指标"
        assert metrics["slide_count"]["value"] == 3, f"应有3页幻灯片,实际: {metrics['slide_count']['value']}"

    def test_read_pptx_slides_content(self, sample_pptx_path):
        """验证幻灯片内容"""
        result = read_pptx(path=sample_pptx_path)
        data = result.get("data", {})
        slides = data["slides"]
        assert len(slides) == 3, "应有3页幻灯片"
        assert slides[0]["slide_num"] == 1, "第一页编号应为1"
        assert "AI智能助手" in slides[0]["text"], "第一页应包含标题文本"

    def test_read_pptx_tables(self, sample_pptx_path):
        """验证PPT中的表格"""
        result = read_pptx(path=sample_pptx_path)
        data = result.get("data", {})
        slide3 = data["slides"][2]
        assert "tables" in slide3, "第三页应包含表格"
        assert len(slide3["tables"]) >= 1, "第三页应至少有1个表格"

    def test_read_nonexistent_pptx(self, temp_dir):
        """读取不存在的PPT文件"""
        path = os.path.join(temp_dir, "nonexistent.pptx")
        result = read_pptx(path=path)
        assert is_error(result), f"读取不存在PPT应失败: {result}"

    def test_read_invalid_extension_pptx(self, temp_dir):
        """读取非法扩展名"""
        path = os.path.join(temp_dir, "test.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_pptx(path=path)
        assert is_error(result), f"读取.txt文件应失败: {result}"

    def test_read_pptx_notes(self, temp_dir):
        """读取幻灯片备注"""
        from pptx import Presentation
        path = os.path.join(temp_dir, "with_notes.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "测试"
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = "这是备注内容"
        prs.save(path)
        result = read_pptx(path=path)
        assert is_success(result), f"读取含备注PPT应成功: {result}"


@pytest.mark.timeout(60)
class TestReadXlsxSingleFunction:
    """read_xlsx单函数功能测试"""

    def test_read_valid_xlsx(self, sample_xlsx_path):
        """读取有效xlsx文件"""
        result = read_xlsx(path=sample_xlsx_path)
        assert is_success(result), f"读取有效xlsx应成功: {result}"
        data = result.get("data", {})
        metrics = result.get("llm_data", {}).get("metrics", {})
        assert "row_count" in metrics, "应包含row_count指标"

    def test_read_xlsx_sheet_names(self, sample_xlsx_path):
        """验证工作表名称"""
        result = read_xlsx(path=sample_xlsx_path)
        data = result.get("data", {})
        assert "sheet_names" in data, "应包含sheet_names字段"
        assert "员工信息表" in data["sheet_names"], "应包含'员工信息表'"
        assert "项目统计" in data["sheet_names"], "应包含'项目统计'"

    def test_read_xlsx_specific_sheet(self, sample_xlsx_path):
        """读取指定工作表"""
        result = read_xlsx(path=sample_xlsx_path, sheet_name="项目统计")
        assert is_success(result), f"读取指定工作表应成功: {result}"
        data = result.get("data", {})
        assert "headers" in data, "应包含headers字段"
        assert data["headers"][0] == "项目名称", f"第一个表头应为'项目名称': {data['headers']}"

    def test_read_xlsx_nonexistent_sheet(self, sample_xlsx_path):
        """读取不存在的工作表"""
        result = read_xlsx(path=sample_xlsx_path, sheet_name="不存在的表")
        assert is_error(result), f"读取不存在工作表应失败: {result}"

    def test_read_valid_csv(self, sample_csv_path):
        """读取有效CSV文件"""
        result = read_xlsx(path=sample_csv_path)
        assert is_success(result), f"读取CSV应成功: {result}"
        data = result.get("data", {})
        assert "headers" in data, "CSV应包含headers字段"
        assert "rows" in data, "CSV应包含rows字段"
        assert len(data["rows"]) >= 1, "CSV应至少有1行数据"

    def test_read_csv_content(self, sample_csv_path):
        """验证CSV内容准认性"""
        result = read_xlsx(path=sample_csv_path)
        data = result.get("data", {})
        assert data["headers"][0] == "日期", f"第一个表头应为'日期': {data['headers']}"
        assert "2026-01-15" in str(data["rows"][0][0]), "第一行日期应包含'2026-01-15'"

    def test_read_nonexistent_xlsx(self, temp_dir):
        """读取不存在的xlsx文件"""
        path = os.path.join(temp_dir, "nonexistent.xlsx")
        result = read_xlsx(path=path)
        assert is_error(result), f"读取不存在xlsx应失败: {result}"


@pytest.mark.timeout(60)
class TestWriteDocxSingleFunction:
    """write_docx单函数功能测试"""

    def test_create_simple_docx(self, temp_dir):
        """创建简单docx文件"""
        path = os.path.join(temp_dir, "simple.docx")
        result = write_docx(path=path, title="简单文档", content="正文内容")
        assert is_success(result), f"创建简单docx应成功: {result}"
        assert os.path.exists(path), "文件应已创建"
        assert os.path.getsize(path) > 0, "文件大小应>0"

    def test_create_docx_with_headings(self, temp_dir):
        """创建带标题层级的docx"""
        path = os.path.join(temp_dir, "headings.docx")
        content = """# 一级标题
## 二级标题
### 三级标题
#### 四级标题
##### 五级标题
正文段落内容."""
        result = write_docx(path=path, content=content)
        assert is_success(result), f"带标题docx应成功: {result}"

    def test_create_docx_with_lists(self, temp_dir):
        """创建带列表的docx"""
        path = os.path.join(temp_dir, "lists.docx")
        content = """# 功能列表

- 文件读写功能
- 代码生成功能
- 数据分析功能

1. 第一项
2. 第二项
3. 第三项"""
        result = write_docx(path=path, content=content)
        assert is_success(result), f"带列表docx应成功: {result}"

    def test_create_docx_empty_content(self, temp_dir):
        """创建空内容docx"""
        path = os.path.join(temp_dir, "empty.docx")
        result = write_docx(path=path)
        assert is_success(result), f"空内容docx应成功: {result}"

    def test_create_docx_nested_dirs(self, temp_dir):
        """在嵌套目录中创建docx"""
        nested_dir = os.path.join(temp_dir, "a", "b", "c")
        os.makedirs(nested_dir, exist_ok=True)
        path = os.path.join(nested_dir, "nested.docx")
        result = write_docx(path=path, title="嵌套文档")
        assert is_success(result), f"嵌套目录docx应成功: {result}"


@pytest.mark.timeout(60)
class TestWritePdfSingleFunction:
    """write_pdf单函数功能测试"""

    def test_create_simple_pdf(self, temp_dir):
        """创建简单PDF文件"""
        path = os.path.join(temp_dir, "simple.pdf")
        result = write_pdf(path=path, title="简单PDF", content="正文内容")
        assert is_success(result), f"创建简单PDF应成功: {result}"
        assert os.path.exists(path), "文件应已创建"

    def test_create_pdf_with_headings(self, temp_dir):
        """创建带标题层级的PDF"""
        path = os.path.join(temp_dir, "headings.pdf")
        content = """# 一级标题
## 二级标题
### 三级标题
#### 四级标题
正文段落内容."""
        result = write_pdf(path=path, content=content)
        assert is_success(result), f"带标题PDF应成功: {result}"

    def test_create_pdf_with_lists(self, temp_dir):
        """创建带列表的PDF"""
        path = os.path.join(temp_dir, "lists.pdf")
        content = """# 功能列表

- 文件读写功能
- 代码生成功能
- 数据分析功能

1. 第一项
2. 第二项
3. 第三项"""
        result = write_pdf(path=path, content=content)
        assert is_success(result), f"带列表PDF应成功: {result}"

    def test_create_pdf_table_data(self, temp_dir):
        """创建含表格的PDF"""
        path = os.path.join(temp_dir, "table.pdf")
        table_data = [
            ["编号", "名称", "状态"],
            ["001", "系统模块", "完成"],
            ["002", "测试模块", "进行中"],
        ]
        result = write_pdf(path=path, table_data=table_data)
        assert is_success(result), f"含表格PDF应成功: {result}"

    def test_create_pdf_nested_dirs(self, temp_dir):
        """在嵌套目录中创建PDF"""
        nested_dir = os.path.join(temp_dir, "x", "y")
        os.makedirs(nested_dir, exist_ok=True)
        path = os.path.join(nested_dir, "nested.pdf")
        result = write_pdf(path=path, title="嵌套PDF")
        assert is_success(result), f"嵌套目录PDF应成功: {result}"


@pytest.mark.timeout(60)
class TestWritePptxSingleFunction:
    """write_pptx单函数功能测试"""

    def test_create_simple_pptx(self, temp_dir):
        """创建简单PPT文件"""
        path = os.path.join(temp_dir, "simple.pptx")
        slides = [{"type": 0, "title": "测试标题"}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"创建简单PPT应成功: {result}"
        assert os.path.exists(path), "文件应已创建"

    def test_create_pptx_with_content_placeholder(self, temp_dir):
        """创建带内容占位符的PPT"""
        path = os.path.join(temp_dir, "content.pptx")
        slides = [{"type": 1, "title": "内容页", "content": "这是正文内容"}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"带内容PPT应成功: {result}"

    def test_create_pptx_empty_slides_error(self, temp_dir):
        """空slides列表应返回错误"""
        path = os.path.join(temp_dir, "empty.pptx")
        result = write_pptx(path=path, slides=[])
        assert is_error(result), f"空slides应失败: {result}"

    def test_create_pptx_none_slides_error(self, temp_dir):
        """None slides应返回错误"""
        path = os.path.join(temp_dir, "none.pptx")
        result = write_pptx(path=path, slides=None)
        assert is_error(result), f"None slides应失败: {result}"

    def test_create_pptx_nested_dirs(self, temp_dir):
        """在嵌套目录中创建PPT"""
        nested_dir = os.path.join(temp_dir, "p", "q")
        os.makedirs(nested_dir, exist_ok=True)
        path = os.path.join(nested_dir, "nested.pptx")
        slides = [{"type": 1, "title": "嵌套"}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"嵌套目录PPT应成功: {result}"


@pytest.mark.timeout(60)
class TestWriteXlsxSingleFunction:
    """write_xlsx单函数功能测试"""

    def test_create_simple_xlsx(self, temp_dir):
        """创建简单xlsx文件"""
        path = os.path.join(temp_dir, "simple.xlsx")
        data = [{"列1": "值1", "列2": "值2"}]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"创建简单xlsx应成功: {result}"
        assert os.path.exists(path), "文件应已创建"

    def test_create_xlsx_empty_data(self, temp_dir):
        """创建空数据xlsx - 内容安全检查应拒绝空数据"""
        path = os.path.join(temp_dir, "empty.xlsx")
        result = write_xlsx(path=path, data=[])
        assert is_error(result), f"空数据xlsx应被内容安全检查拒绝: {result}"

    def test_create_xlsx_multiple_rows(self, temp_dir):
        """创建多行数据xlsx"""
        path = os.path.join(temp_dir, "multi.xlsx")
        data = [
            {"ID": i, "Name": f"Item_{i}", "Value": i * 10}
            for i in range(1, 11)
        ]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"多行xlsx应成功: {result}"

    def test_create_xlsx_nested_dirs(self, temp_dir):
        """在嵌套目录中创建xlsx"""
        nested_dir = os.path.join(temp_dir, "r", "s")
        os.makedirs(nested_dir, exist_ok=True)
        path = os.path.join(nested_dir, "nested.xlsx")
        result = write_xlsx(path=path, data=[{"A": 1}])
        assert is_success(result), f"嵌套目录xlsx应成功: {result}"


# ============================================================
# 3. MixedContent 中英文混合内容测试
# ============================================================

@pytest.mark.timeout(60)
class TestMixedContentDocx:
    """docx中英文混合内容测试"""

    def test_mixed_chinese_english_content(self, temp_dir):
        """中英文混合内容"""
        path = os.path.join(temp_dir, "mixed_docx.docx")
        content = """# AI Assistant System Architecture

## 1. 系统架构 (System Architecture)

本系统采用微服务架构设计,前里使用React 18框架构建现代化用户界面.The backend is built with FastAPI framework, providing high-performance RESTful APIs.

### 1.1 核心模块 (Core Modules)

- **Agent系统**: 实现ReAct循环推理机制,支持多轮对话和工具调用
- **Tool系统**: 支持文件操作,Shell命令,网络请求等多种工具类型
- **Safety模块**: 工具执行前的安全检查和权限控制

The system supports multiple document formats including Word (.docx), PDF, Excel (.xlsx), and PowerPoint (.pptx).

### 1.2 技术栈 (Tech Stack)

| 组件 | 技术选型 | 版本 |
|------|----------|------|
| 前里框架 | React | 18.x |
| 在里框架 | FastAPI | 0.100+ |
| 数据库 | SQLite | 3.x |
| 测试框架 | pytest | 7.x |

## 2. 性能指标 (Performance Metrics)

系统经过严格的性能测试,主要指标如下:

1. **响应时间**: 平均156ms,P95为312ms
2. **并发连接**: 支持1000+并发连接
3. **内存占用**: 稳定运行256MB以内
4. **CPU使用率**: 正常为载下不超过30%

The system has passed all performance benchmarks and is ready for production deployment.
"""
        result = write_docx(path=path, content=content)
        assert is_success(result), f"中英文混合docx应成功: {result}"
        assert os.path.exists(path), "文件应已创建"

    def test_read_mixed_content_docx(self, temp_dir):
        """读取中英文混合内容"""
        from docx import Document
        path = os.path.join(temp_dir, "read_mixed.docx")
        doc = Document()
        doc.add_heading("混合内容文档 Mixed Content Document", 0)
        doc.add_paragraph(
            "This is a bilingual document containing both Chinese and English text. "
            "这是一个包含中英文文本的双语文档."
        )
        doc.add_paragraph(
            "系统支持多种格式:Word,PDF,Excel,PowerPoint."
            "The system supports: Word, PDF, Excel, PowerPoint formats."
        )
        doc.add_paragraph(
            "核心功能包括:文件读写(file read/write),代码生成(code generation),"
            "数据分析 (data analysis),智能对话(intelligent conversation)."
        )
        doc.save(path)
        result = read_docx(path=path)
        assert is_success(result), f"读取混合内容应成功: {result}"
        data = result.get("data", {})
        assert "Chinese" in data["text"] or "Chinese" in str(data.get("text", "")), "应包含英文"
        assert "混合" in data["text"] or "混合" in str(data.get("text", "")), "应包含中文"

    def test_mixed_content_with_table(self, temp_dir):
        """中英文混合表格内容"""
        path = os.path.join(temp_dir, "mixed_table.docx")
        table_data = [
            ["Module模块", "Status状态", "Owner为责人"],
            ["Agent核心", "Completed已完成", "张三San Zhang"],
            ["Tool系统", "In Progress进行中", "李四Si Li"],
            ["Safety模块", "Completed已完成", "王五Wu Wang"],
        ]
        result = write_docx(path=path, table_data=table_data)
        assert is_success(result), f"中英文表格应成功: {result}"


@pytest.mark.timeout(60)
class TestMixedContentPdf:
    """pdf中英文混合内容测试"""

    def test_mixed_chinese_english_pdf(self, temp_dir):
        """中英文混合PDF内容"""
        path = os.path.join(temp_dir, "mixed_pdf.pdf")
        content = """# Technical Documentation 技术文档
## 1. Overview 概述

This document describes the system architecture of the AI Assistant Platform.
本文档描述AI助手平台的系统架构设计.
The platform integrates multiple AI capabilities including natural language processing,
machine learning, and document generation.
平台集成了多种AI能力,包括自然语言处理,机器学习和文档生成.
## 2. Features 功能特点
- 文件操作 File Operations: 支持docx/pdf/xlsx/pptx格式
- 代码生成 Code Generation: 自动生成Python/JavaScript代码
- 数据分析 Data Analysis: 支持CSV/Excel数据分析
- 智能对话 Intelligent Dialogue: 多轮对话上下文管理
The system provides comprehensive error handling and safety checks.
系统提供完善的错误处理和安全检查机制."""
        result = write_pdf(path=path, content=content)
        assert is_success(result), f"中英文混合PDF应成功: {result}"

    def test_read_mixed_pdf(self, temp_dir):
        """读取中英文混合PDF"""
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        path = os.path.join(temp_dir, "read_mixed.pdf")
        doc = SimpleDocTemplate(path, pagesize=A4)
        styles = getSampleStyleSheet()
        try:
            pdfmetrics.registerFont(TTFont('SimSun', 'C:/Windows/Fonts/simsun.ttc', subfontIndex=0))
        except Exception:
            pass
        elements = []
        elements.append(Paragraph("双语文档 Bilingual Document", styles['Title']))
        elements.append(Paragraph(
            "This is a test document with mixed Chinese and English content. "
            "这是一个包含中英文混合内容的测试文档."
            "The system supports multiple languages. 系统支持多语言.",
            styles['Normal']
        ))
        doc.build(elements)
        result = read_pdf(path=path)
        assert is_success(result), f"读取中英文PDF应成功: {result}"


@pytest.mark.timeout(60)
class TestMixedContentXlsx:
    """xlsx中英文混合内容测试"""

    def test_mixed_chinese_english_xlsx(self, temp_dir):
        """中英文混合Excel内容"""
        path = os.path.join(temp_dir, "mixed_xlsx.xlsx")
        data = [
            {"模块名称Module": "Agent系统", "状态Status": "已完成Completed", "为责人Owner": "张三"},
            {"模块名称Module": "Tool系统", "状态Status": "进行中In Progress", "为责人Owner": "李四"},
            {"模块名称Module": "Safety模块", "状态Status": "已完成Completed", "为责人Owner": "王五"},
            {"模块名称Module": "UI组件", "状态Status": "计划中Planned", "为责人Owner": "赵六"},
        ]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"中英文混合xlsx应成功: {result}"

    def test_read_mixed_xlsx(self, temp_dir):
        """读取中英文混合Excel"""
        from openpyxl import Workbook
        path = os.path.join(temp_dir, "read_mixed.xlsx")
        wb = Workbook()
        ws = wb.active
        ws.append(["名称Name", "描述Description", "价值Value"])
        ws.append(["系统架构", "System Architecture", "核心Core"])
        ws.append(["用户界面", "User Interface", "重要Important"])
        ws.append(["数据存储", "Data Storage", "基础Basic"])
        wb.save(path)
        result = read_xlsx(path=path)
        assert is_success(result), f"读取中英文xlsx应成功: {result}"
        data = result.get("data", {})
        assert len(data.get("rows", [])) == 3, f"应有3行数据: {data}"


# ============================================================
# 4. RealScenarios 真实场景测试
# ============================================================

@pytest.mark.timeout(60)
class TestRealScenarios:
    """真实场景测试"""

    def test_write_then_read_docx_roundtrip(self, temp_dir):
        """docx写入在读取往返测试"""
        path = os.path.join(temp_dir, "roundtrip.docx")
        content = """# 项目状态报告
## 一,本周进展
1. 完成Agent核心模块重构
2. 修复Tool系统3个Bug
3. 优化UI组件性能

## 二,下周计划
- 完成安全模块集成测试
- 编写API文档
- 部署测试环境

## 三,风险提示
| 风险项 | 影响 | 应对措施 |
|--------|------|----------|
| 进度延迟 | 中 | 增加人手 |
| 技术债务 | 低 | 持续重构 |"""
        result_write = write_docx(path=path, content=content)
        assert is_success(result_write), f"写入docx应成功: {result_write}"
        result_read = read_docx(path=path)
        assert is_success(result_read), f"读取docx应成功: {result_read}"
        data = result_read.get("data", {})
        assert "Agent" in data["text"], "读取内容应包含'Agent'"
        assert "项目状态报告" in data["text"], "读取内容应包含标题"

    def test_write_then_read_xlsx_roundtrip(self, temp_dir):
        """xlsx写入在读取往返测试"""
        path = os.path.join(temp_dir, "roundtrip.xlsx")
        data = [
            {"产品": "笔记本电脑", "价格": 5999, "库存": 100},
            {"产品": "机械键盘", "价格": 299, "库存": 500},
            {"产品": "显示器", "价格": 1999, "库存": 200},
        ]
        result_write = write_xlsx(path=path, data=data)
        assert is_success(result_write), f"写入xlsx应成功: {result_write}"
        result_read = read_xlsx(path=path)
        assert is_success(result_read), f"读取xlsx应成功: {result_read}"
        read_data = result_read.get("data", {})
        read_rows = read_data.get("rows", [])
        if not read_rows and "sheets" in read_data:
            read_rows = read_data["sheets"][0].get("rows", [])
        assert len(read_rows) == 3, f"应有3行数据: {read_data}"

    def test_write_then_read_csv_roundtrip(self, temp_dir):
        """CSV写入在读取往返测试"""
        csv_path = os.path.join(temp_dir, "roundtrip.csv")
        with open(csv_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerow(["ID", "Name", "Score"])
            writer.writerow(["001", "张三", "95"])
            writer.writerow(["002", "李四", "88"])
        result_read = read_xlsx(path=csv_path)
        assert is_success(result_read), f"读取CSV应成功: {result_read}"
        data = result_read.get("data", {})
        assert len(data.get("rows", [])) == 2, f"应有2行数据: {data}"

    def test_multi_format_document_conversion(self, temp_dir):
        """多格式文档转换场景"""
        content = "# 测试文档\n这是测试内容."
        docx_path = os.path.join(temp_dir, "convert.docx")
        pdf_path = os.path.join(temp_dir, "convert.pdf")
        result1 = write_docx(path=docx_path, content=content)
        assert is_success(result1), f"写入docx应成功: {result1}"
        result2 = write_pdf(path=pdf_path, content=content)
        assert is_success(result2), f"写入pdf应成功: {result2}"
        result3 = read_docx(path=docx_path)
        assert is_success(result3), f"读取docx应成功: {result3}"
        result4 = read_pdf(path=pdf_path)
        assert is_success(result4), f"读取pdf应成功: {result4}"


# ============================================================
# 5. Boundary 边界条件测试
# ============================================================

@pytest.mark.timeout(60)
class TestBoundaryConditions:
    """边界条件测试"""

    def test_write_docx_empty_title(self, temp_dir):
        """空标题"""
        path = os.path.join(temp_dir, "empty_title.docx")
        result = write_docx(path=path, title="", content="正文内容")
        assert is_success(result), f"空标题应成功: {result}"

    def test_write_docx_empty_content(self, temp_dir):
        """空内容 - 内容安全检查应拒绝空内容"""
        path = os.path.join(temp_dir, "empty_content.docx")
        result = write_docx(path=path, title="标题", content="")
        assert is_error(result), f"空内容应被内容安全检查拒绝: {result}"

    def test_write_docx_long_content(self, temp_dir):
        """超长内容"""
        path = os.path.join(temp_dir, "long_content.docx")
        long_content = "这是一段很长很长的内容." * 500
        result = write_docx(path=path, content=long_content)
        assert is_success(result), f"超长内容应成功: {result}"

    def test_write_docx_special_chars(self, temp_dir):
        """特殊字符"""
        path = os.path.join(temp_dir, "special.docx")
        content = "特殊字符测试: !@#$%^&*()_+-=[]{}|;':\",./<>?`~"
        result = write_docx(path=path, content=content)
        assert is_success(result), f"特殊字符应成功: {result}"

    def test_write_pdf_empty_title(self, temp_dir):
        """空标题"""
        path = os.path.join(temp_dir, "empty_title.pdf")
        result = write_pdf(path=path, title="", content="正文内容")
        assert is_success(result), f"空标题应成功: {result}"

    def test_write_pdf_special_chars(self, temp_dir):
        """特殊字符"""
        path = os.path.join(temp_dir, "special.pdf")
        content = "特殊字符测试: !@#$%^&*()_+-=[]{}|;':\",./<>?`~"
        result = write_pdf(path=path, content=content)
        assert is_success(result), f"特殊字符应成功: {result}"

    def test_write_pptx_large_slide_count(self, temp_dir):
        """大量幻灯片"""
        path = os.path.join(temp_dir, "large.pptx")
        slides = [{"type": 1, "title": f"幻灯片{i}", "content": f"第{i}页内容"} for i in range(1, 21)]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"20页幻灯片应成功: {result}"

    def test_write_xlsx_many_rows(self, temp_dir):
        """大量行数据"""
        path = os.path.join(temp_dir, "many_rows.xlsx")
        data = [{"ID": i, "Value": f"data_{i}"} for i in range(1, 101)]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"100行数据应成功: {result}"

    def test_write_xlsx_many_columns(self, temp_dir):
        """大量列"""
        path = os.path.join(temp_dir, "many_cols.xlsx")
        row = {f"col_{i}": f"val_{i}" for i in range(1, 21)}
        result = write_xlsx(path=path, data=[row])
        assert is_success(result), f"20列数据应成功: {result}"

    def test_write_pptx_single_slide(self, temp_dir):
        """单页幻灯片"""
        path = os.path.join(temp_dir, "single.pptx")
        slides = [{"type": 1, "title": "单页", "content": "内容"}]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"单页幻灯片应成功: {result}"


# ============================================================
# 6. Negative 为面用例测试
# ============================================================

@pytest.mark.timeout(60)
class TestNegativeCases:
    """为面用例测试"""

    def test_write_docx_invalid_path(self):
        """无效路径写入docx"""
        result = write_docx(path="Z:\\invalid\\path\\test.docx", title="测试")
        assert is_error(result), f"无效路径应失败: {result}"

    def test_write_pdf_invalid_path(self):
        """无效路径写入pdf"""
        result = write_pdf(path="Z:\\invalid\\path\\test.pdf", title="测试")
        assert is_error(result), f"无效路径应失败: {result}"

    def test_write_pptx_invalid_path(self):
        """无效路径写入pptx"""
        result = write_pptx(path="Z:\\invalid\\path\\test.pptx", slides=[{"type": 1}])
        assert is_error(result), f"无效路径应失败: {result}"

    def test_write_xlsx_invalid_path(self):
        """无效路径写入xlsx"""
        result = write_xlsx(path="Z:\\invalid\\path\\test.xlsx", data=[{"A": 1}])
        assert is_error(result), f"无效路径应失败: {result}"

    def test_read_docx_wrong_extension(self, temp_dir):
        """错误扩展名读取docx"""
        path = os.path.join(temp_dir, "wrong.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_docx(path=path)
        assert is_error(result), f"错误扩展名应失败: {result}"

    def test_read_pdf_wrong_extension(self, temp_dir):
        """错误扩展名读取pdf"""
        path = os.path.join(temp_dir, "wrong.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_pdf(path=path)
        assert is_error(result), f"错误扩展名应失败: {result}"

    def test_read_pptx_wrong_extension(self, temp_dir):
        """错误扩展名读取pptx"""
        path = os.path.join(temp_dir, "wrong.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_pptx(path=path)
        assert is_error(result), f"错误扩展名应失败: {result}"

    def test_read_xlsx_wrong_extension(self, temp_dir):
        """错误扩展名读取xlsx"""
        path = os.path.join(temp_dir, "wrong.txt")
        Path(path).write_text("hello", encoding="utf-8")
        result = read_xlsx(path=path)
        assert is_error(result), f"错误扩展名应失败: {result}"

    def test_write_pptx_invalid_slide_data(self, temp_dir):
        """无效幻灯片数据"""
        path = os.path.join(temp_dir, "invalid.pptx")
        slides = ["not_a_dict", 123, None]
        result = write_pptx(path=path, slides=slides)
        assert is_success(result), f"无效slide数据应跳过但仍成功: {result}"

    def test_write_xlsx_inconsistent_keys(self, temp_dir):
        """不一致的字典键"""
        path = os.path.join(temp_dir, "inconsistent.xlsx")
        data = [
            {"A": 1, "B": 2},
            {"C": 3, "D": 4},
        ]
        result = write_xlsx(path=path, data=data)
        assert is_success(result), f"不一致键应成功: {result}"
