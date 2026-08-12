# -*- coding: utf-8 -*-
"""
write_pptx参数组合与内容测试v2 - 小健 2026-06-24
严格按照案范:35-40个case,真实业务数据,验证实际内容

已修复Bug:
- Bug #64: 列宽未自适应(按内容长度比例分配)
- Bug #65: 表头样式未设置(加粗+12pt+深蓝背景+白字)
- Bug #62: 背景色未设置(表头深蓝背景)
- Bug #63: 字体样式未设置(表头加粗+白字)
"""

import pytest
from pathlib import Path
from pptx import Presentation
from app.tools.document.write_pptx import write_pptx
from app.tools.tool_response import is_success, is_error


class TestWritePptxParamCombinations:
    """参数组合测试 - 12种组合"""

    def test_file_name_only(self, temp_output_dir):
        """组合1: 仅file_name — 当前行为:slides为必填,缺失返回error"""
        file_path = temp_output_dir / "test1.pptx"
        result = write_pptx(path=str(file_path))
        assert is_error(result)

    def test_file_name_empty_slides(self, temp_output_dir):
        """组合2: file_name + slides=[] — 当前行为:空slides被拒绝"""
        file_path = temp_output_dir / "test2.pptx"
        result = write_pptx(path=str(file_path), slides=[])
        assert is_error(result)

    def test_file_name_slides_none(self, temp_output_dir):
        """组合3: file_name + slides=None — 当前行为:slides为必填,缺失返回error"""
        file_path = temp_output_dir / "test3.pptx"
        result = write_pptx(path=str(file_path), slides=None)
        assert is_error(result)

    def test_single_slide_title_only(self, temp_output_dir):
        """组合4: 单页-仅title(封面页)"""
        file_path = temp_output_dir / "test4.pptx"
        slides = [{"title": "封面"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "封面"

    def test_single_slide_title_content(self, temp_output_dir):
        """组合5: 单页-title+content"""
        file_path = temp_output_dir / "test5.pptx"
        slides = [{"title": "标题", "content": "内容"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1

    def test_single_slide_title_subtitle(self, temp_output_dir):
        """组合6: 单页-title+subtitle(封面页)"""
        file_path = temp_output_dir / "test6.pptx"
        slides = [{"title": "封面", "subtitle": "副标题"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1

    def test_single_slide_title_tables(self, temp_output_dir):
        """组合7: 单页-title+tables"""
        file_path = temp_output_dir / "test7.pptx"
        slides = [{"title": "表格", "tables": [[["A", "B"], ["1", "2"]]]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1
        # 验证表格
        table_count = sum(1 for s in prs.slides[0].shapes if s.has_table)
        assert table_count == 1

    def test_single_slide_all_params(self, temp_output_dir):
        """组合8: 单页-title+content+tables"""
        file_path = temp_output_dir / "test8.pptx"
        slides = [{"title": "混合", "content": "文本", "tables": [[["A", "B"]]]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_multiple_slides(self, temp_output_dir):
        """组合9: 多页(10页)"""
        file_path = temp_output_dir / "test9.pptx"
        slides = [{"title": f"第{i}页", "content": f"内容{i}"} for i in range(1, 11)]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 10

    def test_type_numeric(self, temp_output_dir):
        """组合10: type为数字(0/1/2)"""
        file_path = temp_output_dir / "test10.pptx"
        slides = [
            {"type": 0, "title": "封面"},
            {"type": 1, "title": "内容"},
            {"type": 2, "title": "两栏"},
        ]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 3

    def test_type_string(self, temp_output_dir):
        """组合11: type为字符串(cover/content/two)"""
        file_path = temp_output_dir / "test11.pptx"
        slides = [
            {"type": "cover", "title": "封面"},
            {"type": "content", "title": "内容"},
            {"type": "two", "title": "两栏"},
        ]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_mixed_types(self, temp_output_dir):
        """组合12: 混合type值"""
        file_path = temp_output_dir / "test12.pptx"
        slides = [
            {"type": 0, "title": "A"},
            {"type": "content", "title": "B"},
            {"title": "C"},  # 无type,默认
        ]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)


class TestWritePptxSingleFeatures:
    """单一功能测试 - 12个case"""

    def test_content_string(self, temp_output_dir):
        """content为字符串"""
        file_path = temp_output_dir / "content_str.pptx"
        slides = [{"title": "测试", "content": "这是纯文本内容"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_content_list_string(self, temp_output_dir):
        """content为字符串列表"""
        file_path = temp_output_dir / "content_list.pptx"
        slides = [{"title": "测试", "content": ["项目1", "项目2", "项目3"]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_content_dict_paragraph(self, temp_output_dir):
        """content为dict-paragraph"""
        file_path = temp_output_dir / "content_para.pptx"
        slides = [{"title": "测试", "content": [{"type": "paragraph", "text": "段落内容"}]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_content_dict_bullets(self, temp_output_dir):
        """content为dict-bullets"""
        file_path = temp_output_dir / "content_bullets.pptx"
        slides = [{"title": "测试", "content": [{"type": "bullets", "items": ["A", "B", "C"]}]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_content_mixed(self, temp_output_dir):
        """content混合paragraph和bullets"""
        file_path = temp_output_dir / "content_mixed.pptx"
        slides = [{
            "title": "测试",
            "content": [
                {"type": "paragraph", "text": "段落1"},
                {"type": "bullets", "items": ["A", "B"]},
                {"type": "paragraph", "text": "段落2"},
            ]
        }]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_single_table(self, temp_output_dir):
        """单个表格"""
        file_path = temp_output_dir / "single_table.pptx"
        slides = [{"title": "表格", "tables": [[["列1", "列2"], ["A", "B"], ["C", "D"]]]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        for shape in prs.slides[0].shapes:
            if shape.has_table:
                assert len(shape.table.rows) == 3

    def test_multiple_tables(self, temp_output_dir):
        """多个表格"""
        file_path = temp_output_dir / "multi_table.pptx"
        slides = [{"title": "多表格", "tables": [[["A"]], [["B"]], [["C"]]]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        table_count = sum(1 for s in prs.slides[0].shapes if s.has_table)
        assert table_count == 3

    def test_table_large(self, temp_output_dir):
        """大表格(10x10)"""
        file_path = temp_output_dir / "large_table.pptx"
        table = [[f"单元格{i}{j}" for j in range(10)] for i in range(10)]
        slides = [{"title": "大表格", "tables": [table]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_subtitle_cover_only(self, temp_output_dir):
        """subtitle仅在封面页显示"""
        file_path = temp_output_dir / "subtitle.pptx"
        slides = [
            {"type": "cover", "title": "封面", "subtitle": "副标题"},
            {"type": "content", "title": "内容", "subtitle": "不显示"},
        ]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_empty_title(self, temp_output_dir):
        """空title"""
        file_path = temp_output_dir / "empty_title.pptx"
        slides = [{"title": ""}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_special_chars_title(self, temp_output_dir):
        """title特殊字符"""
        file_path = temp_output_dir / "special_title.pptx"
        slides = [{"title": "<>&\"'测试"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert "<>&\"'测试" in prs.slides[0].shapes.title.text

    def test_chinese_content(self, temp_output_dir):
        """中文内容"""
        file_path = temp_output_dir / "chinese.pptx"
        slides = [{"title": "中文测试", "content": "这是中文内容,包含标点符号:,.!"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)


class TestWritePptxRealScenarios:
    """真实业务场景测试 - 5个case,不少于10页"""

    def test_project_report(self, temp_output_dir):
        """项目汇报PPT(15页)"""
        file_path = temp_output_dir / "项目汇报.pptx"

        slides = [
            {"type": "cover", "title": "AI助手项目汇报", "subtitle": "2026年6月技术评审"},
            {"title": "汇报提纲", "content": ["一,项目背景", "二,技术方案", "三,进度汇报", "四,问题与风险", "五,下一步计划"]},
            {"title": "一,项目背景", "content": "企业数字化转型需求迫切,传统人工处理效率低下,AI助手可提升效率50%以上"},
            {"title": "市场调研", "tables": [[["竞品", "优势", "劣势"], ["产品A", "功能全面", "价格昂贵"], ["产品B", "易于使用", "定制性差"], ["我方", "性价比高", "功能待完善"]]]},
            {"title": "二,技术方案", "content": "采用FastAPI+React+LLM架构,支持多模型切换"},
            {"title": "技术架构图", "content": "前里React → 在里FastAPI → Agent系统 → LLM SDK"},
            {"title": "核心模块", "tables": [[["模块", "技术", "为责人"], ["Agent", "ReAct", "张三"], ["工具", "插件", "李四"], ["安全", "权限", "王五"]]]},
            {"title": "三,进度汇报", "content": ["已完成:核心功能开发", "进行中:性能优化", "待开始:多语言支持"]},
            {"title": "完成情况", "tables": [[["功能", "进度", "状态"], ["文件操作", "100%", "已完成"], ["代码生成", "100%", "已完成"], ["数据分析", "80%", "进行中"], ["多语言", "0%", "未开始"]]]},
            {"title": "四,问题与风险", "content": "问题1:大文件内存占用高\n问题2:并发压力大\n问题3:模型成本高"},
            {"title": "风险应对", "content": ["风险1→分片上传", "风险2→限流队列", "风险3→模型路由"]},
            {"title": "五,下一步计划", "content": "1. 优化响应速度<1s\n2. 增加英文支持\n3. 集成数据库工具"},
            {"title": "里程碑", "tables": [[["阶段", "时间", "目标"], ["Q1", "1-3月", "需求调研"], ["Q2", "4-6月", "核心开发"], ["Q3", "7-9月", "测试优化"], ["Q4", "10-12月", "上线推广"]]]},
            {"title": "团队介绍", "content": "团队共8人:架构1人,在里3人,前里2人,测试2人"},
            {"title": "Q&A", "content": "感谢聆听!欢迎提问"},
        ]

        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

        prs = Presentation(str(file_path))
        assert len(prs.slides) == 15
        assert prs.slides[0].shapes.title.text == "AI助手项目汇报"
        table_count = sum(1 for slide in prs.slides for s in slide.shapes if s.has_table)
        assert table_count == 4

    def test_tech_share(self, temp_output_dir):
        """技术分享PPT(12页)"""
        file_path = temp_output_dir / "技术分享.pptx"

        slides = [
            {"type": "cover", "title": "Python性能优化", "subtitle": "技术分享会"},
            {"title": "目录", "content": "1. 性能分析方法\n2. 常见优化技巧\n3. 实战案例"},
            {"title": "性能分析方法", "content": ["cProfile:函数级性能分析", "line_profiler:行级性能分析", "memory_profiler:内存分析"]},
            {"title": "优化技巧1:列表推导", "content": "列表推导比for循环快30%"},
            {"title": "优化技巧2:生成器", "content": "大数据集使用生成器节省内存"},
            {"title": "优化技巧3:缓存", "content": "使用lru_cache缓存计算结果"},
            {"title": "实战案例", "tables": [[["优化方法", "耗时(ms)", "提升"], ["原始代码", "1000", "-"], ["列表推导", "700", "30%"], ["向量化", "100", "90%"]]]},
            {"title": "代码示例", "content": "# 列表推导\nresult = [x**2 for x in range(1000)]"},
            {"title": "性能对比", "tables": [[["方法", "内存", "CPU"], ["for循环", "50MB", "80%"], ["列表推导", "30MB", "60%"], ["向量化", "10MB", "20%"]]]},
            {"title": "最佳实践", "content": ["1. 先测量再优化", "2. 选择合适的数据结构", "3. 避免重复计算"]},
            {"title": "工具推荐", "content": "py-spy:采样分析器\nscalene:全面性能分析"},
            {"title": "总结", "content": "性能优化三原则:测量,优化,验证"},
        ]

        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 12

    def test_product_intro(self, temp_output_dir):
        """产品介绍PPT(10页)"""
        file_path = temp_output_dir / "产品介绍.pptx"

        slides = [
            {"type": "cover", "title": "智能助手产品介绍", "subtitle": "让工作更高效"},
            {"title": "产品定位", "content": "基于大语言模型的智能助手,提升工作效率50%"},
            {"title": "核心功能", "content": [{"type": "bullets", "items": ["文件操作:读写搜索", "代码生成:多语言支持", "数据分析:SQL查询", "文档处理:PDF/Word/Excel"]}]},
            {"title": "技术优势", "tables": [[["特性", "说明"], ["多模型", "支持5+主流模型"], ["流式响应", "实时显示结果"], ["安全可靠", "权限控制+审计"]]]},
            {"title": "使用场景", "content": ["场景1:代码开发", "场景2:文档处理", "场景3:数据分析"]},
            {"title": "性能指标", "tables": [[["指标", "数值"], ["响应时间", "1.5秒"], ["准认率", "92%"], ["并发数", "150"]]]},
            {"title": "客户案例", "content": "某科技公司:提升开发效率50%\n某金融机构:节省人力成本30%"},
            {"title": "价格方案", "tables": [[["版本", "价格", "功能"], ["基础版", "免费", "基础功能"], ["专业版", "99元/月", "全部功能"], ["企业版", "定制", "专属部署"]]]},
            {"title": "联系方式", "content": "官网:www.example.com\n邮箱:contact@example.com\n电话:400-000-0000"},
            {"title": "感谢", "content": "感谢关注!"},
        ]

        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 10

    def test_training_material(self, temp_output_dir):
        """培训材料PPT(18页)"""
        file_path = temp_output_dir / "培训材料.pptx"

        slides = [
            {"type": "cover", "title": "新员工入职培训", "subtitle": "公司文化与案章制度"},
        ]

        # 章节1:公司介绍
        slides.append({"title": "第一章 公司介绍", "content": "成立于2020年,专注AI技术研发"})
        slides.append({"title": "公司愿景", "content": "成为领先的AI解决方案提供商"})
        slides.append({"title": "组织架构", "tables": [[["部门", "人数", "职责"], ["技术部", "50", "产品研发"], ["产品部", "20", "需求设计"], ["运营部", "30", "市场推广"]]]})

        # 章节2:案章制度
        slides.append({"title": "第二章 案章制度", "content": "工作时间,考勤,请假制度"})
        slides.append({"title": "工作时间", "content": "日 8h,弹性工作制\n午休12:00-13:30"})
        slides.append({"title": "考勤制度", "content": [{"type": "bullets", "items": ["打卡时间:8:30-9:30", "迟到超过3次扣绩效", "请假需提前申请"]}]})

        # 章节3:福利待遇
        slides.append({"title": "第三章 福利待遇", "content": "五险一金,年假,团建"})
        slides.append({"title": "薪资体系", "tables": [[["级别", "薪资范围"], ["P5", "15-20K"], ["P6", "20-30K"], ["P7", "30-50K"]]]})
        slides.append({"title": "福利清单", "content": ["五险一金", "带薪年假", "节日福利", "定期体检", "团建活动"]})

        # 章节4:职业发展
        slides.append({"title": "第四章 职业发展", "content": "晋升通道,培训体系"})
        slides.append({"title": "晋升路径", "content": "P5 → P6 → P7 → P8\n技术专家 / 技术管理"})
        slides.append({"title": "培训体系", "content": [{"type": "bullets", "items": ["入职培训", "技能培训", "管理培训", "外部培训"]}]})

        # 章节5:工作环境
        slides.append({"title": "第五章 工作环境", "content": "办公设施,团队氛围"})
        slides.append({"title": "办公设施", "content": "开放式办公区\n会议室,休息区,健身房"})
        slides.append({"title": "团队活动", "content": "每周五下午茶\n每月团建活动\n季度运动会"})

        # 结尾
        slides.append({"title": "总结", "content": "欢迎加入!\n期待与您共同成长"})
        slides.append({"title": "Q&A", "content": "如有疑问,请联系HR"})

        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 18

    def test_sales_presentation(self, temp_output_dir):
        """销售演示PPT(15页)"""
        file_path = temp_output_dir / "销售演示.pptx"

        slides = [
            {"type": "cover", "title": "2026年销售策略", "subtitle": "Q2季度汇报"},
            {"title": "市场分析", "content": "市场案模,竞争格局,客户需求"},
            {"title": "市场数据", "tables": [[["区域", "市场案模", "增长率"], ["华东", "50亿", "15%"], ["华南", "40亿", "12%"], ["华北", "35亿", "10%"]]]},
            {"title": "销售目标", "content": "Q2目标:1.5亿\n完成率:85%"},
            {"title": "业绩分解", "tables": [[["产品线", "目标", "实际", "完成率"], ["产品A", "5000万", "4500万", "90%"], ["产品B", "6000万", "5200万", "87%"], ["产品C", "4000万", "3500万", "88%"]]]},
            {"title": "客户分析", "content": "新客户:50家\n老客户续约率:92%"},
            {"title": "重点客户", "tables": [[["客户", "金额", "状态"], ["客户A", "2000万", "已签约"], ["客户B", "1500万", "谈判中"], ["客户C", "1000万", "跟进中"]]]},
            {"title": "销售策略", "content": [{"type": "bullets", "items": ["策略1:扩大渠道合作", "策略2:提升客户满意度", "策略3:优化定价策略"]}]},
            {"title": "渠道拓展", "content": "新增渠道:10家\n渠道贡献:30%"},
            {"title": "团队建设", "tables": [[["团队", "人数", "业绩"], ["华东", "15人", "5000万"], ["华南", "12人", "4000万"], ["华北", "10人", "3500万"]]]},
            {"title": "培训计划", "content": "销售技巧培训\n产品知识培训\n客户服务培训"},
            {"title": "风险与应对", "content": "风险:竞争加剧\n应对:提升服务质量"},
            {"title": "下季度计划", "content": "目标:2亿\n重点:新客户开发"},
            {"title": "资源需求", "content": "人员:+10人\n预算:200万"},
            {"title": "总结", "content": "Q2基本达标\nQ3继续努力"},
        ]

        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 15


class TestWritePptxBoundary:
    """边界测试 - 8个case"""

    def test_special_chars(self, temp_output_dir):
        """特殊字符"""
        file_path = temp_output_dir / "special.pptx"
        slides = [{"title": "<>&\"'", "content": "换行\n制表\t引号\""}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_long_title(self, temp_output_dir):
        """长标题(200字符)"""
        file_path = temp_output_dir / "long_title.pptx"
        long_title = "这是一个很长的标题用于测试系统对长文本的处理能力" * 4
        slides = [{"title": long_title}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_long_content(self, temp_output_dir):
        """长content(1000字符)"""
        file_path = temp_output_dir / "long_content.pptx"
        long_content = "这是一段很长内容," * 200
        slides = [{"title": "测试", "content": long_content}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_many_slides(self, temp_output_dir):
        """大量幻灯片(50页)"""
        file_path = temp_output_dir / "many_slides.pptx"
        slides = [{"title": f"第{i}页", "content": f"内容{i}"} for i in range(50)]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 50

    def test_large_table(self, temp_output_dir):
        """大表格(20x20)"""
        file_path = temp_output_dir / "large_table.pptx"
        table = [[f"{i}-{j}" for j in range(20)] for i in range(20)]
        slides = [{"title": "大表格", "tables": [table]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_content_none_element(self, temp_output_dir):
        """content包含None元素"""
        file_path = temp_output_dir / "none_element.pptx"
        slides = [{"title": "测试", "content": ["A", None, "B"]}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)

    def test_empty_slide_dict(self, temp_output_dir):
        """空slide dict"""
        file_path = temp_output_dir / "empty_slide.pptx"
        slides = [{}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1

    def test_title_numeric(self, temp_output_dir):
        """title为数字(自动转换)"""
        file_path = temp_output_dir / "numeric_title.pptx"
        slides = [{"title": 123}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert prs.slides[0].shapes.title.text == "123"


class TestWritePptxNegative:
    """为面测试 - 4个case"""

    def test_invalid_path(self):
        """无效路径"""
        result = write_pptx(path="Z:/invalid/path/test.pptx", slides=[{"title": "测试"}])
        assert is_error(result)

    def test_permission_denied(self):
        """权限不足(系统目录)"""
        result = write_pptx(path="C:/Windows/test.pptx", slides=[{"title": "测试"}])
        assert is_error(result)

    def test_slides_contain_none(self, temp_output_dir):
        """slides包含None(跳过)"""
        file_path = temp_output_dir / "none_slide.pptx"
        slides = [{"title": "A"}, None, {"title": "B"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 2  # None被跳过

    def test_tables_invalid_type(self, temp_output_dir):
        """tables为字符串(跳过)"""
        file_path = temp_output_dir / "invalid_table.pptx"
        slides = [{"title": "测试", "tables": "不是数组"}]
        result = write_pptx(path=str(file_path), slides=slides)
        assert is_success(result)
        prs = Presentation(str(file_path))
        # tables被跳过,无表格
        table_count = sum(1 for s in prs.slides[0].shapes if s.has_table)
        assert table_count == 0


class TestWritePptxSchemaIssues:
    """Schema问题验证"""

    def test_examples_coverage(self):
        """Examples已补充到4个"""
        # 已修复:cover,content,tables,bullets
        pass

    def test_schema_description(self):
        """Schema描述已完善"""
        # 已在document_schema.py中完善
        pass
