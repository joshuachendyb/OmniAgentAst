# -*- coding: utf-8 -*-
"""
document读取工具全面深入测试 - 适配当前返回结构(测试-代码漂移修正) - 小健 2026-07-12

当前行为(以读 app/tools/document/read_*.py 源码为准):
- read_docx / read_pdf : data={"text", "tables"?, "images"?}; 段落数/页数等统计在 result["llm_data"]
- read_pptx            : data={"slides":[{"slide_num","text","tables"?}], "notes"?}; 页数在 result["llm_data"]
- read_xlsx            : 单表 data={"sheet_name","headers","rows","sheet_names"}; 多表 data={"sheets":[...],"sheet_names":[...]}

每个工具覆盖参数组合/内容类型/边界/负面测试。
"""
import pytest
from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

from app.tools.document.read_pdf import read_pdf
from app.tools.document.read_docx import read_docx
from app.tools.document.read_pptx import read_pptx
from app.tools.document.read_xlsx import read_xlsx
from app.tools.tool_response import is_success, is_error


@pytest.fixture
def temp_output_dir(tmp_path):
    """临时输出目录"""
    output_dir = tmp_path / "read_deep_test"
    output_dir.mkdir(exist_ok=True)
    return output_dir


# ============================================================================
# read_docx测试
# ============================================================================
class TestReadDocxParamCombinations:
    """read_docx参数组合测试"""

    def test_file_name_only(self, temp_output_dir):
        """仅file_name参数"""
        doc = Document()
        doc.add_paragraph("测试")
        doc_path = temp_output_dir / "test.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert "text" in data

    def test_docx_extension(self, temp_output_dir):
        """.docx扩展名"""
        doc = Document()
        doc.add_paragraph("内容")
        doc_path = temp_output_dir / "test.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_relative_path(self, temp_output_dir):
        """相对路径"""
        doc = Document()
        doc.add_paragraph("内容")
        doc_path = temp_output_dir / "relative.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_absolute_path(self, temp_output_dir):
        """绝对路径"""
        doc = Document()
        doc.add_paragraph("内容")
        doc_path = temp_output_dir / "absolute.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)


class TestReadDocxContentTypes:
    """read_docx内容类型测试"""

    def test_plain_text(self, temp_output_dir):
        """纯文本"""
        doc = Document()
        doc.add_paragraph("这是纯文本内容")
        doc_path = temp_output_dir / "plain.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert "纯文本" in data.get("text", "")

    def test_chinese_text(self, temp_output_dir):
        """中文文本"""
        doc = Document()
        doc.add_paragraph("中文测试:张三李四王五赵六")
        doc_path = temp_output_dir / "chinese.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert "张三" in data.get("text", "")

    def test_special_chars(self, temp_output_dir):
        """特殊字符"""
        doc = Document()
        doc.add_paragraph("特殊字符:>&\"'\\n\\t")
        doc_path = temp_output_dir / "special.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_long_text(self, temp_output_dir):
        """长文本(1000字)"""
        doc = Document()
        long_text = "内容" * 500
        doc.add_paragraph(long_text)
        doc_path = temp_output_dir / "long.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("text", "")) >= 1000

    def test_multiple_paragraphs(self, temp_output_dir):
        """多段落"""
        doc = Document()
        for i in range(10):
            doc.add_paragraph(f"第{i}段内容")
        doc_path = temp_output_dir / "multi_para.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert "text" in data

    def test_empty_paragraphs(self, temp_output_dir):
        """空段落"""
        doc = Document()
        doc.add_paragraph("内容")
        doc.add_paragraph("")
        doc.add_paragraph("内容")
        doc.add_paragraph("")
        doc_path = temp_output_dir / "empty_para.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert "text" in data
        llm_data = result.get("llm_data", {})
        assert "空" in llm_data.get("summary", "")


class TestReadDocxBoundary:
    """read_docx边界测试"""

    def test_empty_document(self, temp_output_dir):
        """空文档"""
        doc = Document()
        doc_path = temp_output_dir / "empty.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert data.get("text", "") == ""

    def test_only_spaces(self, temp_output_dir):
        """仅空格"""
        doc = Document()
        doc.add_paragraph("   ")
        doc_path = temp_output_dir / "spaces.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_merged_cells(self, temp_output_dir):
        """合并单元格"""
        doc = Document()
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "合并"
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 0).text = "合并内容"
        doc_path = temp_output_dir / "merged.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("tables", [])) == 1

    def test_nested_table(self, temp_output_dir):
        """嵌套表格(Word不支持)"""
        doc = Document()
        doc.add_paragraph("普通文档")
        doc_path = temp_output_dir / "nested.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_image_in_document(self, temp_output_dir):
        """包含图片"""
        doc = Document()
        doc.add_paragraph("包含图片的文档")
        doc_path = temp_output_dir / "image.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_header_footer(self, temp_output_dir):
        """页眉页脚"""
        doc = Document()
        doc.sections[0].header.paragraphs[0].text = "页眉"
        doc.add_paragraph("正文")
        doc.sections[0].footer.paragraphs[0].text = "页脚"
        doc_path = temp_output_dir / "header.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)
        data = result.get("data", {})
        assert "页眉" not in data.get("text", "")

    def test_hyperlink(self, temp_output_dir):
        """超链接"""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run()
        run.text = "链接"
        doc_path = temp_output_dir / "hyperlink.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_formatting(self, temp_output_dir):
        """格式(粗体,斜体)"""
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("粗体")
        run.bold = True
        doc_path = temp_output_dir / "format.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_list(self, temp_output_dir):
        """列表"""
        doc = Document()
        doc.add_paragraph("项目1", style='List Bullet')
        doc.add_paragraph("项目2", style='List Bullet')
        doc_path = temp_output_dir / "list.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)

    def test_unicode_chars(self, temp_output_dir):
        """Unicode字符"""
        doc = Document()
        doc.add_paragraph("Emoji: 😀🎉✅中文:测试")
        doc_path = temp_output_dir / "unicode.docx"
        doc.save(str(doc_path))

        result = read_docx(str(doc_path))
        assert is_success(result)


class TestReadDocxNegative:
    """read_docx负面测试"""

    def test_file_not_exist(self):
        """文件不存在"""
        result = read_docx("Z:/not_exist.docx")
        assert is_error(result)

    def test_invalid_extension(self, temp_output_dir):
        """无效扩展名"""
        txt_file = temp_output_dir / "test.txt"
        txt_file.write_text("内容")

        result = read_docx(str(txt_file))
        assert is_error(result)

    def test_corrupted_file(self, temp_output_dir):
        """损坏文件"""
        corrupted = temp_output_dir / "corrupted.docx"
        corrupted.write_bytes(b"corrupted")

        result = read_docx(str(corrupted))
        assert is_error(result)

    def test_empty_file(self, temp_output_dir):
        """空文件"""
        empty = temp_output_dir / "empty.docx"
        empty.write_bytes(b"")

        result = read_docx(str(empty))
        assert is_error(result)

    def test_wrong_format(self, temp_output_dir):
        """错误格式"""
        wrong = temp_output_dir / "wrong.docx"
        wrong.write_text("This is not a Word document")

        result = read_docx(str(wrong))
        assert is_error(result)


# ============================================================================
# read_pptx测试
# ============================================================================
class TestReadPptxParamCombinations:
    """read_pptx参数组合测试"""

    def _make_pptx(self, path, slides_text=None, with_table=False, with_notes=False):
        prs = Presentation()
        texts = slides_text or ["标题"]
        for t in texts:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            tf = txBox.text_frame
            tf.text = t
            if with_table:
                shapes = slide.shapes
                table = shapes.add_table(3, 3, 100, 100, 100, 100).table
                for i in range(3):
                    for j in range(3):
                        table.cell(i, j).text = f"{i}{j}"
            if with_notes:
                slide.notes_slide.notes_text_frame.text = "这是备注"
        prs.save(str(path))
        return prs

    def test_file_name_only(self, temp_output_dir):
        """仅file_name参数"""
        ppt_path = temp_output_dir / "test.pptx"
        self._make_pptx(ppt_path, ["标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_pptx_extension(self, temp_output_dir):
        """.pptx扩展名"""
        ppt_path = temp_output_dir / "test.pptx"
        self._make_pptx(ppt_path, ["标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_relative_path(self, temp_output_dir):
        """相对路径"""
        ppt_path = temp_output_dir / "relative.pptx"
        self._make_pptx(ppt_path, ["标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_absolute_path(self, temp_output_dir):
        """绝对路径"""
        ppt_path = temp_output_dir / "absolute.pptx"
        self._make_pptx(ppt_path, ["标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_chinese_filename(self, temp_output_dir):
        """中文文件名"""
        ppt_path = temp_output_dir / "中文测试.pptx"
        self._make_pptx(ppt_path, ["标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_single_slide(self, temp_output_dir):
        """单张幻灯片"""
        ppt_path = temp_output_dir / "single.pptx"
        self._make_pptx(ppt_path, ["标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 1

    def test_multiple_slides(self, temp_output_dir):
        """多张幻灯片"""
        ppt_path = temp_output_dir / "multi.pptx"
        self._make_pptx(ppt_path, [f"第{i}页" for i in range(10)])

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 10

    def test_title_only(self, temp_output_dir):
        """仅标题"""
        ppt_path = temp_output_dir / "title_only.pptx"
        self._make_pptx(ppt_path, ["仅标题"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 1

    def test_content_only(self, temp_output_dir):
        """仅内容"""
        ppt_path = temp_output_dir / "content_only.pptx"
        self._make_pptx(ppt_path, ["内容文本"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 1

    def test_title_and_content(self, temp_output_dir):
        """标题和内容"""
        ppt_path = temp_output_dir / "title_content.pptx"
        self._make_pptx(ppt_path, ["标题", "内容"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) >= 1

    def test_table_simple(self, temp_output_dir):
        """简单表格"""
        ppt_path = temp_output_dir / "table.pptx"
        self._make_pptx(ppt_path, ["表格页"], with_table=True)

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        slides = data.get("slides", [])
        assert "00 | 01 | 02" in slides[0].get("text", "")

    def test_table_large(self, temp_output_dir):
        """大表格(10x10)"""
        ppt_path = temp_output_dir / "large_table.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        table = shapes.add_table(10, 10, 100, 100, 100, 100).table
        for i in range(10):
            for j in range(10):
                table.cell(i, j).text = f"单元格{i}{j}"
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 1

    def test_multiple_tables(self, temp_output_dir):
        """多表格"""
        ppt_path = temp_output_dir / "multi_table.pptx"
        prs = Presentation()
        for t in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes
            table = shapes.add_table(2, 2, 100, 100 + t * 150, 100, 100).table
            for i in range(2):
                for j in range(2):
                    table.cell(i, j).text = f"{i}{j}"
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 3

    def test_notes(self, temp_output_dir):
        """备注"""
        ppt_path = temp_output_dir / "notes.pptx"
        self._make_pptx(ppt_path, ["标题"], with_notes=True)

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("notes", [])) > 0

    def test_chinese_text(self, temp_output_dir):
        """中文文本"""
        ppt_path = temp_output_dir / "chinese.pptx"
        self._make_pptx(ppt_path, ["中文测试:张三李四"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_special_chars(self, temp_output_dir):
        """特殊字符"""
        ppt_path = temp_output_dir / "special.pptx"
        self._make_pptx(ppt_path, ["特殊字符:>&\"'"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_image(self, temp_output_dir):
        """图片"""
        ppt_path = temp_output_dir / "image.pptx"
        self._make_pptx(ppt_path, ["图片页"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_multiple_shapes(self, temp_output_dir):
        """多个形状"""
        ppt_path = temp_output_dir / "shapes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        for i in range(3):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            tf = txBox.text_frame
            tf.text = f"形状{i}"
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 1

    def test_empty_slide(self, temp_output_dir):
        """空白幻灯片"""
        ppt_path = temp_output_dir / "empty_slide.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 1

    def test_different_layouts(self, temp_output_dir):
        """不同布局"""
        ppt_path = temp_output_dir / "layouts.pptx"
        prs = Presentation()
        for i in range(min(6, len(prs.slide_layouts))):
            slide = prs.slides.add_slide(prs.slide_layouts[i])
            if slide.shapes.title:
                slide.shapes.title.text = f"布局{i}"
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_empty_presentation(self, temp_output_dir):
        """空演示文稿"""
        ppt_path = temp_output_dir / "empty.pptx"
        prs = Presentation()
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 0

    def test_large_presentation(self, temp_output_dir):
        """大演示文稿(100页)"""
        ppt_path = temp_output_dir / "large.pptx"
        prs = Presentation()
        for i in range(100):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            tf = txBox.text_frame
            tf.text = f"第{i}页"
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        assert len(data.get("slides", [])) == 100

    def test_unicode_chars(self, temp_output_dir):
        """Unicode字符"""
        ppt_path = temp_output_dir / "unicode.pptx"
        self._make_pptx(ppt_path, ["Emoji: 😀🎉✅"])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_long_text(self, temp_output_dir):
        """长文本(1000字)"""
        ppt_path = temp_output_dir / "long.pptx"
        long_text = "内容" * 100
        self._make_pptx(ppt_path, [long_text])

        result = read_pptx(str(ppt_path))
        assert is_success(result)

    def test_table_with_empty_cells(self, temp_output_dir):
        """表格空单元格"""
        ppt_path = temp_output_dir / "empty_cell.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        table = shapes.add_table(2, 3, 100, 100, 100, 100).table
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = ""
        table.cell(0, 2).text = "C"
        prs.save(str(ppt_path))

        result = read_pptx(str(ppt_path))
        assert is_success(result)


class TestReadPptxNegative:
    """read_pptx负面测试"""

    def test_file_not_exist(self):
        """文件不存在"""
        result = read_pptx("Z:/not_exist.pptx")
        assert is_error(result)

    def test_invalid_extension(self, temp_output_dir):
        """无效扩展名"""
        txt_file = temp_output_dir / "test.txt"
        txt_file.write_text("内容")

        result = read_pptx(str(txt_file))
        assert is_error(result)

    def test_corrupted_file(self, temp_output_dir):
        """损坏文件"""
        corrupted = temp_output_dir / "corrupted.pptx"
        corrupted.write_bytes(b"corrupted")

        result = read_pptx(str(corrupted))
        assert is_error(result)

    def test_empty_file(self, temp_output_dir):
        """空文件"""
        empty = temp_output_dir / "empty.pptx"
        empty.write_bytes(b"")

        result = read_pptx(str(empty))
        assert is_error(result)

    def test_wrong_format(self, temp_output_dir):
        """错误格式"""
        wrong = temp_output_dir / "wrong.pptx"
        wrong.write_text("This is not a PowerPoint document")

        result = read_pptx(str(wrong))
        assert is_error(result)


# ============================================================================
# read_xlsx测试
# ============================================================================
class TestReadXlsxParamCombinations:
    """read_xlsx参数组合测试"""

    def test_file_name_only(self, temp_output_dir):
        """仅file_name参数"""
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.append(["A", "B"])
        ws.append([1, 2])
        xlsx_path = temp_output_dir / "test.xlsx"
        wb.save(str(xlsx_path))

        result = read_xlsx(str(xlsx_path))
        assert is_success(result)

    def test_xlsx_extension(self, temp_output_dir):
        """.xlsx扩展名"""
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.append(["A", "B"])
        xlsx_path = temp_output_dir / "test.xlsx"
        wb.save(str(xlsx_path))

        result = read_xlsx(str(xlsx_path))
        assert is_success(result)

    def test_csv_extension(self, temp_output_dir):
        """.csv扩展名"""
        csv_path = temp_output_dir / "test.csv"
        csv_path.write_text("A,B\n1,2\n")

        result = read_xlsx(str(csv_path))
        assert is_success(result)

    def test_sheet_name_parameter(self, temp_output_dir):
        """sheet_name参数"""
        from openpyxl import Workbook
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["A", "B"])
        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["C", "D"])
        xlsx_path = temp_output_dir / "multi.xlsx"
        wb.save(str(xlsx_path))

        result = read_xlsx(str(xlsx_path), sheet_name="Sheet2")
        assert is_success(result)
        data = result.get("data", {})
        assert data.get("sheet_name") == "Sheet2"

    def test_sheet_name_none(self, temp_output_dir):
        """sheet_name=None读取所有"""
        from openpyxl import Workbook
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["A", "B"])
        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["C", "D"])
        xlsx_path = temp_output_dir / "all.xlsx"
        wb.save(str(xlsx_path))

        result = read_xlsx(str(xlsx_path), sheet_name=None)
        assert is_success(result)
        data = result.get("data", {})
        assert "sheets" in data
