# -*- coding: utf-8 -*-
"""
write_pptx参数组合与内容测试 - 小健 2026-06-24

测试目标:
1. 参数组合:file_name(必填), slides(可选)
2. 功能点:封面页,内容页,多页,标题,正文,列表,表格
3. 真实场景:项目汇报,技术分享,产品介绍
4. 边界测试:空slides,长标题,特殊字符
5. 负面测试:无效路径
"""

import pytest
from pathlib import Path
from pptx import Presentation
from app.tools.document.write_pptx import write_pptx
from app.tools.tool_response import is_success, is_error


class TestWritePptxParamCombinations:
    """参数组合测试 - 4种组合"""

    def test_file_name_only(self, temp_output_dir):
        """组合1: 仅必填参数file_name — 当前行为:slides为必填,缺失返回error"""
        file_path = temp_output_dir / "empty.pptx"
        result = write_pptx(path=str(file_path))

        assert is_error(result)

    def test_file_name_with_empty_slides(self, temp_output_dir):
        """组合2: file_name + slides=[] — 当前行为:空slides被拒绝"""
        file_path = temp_output_dir / "empty_slides.pptx"
        result = write_pptx(path=str(file_path), slides=[])

        assert is_error(result)

    def test_file_name_with_slides(self, temp_output_dir):
        """组合3: file_name + slides"""
        file_path = temp_output_dir / "slides.pptx"
        slides = [{"title": "测试标题", "content": "测试内容"}]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1

    def test_slides_none(self, temp_output_dir):
        """组合4: slides=None — 当前行为:slides为必填,缺失返回error"""
        file_path = temp_output_dir / "none.pptx"
        result = write_pptx(path=str(file_path), slides=None)

        assert is_error(result)


class TestWritePptxSingleFeatures:
    """单一功能测试"""

    def test_cover_page(self, temp_output_dir):
        """封面页(仅有title)"""
        file_path = temp_output_dir / "cover.pptx"
        slides = [{"title": "项目汇报"}]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        assert slide.shapes.title.text == "项目汇报"

    def test_content_page(self, temp_output_dir):
        """内容页(title + content)"""
        file_path = temp_output_dir / "content.pptx"
        slides = [{"title": "业绩概览", "content": "本季度销售额增长20%"}]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1

    def test_multiple_slides(self, temp_output_dir):
        """多页PPT(10页)"""
        file_path = temp_output_dir / "multi.pptx"
        slides = [{"title": f"第{i}页", "content": f"第{i}页内容"} for i in range(1, 11)]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 10

    def test_subtitle(self, temp_output_dir):
        """副标题(封面页)"""
        file_path = temp_output_dir / "subtitle.pptx"
        slides = [{"title": "项目汇报", "subtitle": "2026年度总结"}]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        # 验证:副标题是否显示

    def test_content_list(self, temp_output_dir):
        """content为列表"""
        file_path = temp_output_dir / "list.pptx"
        slides = [{
            "title": "功能列表",
            "content": ["功能1:数据导入", "功能2:数据分析", "功能3:报告生成"]
        }]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 1

    def test_content_dict_bullets(self, temp_output_dir):
        """content为dict,包含bullets"""
        file_path = temp_output_dir / "bullets.pptx"
        slides = [{
            "title": "要点总结",
            "content": [
                {"type": "paragraph", "text": "主要成果:"},
                {"type": "bullets", "items": ["完成100%目标", "提升效率30%", "降低成本20%"]}
            ]
        }]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)

    def test_tables(self, temp_output_dir):
        """表格功能"""
        file_path = temp_output_dir / "table.pptx"
        slides = [{
            "title": "数据对比",
            "tables": [
                [["项目", "数值"], ["A", "100"], ["B", "200"]]
            ]
        }]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        slide = prs.slides[0]
        # 验证:是否有表格shape
        table_shapes = [s for s in slide.shapes if s.has_table]
        assert len(table_shapes) == 1


class TestWritePptxRealScenarios:
    """真实业务场景测试 - 不少于10页"""

    def test_project_report(self, temp_output_dir):
        """项目汇报PPT(真实场景)"""
        file_path = temp_output_dir / "项目汇报.pptx"

        slides = [
            {"title": "AI助手项目汇报", "subtitle": "2026年1月"},
            {"title": "项目背景", "content": "随着AI技术发展,企业对智能助手需求日益增长"},
            {"title": "项目目标", "content": [
                "目标1:提升用户效率50%",
                "目标2:降低人工成本30%",
                "目标3:实现7x24小时服务"
            ]},
            {"title": "技术架构", "content": "采用FastAPI + React + LLM架构,支持多模型切换"},
            {"title": "核心功能", "content": [
                {"type": "paragraph", "text": "已实现功能:"},
                {"type": "bullets", "items": ["文件操作", "代码生成", "数据分析", "文档处理"]}
            ]},
            {"title": "性能指标", "tables": [[
                ["指标", "目标", "实际"],
                ["响应时间", "<2s", "1.5s"],
                ["准认率", ">90%", "92%"],
                ["并发数", ">100", "150"]
            ]]},
            {"title": "团队介绍", "content": "团队共8人:架构师1人,在里3人,前里4人"},
            {"title": "里程碑", "content": [
                "Q1:需求调研与技术选型",
                "Q2:核心功能开发",
                "Q3:测试与优化",
                "Q4:上线推广"
            ]},
            {"title": "风险与应对", "content": "主要风险:模型成本高\n应对策略:引入多模型路由,根据任务复杂度选择模型"},
            {"title": "下一步计划", "content": "1. 优化响应速度\n2. 增加多语言支持\n3. 集成更多工具"},
            {"title": "Q&A", "content": "感谢聆听!"},
        ]

        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 11
        # 验证标题
        assert prs.slides[0].shapes.title.text == "AI助手项目汇报"

    def test_tech_share(self, temp_output_dir):
        """技术分享PPT"""
        file_path = temp_output_dir / "技术分享.pptx"

        slides = [
            {"title": "Python性能优化", "subtitle": "技术分享会"},
            {"title": "目录", "content": "1. 性能分析方法\n2. 常见优化技巧\n3. 实战案例"},
            {"title": "性能分析方法", "content": [
                "cProfile:函数级性能分析",
                "line_profiler:行级性能分析",
                "memory_profiler:内存分析"
            ]},
            {"title": "优化技巧1:列表推导", "content": "列表推导比for循环快50%"},
            {"title": "优化技巧2:生成器", "content": "大数据集使用生成器节省内存"},
            {"title": "优化技巧3:缓存", "content": "使用lru_cache缓存计算结果"},
            {"title": "实战案例", "tables": [[
                ["优化方法", "耗时(ms)", "提升"],
                ["原始代码", "1000", "-"],
                ["列表推导", "700", "30%"],
                ["向量化", "100", "90%"]
            ]]},
            {"title": "总结", "content": "性能优化三原则:测量,优化,验证"},
        ]

        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 8


class TestWritePptxBoundary:
    """边界测试"""

    def test_special_chars(self, temp_output_dir):
        """特殊字符:<>&"'"""
        file_path = temp_output_dir / "special.pptx"
        slides = [{"title": "特殊字符<>&\"'", "content": "测试内容<>&\"'"}]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert "<>&\"'" in prs.slides[0].shapes.title.text

    def test_long_title(self, temp_output_dir):
        """长标题(100字符)"""
        file_path = temp_output_dir / "long.pptx"
        long_title = "这是一个很长的标题用于测试系统对长文本的处理能力" * 2
        slides = [{"title": long_title}]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)

    def test_many_slides(self, temp_output_dir):
        """大量幻灯片(50页)"""
        file_path = temp_output_dir / "many.pptx"
        slides = [{"title": f"第{i}页", "content": f"内容{i}"} for i in range(50)]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 50


class TestWritePptxNegative:
    """负面测试"""

    def test_invalid_path(self):
        """无效路径"""
        result = write_pptx(path="Z:/invalid/path/test.pptx", slides=[{"title": "测试"}])
        assert is_error(result)


class TestWritePptxSchemaIssues:
    """Schema问题验证"""

    def test_examples_coverage(self):
        """Examples覆盖不足:只有1个示例"""
        # Schema有4个参数,Examples有1个,但:
        # 1. 没有展示tables用法
        # 2. 没有展示content为list的用法
        # 3. 没有展示subtitle用法
        # Bug:Examples不够丰富,LLM不知道高级用法
        pass

    def test_slides_description_clarity(self):
        """slides参数描述不够清晰"""
        # 描述说"幻灯片列表",但没有说明:
        # 1. type字段的可选值(0/1/2, "cover"/"content"/"two")
        # 2. content支持哪些格式(str, list, dict)
        # 3. tables格式要求
        # Bug:Schema描述不完整
        pass

    def test_slide_type_values(self, temp_output_dir):
        """验证slide type的不同值"""
        file_path = temp_output_dir / "types.pptx"
        slides = [
            {"type": 0, "title": "封面"},  # type=0
            {"type": 1, "title": "内容页", "content": "内容"},  # type=1
            {"type": "cover", "title": "封面子符串"},  # type="cover"
            {"type": "content", "title": "内容页字符串", "content": "内容"},  # type="content"
        ]
        result = write_pptx(path=str(file_path), slides=slides)

        assert is_success(result)
        prs = Presentation(str(file_path))
        assert len(prs.slides) == 4
