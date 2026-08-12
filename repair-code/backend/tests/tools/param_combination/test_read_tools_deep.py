# -*- coding: utf-8 -*-
"""
文档读取工具深度测试 - 发现Bug - 小健 2026-06-24

严格按照测试案范:
1. Schema驱动原则 - 覆盖所有参数组合
2. 内容丰富性原则 - 真实业务数据
3. 验证完整性原则 - 必须验证实际结果
4. 问题发现原则 - 测试目的是发现问题

测试工具:read_pdf, read_docx, read_pptx, read_xlsx
"""
import pytest
from pathlib import Path
from app.tools.document.read_pdf import read_pdf
from app.tools.document.read_docx import read_docx
from app.tools.document.read_pptx import read_pptx
from app.tools.document.read_xlsx import read_xlsx
from app.tools.tool_response import is_success, is_error


@pytest.fixture
def test_data_dir():
    """测试数据目录"""
    return Path("G:/OmniAgentAs-desk/backend/tests/data")


class TestReadPdfBugs:
    """read_pdf Bug测试"""
    
    def test_bug1_file_not_exist(self):
        """Bug #1: 文件不存在时的错误信息不够清晰"""
        result = read_pdf("Z:/not_exist.pdf")
        assert is_error(result)
        # Bug: 错误信息应该包含文件路径
        llm_data = result.get("llm_data", {})
        # 当前: 文件路径位于 summary 与 action.target 中(status.detail 为驱动器校验信息)
        assert "not_exist.pdf" in llm_data.get("summary", "") or \
            "not_exist.pdf" in llm_data.get("action", {}).get("target", "")
    
    def test_bug2_empty_pdf(self, test_data_dir, tmp_path):
        """Bug #2: 空PDF文件处理"""
        # 创建空PDF(需要测试数据)
        pass
    
    def test_bug3_invalid_extension(self, test_data_dir, tmp_path):
        """Bug #3: 无效扩展名处理"""
        # 创建一个非PDF文件但命名为.pdf
        fake_pdf = tmp_path / "fake.pdf"
        fake_pdf.write_text("This is not a PDF")
        result = read_pdf(str(fake_pdf))
        # Bug: 应该报错,但可能返回success
        assert is_error(result)
    
    def test_bug4_large_pdf(self, test_data_dir):
        """Bug #4: 大PDF文件性能"""
        # 需要大PDF测试数据
        pass
    
    def test_bug5_corrupted_pdf(self, test_data_dir, tmp_path):
        """Bug #5: 损坏的PDF文件"""
        # 创建损坏的PDF
        corrupted = tmp_path / "corrupted.pdf"
        corrupted.write_bytes(b"%PDF-1.4\n%corrupted")
        result = read_pdf(str(corrupted))
        # Bug: 应该返回明认的错误信息
        assert is_error(result)
    
    def test_bug6_chinese_filename(self, tmp_path):
        """Bug #6: 中文文件名"""
        # 创建测试PDF(需要真实PDF)
        pass
    
    def test_bug7_special_chars_filename(self, tmp_path):
        """Bug #7: 特殊字符文件名"""
        # 创建文件名包含特殊字符的PDF
        pass
    
    def test_bug8_password_protected(self, test_data_dir):
        """Bug #8: 密码保护的PDF"""
        # 需要密码保护的PDF测试数据
        pass
    
    def test_bug9_scanned_pdf(self, test_data_dir):
        """Bug #9: 扫描版PDF(图片)"""
        # 需要扫描版PDF测试数据
        pass
    
    def test_bug10_multi_column_pdf(self, test_data_dir):
        """Bug #10: 多栏PDF布局"""
        # 需要多栏PDF测试数据
        pass


class TestReadDocxBugs:
    """read_docx Bug测试"""
    
    def test_bug1_file_not_exist(self):
        """Bug #1: 文件不存在"""
        result = read_docx("Z:/not_exist.docx")
        assert is_error(result)
    
    def test_bug2_doc_format(self, test_data_dir, tmp_path):
        """Bug #2: .doc格式处理(需要pandoc)"""
        # 需要测.doc文件
        pass
    
    def test_bug3_empty_docx(self, tmp_path):
        """Bug #3: 空Word文档"""
        # 创建空Word文档
        from docx import Document
        empty_doc = tmp_path / "empty.docx"
        doc = Document()
        doc.save(str(empty_doc))
        result = read_docx(str(empty_doc))
        assert is_success(result)
        # Bug: 空文档应该返回什么?
        data = result.get("data", {})
        assert data.get("text", "") == ""
    
    def test_bug4_large_docx(self, test_data_dir):
        """Bug #4: 大Word文档性能"""
        pass
    
    def test_bug5_corrupted_docx(self, tmp_path):
        """Bug #5: 损坏的Word文档"""
        corrupted = tmp_path / "corrupted.docx"
        corrupted.write_bytes(b"PK\x03\x04corrupted")
        result = read_docx(str(corrupted))
        assert is_error(result)
    
    def test_bug6_complex_tables(self, tmp_path):
        """Bug #6: 复杂表格(合并单元格)"""
        from docx import Document
        doc = Document()
        table = doc.add_table(rows=3, cols=3)
        # 合并单元格
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 0).text = "合并"
        doc_path = tmp_path / "merged.docx"
        doc.save(str(doc_path))
        result = read_docx(str(doc_path))
        assert is_success(result)
        # Bug: 合并单元格如何处理?
        data = result.get("data", {})
        tables = data.get("tables", [])
        # 合并单元格可能导致数据重复或丢失
    
    def test_bug7_images_in_docx(self, tmp_path):
        """Bug #7: 包含图片的Word文档"""
        # 需要包含图片的Word文档
        pass
    
    def test_bug8_nested_tables(self, tmp_path):
        """Bug #8: 嵌套表格"""
        # Word不支持嵌套表格,跳过
        pass
    
    def test_bug9_special_formatting(self, tmp_path):
        """Bug #9: 特殊格式(上标,下标,公式)"""
        from docx import Document
        doc = Document()
        para = doc.add_paragraph()
        run = para.add_run("正常")
        run2 = para.add_run("上标")
        run2.font.superscript = True
        doc_path = tmp_path / "format.docx"
        doc.save(str(doc_path))
        result = read_docx(str(doc_path))
        assert is_success(result)
        # Bug: 特殊格式信息丢失
    
    def test_bug10_headers_footers(self, tmp_path):
        """Bug #10: 页眉页脚"""
        from docx import Document
        doc = Document()
        section = doc.sections[0]
        header = section.header
        header_para = header.paragraphs[0]
        header_para.text = "页眉内容"
        doc_path = tmp_path / "header.docx"
        doc.save(str(doc_path))
        result = read_docx(str(doc_path))
        assert is_success(result)
        # Bug: 页眉页脚内容未提取


class TestReadPptxBugs:
    """read_pptx Bug测试"""
    
    def test_bug1_file_not_exist(self):
        """Bug #1: 文件不存在"""
        result = read_pptx("Z:/not_exist.pptx")
        assert is_error(result)
    
    def test_bug2_empty_pptx(self, tmp_path):
        """Bug #2: 空PPT"""
        from pptx import Presentation
        prs = Presentation()
        empty_ppt = tmp_path / "empty.pptx"
        prs.save(str(empty_ppt))
        result = read_pptx(str(empty_ppt))
        assert is_success(result)
        data = result.get("data", {})
        assert data.get("slide_count", 0) == 0
    
    def test_bug3_large_pptx(self, test_data_dir):
        """Bug #3: 大PPT性能"""
        pass
    
    def test_bug4_corrupted_pptx(self, tmp_path):
        """Bug #4: 损坏的PPT"""
        corrupted = tmp_path / "corrupted.pptx"
        corrupted.write_bytes(b"PK\x03\x04corrupted")
        result = read_pptx(str(corrupted))
        assert is_error(result)
    
    def test_bug5_images_in_pptx(self, tmp_path):
        """Bug #5: 包含图片的PPT"""
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # 添加图片(需要真实图片)
        # slide.shapes.add_picture(...)
        ppt_path = tmp_path / "image.pptx"
        prs.save(str(ppt_path))
        result = read_pptx(str(ppt_path))
        # Bug: 图片信息未提取
    
    def test_bug6_tables_in_pptx(self, tmp_path):
        """Bug #6: 包含表格的PPT"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        table = shapes.add_table(3, 3, 0, 0, 100, 100).table
        table.cell(0, 0).text = "表格内容"
        ppt_path = tmp_path / "table.pptx"
        prs.save(str(ppt_path))
        result = read_pptx(str(ppt_path))
        assert is_success(result)
        # Bug: 表格内容未提取
    
    def test_bug7_notes_extraction(self, tmp_path):
        """Bug #7: 备注提取"""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "这是备注"
        ppt_path = tmp_path / "notes.pptx"
        prs.save(str(ppt_path))
        result = read_pptx(str(ppt_path))
        assert is_success(result)
        data = result.get("data", {})
        # Bug: 备注应该被提取
        notes = data.get("notes", [])
        assert len(notes) > 0
    
    def test_bug8_master_slides(self, tmp_path):
        """Bug #8: 母版幻灯片"""
        pass
    
    def test_bug9_animations(self, tmp_path):
        """Bug #9: 动画效果"""
        pass
    
    def test_bug10_hidden_slides(self, tmp_path):
        """Bug #10: 隐藏幻灯片"""
        pass


class TestReadXlsxBugs:
    """read_xlsx Bug测试"""
    
    def test_bug1_file_not_exist(self):
        """Bug #1: 文件不存在"""
        result = read_xlsx("Z:/not_exist.xlsx")
        assert is_error(result)
    
    def test_bug2_empty_xlsx(self, tmp_path):
        """Bug #2: 空Excel"""
        from openpyxl import Workbook
        wb = Workbook()
        empty_xlsx = tmp_path / "empty.xlsx"
        wb.save(str(empty_xlsx))
        result = read_xlsx(str(empty_xlsx))
        assert is_success(result)
        data = result.get("data", {})
        # Bug: 空Excel应该返回什么?
        assert data.get("row_count", 0) == 0
    
    def test_bug3_csv_format(self, tmp_path):
        """Bug #3: CSV格式"""
        csv_file = tmp_path / "test.csv"
        csv_file.write_text("A,B,C\n1,2,3\n4,5,6", encoding="utf-8")
        result = read_xlsx(str(csv_file))
        assert is_success(result)
        data = result.get("data", {})
        # 当前: row_count 已从 data 移入 llm_data.metrics, data 直接含 rows
        assert len(data.get("rows", [])) == 2
    
    def test_bug4_xls_format(self, tmp_path):
        """Bug #4: .xls格式"""
        # 需要xls测试数据
        pass
    
    def test_bug5_large_xlsx(self, test_data_dir):
        """Bug #5: 大Excel性能"""
        pass
    
    def test_bug6_corrupted_xlsx(self, tmp_path):
        """Bug #6: 损坏的Excel"""
        corrupted = tmp_path / "corrupted.xlsx"
        corrupted.write_bytes(b"PK\x03\x04corrupted")
        result = read_xlsx(str(corrupted))
        assert is_error(result)
    
    def test_bug7_multi_sheet(self, tmp_path):
        """Bug #7: 多工作表"""
        from openpyxl import Workbook
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["A", "B"])
        ws1.append([1, 2])
        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["C", "D"])
        ws2.append([3, 4])
        multi_xlsx = tmp_path / "multi.xlsx"
        wb.save(str(multi_xlsx))
        result = read_xlsx(str(multi_xlsx))
        assert is_success(result)
        data = result.get("data", {})
        # Bug: 只读取了第一个工作表
        assert len(data.get("sheet_names", [])) == 2
    
    def test_bug8_formula_cells(self, tmp_path):
        """Bug #8: 公式单元格"""
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.append(["A", "B", "C"])
        ws.append([1, 2, "=A2+B2"])  # 公式
        formula_xlsx = tmp_path / "formula.xlsx"
        wb.save(str(formula_xlsx))
        result = read_xlsx(str(formula_xlsx))
        assert is_success(result)
        data = result.get("data", {})
        rows = data.get("rows", [])
        # Bug: 公式结果应该是3,但可能返回公式字符串
        # data_only=True应该返回计算结果
    
    def test_bug9_merged_cells(self, tmp_path):
        """Bug #9: 合并单元格"""
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.append(["A", "B", "C"])
        ws.append([1, 2, 3])
        ws.merge_cells('A2:B2')
        merged_xlsx = tmp_path / "merged.xlsx"
        wb.save(str(merged_xlsx))
        result = read_xlsx(str(merged_xlsx))
        assert is_success(result)
        # Bug: 合并单元格如何处理?
    
    def test_bug10_different_encodings_csv(self, tmp_path):
        """Bug #10: 不同编码的CSV"""
        # GBK编码的CSV
        gbk_csv = tmp_path / "gbk.csv"
        gbk_csv.write_bytes("姓名,年龄\n张三,25".encode("gbk"))
        result = read_xlsx(str(gbk_csv))
        assert is_success(result)
        data = result.get("data", {})
        # Bug: 应该自动检测编码并正认读取


class TestReadToolsNegative:
    """为面测试"""
    
    def test_read_pdf_invalid_path(self):
        """无效路径"""
        assert is_error(read_pdf("invalid/path.pdf"))
    
    def test_read_docx_invalid_path(self):
        """无效路径"""
        assert is_error(read_docx("invalid/path.docx"))
    
    def test_read_pptx_invalid_path(self):
        """无效路径"""
        assert is_error(read_pptx("invalid/path.pptx"))
    
    def test_read_xlsx_invalid_path(self):
        """无效路径"""
        assert is_error(read_xlsx("invalid/path.xlsx"))


class TestReadToolsSchemaIssues:
    """Schema问题测试"""
    
    def test_read_pdf_schema(self):
        """read_pdf Schema"""
        from app.tools.document.document_schema import ReadPdfInput
        field = ReadPdfInput.model_fields['path']
        # Bug: Schema描述是否清晰?
        assert field.description is not None
    
    def test_read_docx_schema(self):
        """read_docx Schema"""
        from app.tools.document.document_schema import ReadDocxInput
        field = ReadDocxInput.model_fields['path']
        assert field.description is not None
    
    def test_read_pptx_schema(self):
        """read_pptx Schema"""
        from app.tools.document.document_schema import ReadPptxInput
        field = ReadPptxInput.model_fields['path']
        assert field.description is not None
    
    def test_read_xlsx_schema(self):
        """read_xlsx Schema"""
        from app.tools.document.document_schema import ReadXlsxInput
        field = ReadXlsxInput.model_fields['path']
        assert field.description is not None
